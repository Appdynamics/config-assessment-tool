import json
import logging
import os
from datetime import datetime
from enum import Enum
from typing import List

import click
from openpyxl import load_workbook
from tzlocal import get_localzone
from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx.slide import Slide
from pptx.util import Inches, Pt


class Color(Enum):
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    RED = RGBColor(255, 0, 0)
    GREEN = RGBColor(0, 255, 0)
    BLUE = RGBColor(0, 0, 255)


def setBackgroundImage(root: Presentation, slide: Slide, image_path: str):
    left = top = Inches(0)
    img_path = image_path
    pic = slide.shapes.add_picture(img_path, left, top, width=root.slide_width, height=root.slide_height)

    # Move it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def setTitle(slide: Slide, text: str, color: Color = Color.BLACK, fontSize: int = 32):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(fontSize)
    title.text_frame.paragraphs[0].font.color.rgb = color.value


def addBulletedText(slide: Slide, text: List[str], color: Color = Color.BLACK, fontSize: int = 12):
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    if len(text) == 0:
        return

    tf = body_shape.text_frame
    tf.text = text[0]

    for i in range(1, len(text)):
        p = tf.add_paragraph()
        p.text = text[i]

    for paragraph in body_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(fontSize)
            run.font.color.rgb = color.value


def addTable(slide, data, color: Color = Color.BLACK, fontSize: int = 16, left: int = 0.25, top: int = 3.5, width: int = 9.5, height: int = 1.5):
    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table

    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            for paragraph in table.cell(i, j).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fontSize)
                    run.font.color.rgb = color.value


def getValuesInColumn(sheet, param):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == param:
            j = 0
            for data in column_cell[1:]:
                values.append(data.value)
            break
    return values


@click.command()
@click.option("--folder", "-f", help="Output folder to read from", required=True)
def main(folder: str):
    logging.info(f"Creating presentation from output folder: {folder}")

    root = Presentation()

    """ Ref for slide types: 
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only 
    6 ->  Blank
    7 ->  Content with caption
    8 ->  Pic with caption
    """

    # Title Slide
    logging.info(f"Creating Title Slide")
    slide = root.slides.add_slide(root.slide_layouts[0])
    setBackgroundImage(root, slide, "bin/ppt-assets/background.jpg")
    setTitle(slide, f"{folder} Configuration Assessment Highlights", Color.WHITE, fontSize=48)
    info = json.loads(open(f"output/{folder}/info.json").read())
    slide.placeholders[1].text = f'Data As Of: {datetime.fromtimestamp(info["lastRun"], get_localzone()).strftime("%m-%d-%Y at %H:%M:%S")}'

    # Criteria Slide
    logging.info(f"Creating Criteria Slide")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"Configuration Assessment Tool Criteria")
    slide.shapes.add_picture("bin/ppt-assets/criteria.png", Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(5))

    # Criteria ctd... Slide
    logging.info(f"Creating Criteria Slide ctd...")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"Configuration Assessment Tool Criteria ctd...")
    slide.shapes.add_picture("bin/ppt-assets/criteria2.png", Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(4))

    # S/G/P Criteria & Scoring Slide
    logging.info(f"Creating S/G/P Criteria & Scoring Slide")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"B/S/G/P Criteria & Scoring")
    text = [
        "Platinum - 70% of criteria for an application must be Platinum, all remaining criteria must be at least Gold",
        "Gold - 70% of criteria for an application must be Gold or higher, all remaining criteria must be at least Silver",
        "Silver - 70% of criteria for an application must be Silver or higher",
        "Bronze - All remaining applications",
    ]
    addBulletedText(slide, text)
    wb = load_workbook(filename=f"output/{folder}/{folder}-MaturityAssessment-apm.xlsx")
    totalApplications = wb["Analysis"].max_row - 1
    sheet = wb["Analysis"]
    scores = getValuesInColumn(sheet, "OverallAssessment")
    data = [
        ["Controller", "Apps", "Bronze% (#)", "Silver% (#)", "Gold% (#)", "Platinum% (#)"],
        [
            folder,
            str(totalApplications),
            f"{format(scores.count('bronze') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('bronze')})",
            f"{format(scores.count('silver') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('silver')})",
            f"{format(scores.count('gold') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('gold')})",
            f"{format(scores.count('platinum') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('platinum')})",
        ],
    ]
    addTable(slide, data)

    # App & Machine Agents
    logging.info(f"Creating App & Machine Agents Slide")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"App & Machine Agents")
    text = [
        "AppD Agents are supported for 1 year from release",
        "*Machine Agents Reporting No Data includes uninstrumented apps",
    ]
    addBulletedText(slide, text)
    percentAgentsLessThan1YearOld = getValuesInColumn(wb["AppAgentsAPM"], "percentAgentsLessThan1YearOld")
    percentAgentsReportingData = getValuesInColumn(wb["AppAgentsAPM"], "percentAgentsReportingData")
    percentMachineAgentsLessThan1YearOld = getValuesInColumn(wb["MachineAgentsAPM"], "percentAgentsLessThan1YearOld")
    percentMachineAgentsReportingData = getValuesInColumn(wb["MachineAgentsAPM"], "percentAgentsReportingData")
    data = [
        [
            "Controller",
            "100% App Agents Older Than 1yr",
            "% App Agents Reporting No Data",
            "100% Machine Agents Older Than 1yr",
            "% Machine Agents Reporting No Data*",
        ],
        [
            folder,
            str(format(len([x for x in percentAgentsLessThan1YearOld if x != 100]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentAgentsReportingData if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentMachineAgentsLessThan1YearOld if x != 100]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentMachineAgentsReportingData if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Business Transactions
    logging.info(f"Creating Business Transactions Slide")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"Business Transactions")
    text = [
        "The more BTs monitored, the less meaningful they become",
        "AppD recommends limiting monitored BTs per app to fewer than 200",
    ]
    addBulletedText(slide, text)
    numberOfBTs = getValuesInColumn(wb["BusinessTransactionsAPM"], "numberOfBTs")
    percentBTsWithLoad = getValuesInColumn(wb["BusinessTransactionsAPM"], "percentBTsWithLoad")
    btLockdownEnabled = getValuesInColumn(wb["BusinessTransactionsAPM"], "btLockdownEnabled")
    numberCustomMatchRules = getValuesInColumn(wb["BusinessTransactionsAPM"], "numberCustomMatchRules")
    data = [
        [
            "Controller",
            "% Apps with >200 BTs",
            "% Apps w/BTs without Any Load",
            "% Apps without BT Lockdown Enabled",
            "% Apps without Custom Match Rules",
        ],
        [
            folder,
            str(format(len([x for x in numberOfBTs if x >= 200]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentBTsWithLoad if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in btLockdownEnabled if x == "=FALSE()"]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in numberCustomMatchRules if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Backends
    logging.info(f"Creating Backends Slide")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"Backends")
    text = [
        "Backends represent external service calls made by an application",
        "While the out of the box default configurations are good, it is recommended to configure them manually for maximum visibility",
    ]
    addBulletedText(slide, text)
    percentBackendsWithLoad = getValuesInColumn(wb["BackendsAPM"], "percentBackendsWithLoad")
    backendLimitNotHit = getValuesInColumn(wb["BackendsAPM"], "backendLimitNotHit")
    numberOfCustomBackendRules = getValuesInColumn(wb["BackendsAPM"], "numberOfCustomBackendRules")
    data = [
        [
            "Controller",
            "% Apps without Backends Detected",
            "% Apps Hitting Backend Limit",
            "% Apps without Custom Detection Rules",
        ],
        [
            folder,
            str(format(len([x for x in percentBackendsWithLoad if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in backendLimitNotHit if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in numberOfCustomBackendRules if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Overhead
    logging.info(f"Creating Overhead Slide")
    slide = root.slides.add_slide(root.slide_layouts[1])
    setTitle(slide, f"Overhead")
    text = [
        "AppDynamics allows several different means to extract additional information from the APM data",
        "While okay to use in lower environments, it is never recommended to use these features in a production environment",
        "Leaving some of these settings enabled can cause significant overhead",
    ]
    addBulletedText(slide, text)
    developerModeNotEnabledForAnyBT = getValuesInColumn(wb["OverheadAPM"], "developerModeNotEnabledForAnyBT")
    findEntryPointsNotEnabled = getValuesInColumn(wb["OverheadAPM"], "findEntryPointsNotEnabled")
    aggressiveSnapshottingNotEnabled = getValuesInColumn(wb["OverheadAPM"], "aggressiveSnapshottingNotEnabled")
    developerModeNotEnabledForApplication = getValuesInColumn(wb["OverheadAPM"], "developerModeNotEnabledForApplication")
    data = [
        [
            "Controller",
            "% Apps with BT Developer Mode Enabled",
            "% Apps with Find Entry Points Enabled",
            "% Apps with Aggressive Snapshotting Enabled",
            "% Apps with Developer Mode Enabled",
        ],
        [
            folder,
            str(format(len([x for x in developerModeNotEnabledForAnyBT if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in findEntryPointsNotEnabled if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in aggressiveSnapshottingNotEnabled if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in developerModeNotEnabledForApplication if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Saving file
    logging.info(f"Saving presentation to output/{folder}/{folder}-cx-presentation.pptx")
    root.save(f"output/{folder}/{folder}-cx-presentation.pptx")


if __name__ == "__main__":
    # cd to config-assessment-tool root directory
    path = os.path.realpath(f"{__file__}/../..")
    os.chdir(path)

    # create logs and output directories
    if not os.path.exists("logs"):
        os.makedirs("logs")
    if not os.path.exists("output"):
        os.makedirs("output")

    # init logging
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        handlers=[
            logging.FileHandler("logs/asd.log"),
            logging.StreamHandler(),
        ],
    )

    logging.info(f"Working directory is {os.getcwd()}")

    main()
