import json
import logging
import os
from datetime import datetime
from enum import Enum
from typing import List

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.util import Inches, Pt
from tzlocal import get_localzone

class Color(Enum):
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    RED = RGBColor(255, 0, 0)
    GREEN = RGBColor(0, 255, 0)
    BLUE = RGBColor(0, 0, 255)

def delete_shape(shape):
    sp = shape.element
    shape.element.getparent().remove(sp)

def remove_overlapping_shapes(slide, left_inch, top_inch, width_inch, height_inch, exclude_shapes=None):
    if exclude_shapes is None:
        exclude_shapes = []

    # Convert to EMUs for comparison
    l = Inches(left_inch)
    t = Inches(top_inch)
    w = Inches(width_inch)
    h = Inches(height_inch)
    r = l + w
    b = t + h

    # Collect shapes to remove
    to_remove = []
    for shape in slide.shapes:
        if shape in exclude_shapes:
            continue

        # Check intersection
        # We need to handle shapes that might not have simple left/top attributes exposed directly or correctly for some types
        # but generally shapes in slide.shapes do.
        if not hasattr(shape, 'left'):
            continue

        sl = shape.left
        st = shape.top
        sw = shape.width
        sh = shape.height
        sr = sl + sw
        sb = st + sh

        # Check intersection
        # No intersection if: R1 < L2 or L1 > R2 or B1 < T2 or T1 > B2
        if not (r < sl or l > sr or b < st or t > sb):
            # Intersection detected.
            # Only remove if it looks like content we want to replace (Charts, Tables, Pictures, empty TextBoxes)
            # Preserve TextBoxes that have text (titles, descriptions)

            keep_shape = False
            if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                # It has text. Unless it's just "Click to edit...", we keep it.
                # Simplistic check: Keep all text shapes.
                keep_shape = True

            if not keep_shape:
                to_remove.append(shape)

    for shape in to_remove:
        try:
            delete_shape(shape)
        except Exception as e:
             logging.warning(f"Could not delete overlapping shape: {e}")

def updateTitle(slide: Slide, text: str):
    if slide.shapes.title:
        slide.shapes.title.text = text

def addBulletedText(slide: Slide, text: List[str], fontSize: int = 12):
    if len(slide.shapes.placeholders) < 2:
        return

    body_shape = slide.shapes.placeholders[1]
    if not body_shape.has_text_frame:
        return

    tf = body_shape.text_frame
    if not text:
        tf.clear()
        return

    tf.text = text[0]

    for i in range(1, len(text)):
        p = tf.add_paragraph()
        p.text = text[i]

    # Preserve template formatting where possible, but if needed we can force size
    # Uncommenting font size enforcement to match cxPpt.py behavior if template defaults are wrong
    for paragraph in body_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(fontSize)

def addNestedBulletedText(slide: Slide, text: dict, headerFontSize: int = 24, subHeaderFontSize: int = 16, color: Color = Color.WHITE):
    if len(slide.shapes.placeholders) < 2:
        return

    body_shape = slide.shapes.placeholders[1]
    if not body_shape.has_text_frame:
        return

    # Lower the text box to avoid crowding the title
    body_shape.top = Inches(2.0)

    tf = body_shape.text_frame

    # helper to clear and start fresh
    tf.clear()

    if not text:
        return

    first_key = next(iter(text))
    p = tf.paragraphs[0]
    p.text = first_key
    for subheader in text[first_key]:
        p = tf.add_paragraph()
        p.text = subheader
        p.level = 1

    firstIter = True
    for header in text.keys():
        if header == first_key:
            continue

        p = tf.add_paragraph()
        p.text = header
        for subheader in text[header]:
            p = tf.add_paragraph()
            p.text = subheader
            p.level = 1

    for paragraph in body_shape.text_frame.paragraphs:
        fontSize = headerFontSize if paragraph.level == 0 else subHeaderFontSize
        for run in paragraph.runs:
            run.font.size = Pt(fontSize)
            run.font.color.rgb = color.value

def addTable(slide, data, fontSize: int = 16, left: float = None, top: float = 3.5, width: float = 9.5, height: float = 1.5):
    # Calculate centering if left is None
    if left is None:
        try:
             # Try to access presentation slide width
             slide_width = slide.part.package.presentation_part.presentation.slide_width
        except:
             slide_width = Inches(13.333) # Fallback to 16:9

        slide_width_inches = slide_width / 914400
        left_inches = (slide_width_inches - width) / 2
    else:
        left_inches = left

    # Remove overlapping shapes essentially clearing the area for the table
    exclude = []
    if slide.shapes.title:
        exclude.append(slide.shapes.title)

    # Try to identify the body placeholder to exclude it (usually idx 1)
    if len(slide.placeholders) > 1:
        try:
            exclude.append(slide.placeholders[1])
        except KeyError:
            pass

    remove_overlapping_shapes(slide, left_inches, top, width, height, exclude_shapes=exclude)

    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(left_inches), Inches(top), Inches(width), Inches(height))
    table = shape.table

    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            for paragraph in table.cell(i, j).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fontSize)
                    run.font.color.rgb = Color.BLACK.value

def getValuesInColumn(sheet, param):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == param:
            for data in column_cell[1:]:
                values.append(data.value)
            break
    return values

def getAppsWithScore(sheet, assessmentScore):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == "OverallAssessment":
            for idx, data in enumerate(column_cell[1:]):
                if data.value == assessmentScore:
                    # +2 because idx is 0-based and rows are 1-based, and header is row 1
                    values.append(sheet[f"C{idx + 2}"].value)
            break
    return values

def createCxPpt(folder, output_dir="output"):
    logging.info(f"Creating presentation from template for output folder: {folder}")

    template_path = "backend/resources/pptAssets/cxPpt_template.pptx"
    try:
        root = Presentation(template_path)
    except Exception as e:
        logging.error(f"Failed to load template {template_path}: {e}")
        return

    job_dir = os.path.join(output_dir, folder)
    info_path = os.path.join(job_dir, "info.json")
    if not os.path.exists(info_path):
        logging.warning(f"Info file {info_path} does not exist. Skipping PPT generation.")
        return

    info = json.loads(open(info_path).read())

    ma_path = os.path.join(job_dir, f"{folder}-MaturityAssessment-apm.xlsx")
    if not os.path.exists(ma_path):
        logging.warning(f"Maturity Assessment file {ma_path} does not exist. Skipping PPT generation.")
        return

    wb = load_workbook(filename=ma_path)
    totalApplications = wb["Analysis"].max_row - 1
    sheet = wb["Analysis"]
    scores = getValuesInColumn(sheet, "OverallAssessment")

    # 1. Title Slide (Slide 0)
    slide = root.slides[0]
    updateTitle(slide, f"{folder} Configuration Assessment Highlights")
    if len(slide.shapes.placeholders) > 1:
        # Assuming placeholder 1 is subtitle/date
        slide.shapes.placeholders[1].text = f'Data As Of: {datetime.fromtimestamp(info["lastRun"], get_localzone()).strftime("%m-%d-%Y at %H:%M:%S")}'

    # 2. Current State (Slide 2 - Section Header)
    slide = root.slides[2]
    updateTitle(slide, "Current State")

    # 3. Criteria & Scoring (Slide 3)
    slide = root.slides[3]
    # Keep title from template: "Criteria & Scoring"

    data = [
        ["Controller", "Apps", "Bronze% (#)", "Silver% (#)", "Gold% (#)", "Platinum% (#)"],
        [
            folder,
            str(totalApplications),
            f"{format((scores.count('bronze') / (wb['Analysis'].max_row - 1) if wb['Analysis'].max_row - 1 > 0 else 0) * 100, '.0f')}% ({scores.count('bronze')})",
            f"{format((scores.count('silver') / (wb['Analysis'].max_row - 1) if wb['Analysis'].max_row - 1 > 0 else 0) * 100, '.0f')}% ({scores.count('silver')})",
            f"{format((scores.count('gold') / (wb['Analysis'].max_row - 1) if wb['Analysis'].max_row - 1 > 0 else 0) * 100, '.0f')}% ({scores.count('gold')})",
            f"{format((scores.count('platinum') / (wb['Analysis'].max_row - 1) if wb['Analysis'].max_row - 1 > 0 else 0) * 100, '.0f')}% ({scores.count('platinum')})",
        ],
    ]
    addTable(slide, data, top=5.0)

    # 4. App & Machine Agents (Slide 4)
    slide = root.slides[4]
    # Keep title: "Application and Machine Agents"

    percentAgentsLessThan1YearOld = getValuesInColumn(wb["AppAgentsAPM"], "percentAgentsLessThan1YearOld")
    percentAgentsReportingData = getValuesInColumn(wb["AppAgentsAPM"], "percentAgentsReportingData")
    percentMachineAgentsLessThan1YearOld = getValuesInColumn(wb["MachineAgentsAPM"], "percentAgentsLessThan1YearOld")
    percentMachineAgentsReportingData = getValuesInColumn(wb["MachineAgentsAPM"], "percentAgentsReportingData")

    data_agents = [
        [
            "Controller",
            "100% App Agents Older Than 1yr",
            "% App Agents Reporting No Data",
            "100% Machine Agents Older Than 1yr",
            "% Machine Agents Reporting No Data*",
        ],
        [
            folder,
            str(format((len([x for x in percentAgentsLessThan1YearOld if x != 100]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in percentAgentsReportingData if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in percentMachineAgentsLessThan1YearOld if x != 100]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in percentMachineAgentsReportingData if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data_agents, top=5.0)

    # 5. Low Hanging Fruit (Slide 6 - Section Header)
    # slide = root.slides[6]
    # updateTitle(slide, "Low Hanging Fruit")

    # 6. Overhead (Slide 7)
    slide = root.slides[7]

    developerModeNotEnabledForAnyBT = getValuesInColumn(wb["OverheadAPM"], "developerModeNotEnabledForAnyBT")
    findEntryPointsNotEnabled = getValuesInColumn(wb["OverheadAPM"], "findEntryPointsNotEnabled")
    aggressiveSnapshottingNotEnabled = getValuesInColumn(wb["OverheadAPM"], "aggressiveSnapshottingNotEnabled")
    developerModeNotEnabledForApplication = getValuesInColumn(wb["OverheadAPM"], "developerModeNotEnabledForApplication")

    data_overhead = [
        [
            "Controller",
            "% Apps with BT Developer Mode Enabled",
            "% Apps with Find Entry Points Enabled",
            "% Apps with Aggressive Snapshotting Enabled",
            "% Apps with Developer Mode Enabled",
        ],
        [
            folder,
            str(format((len([x for x in developerModeNotEnabledForAnyBT if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in findEntryPointsNotEnabled if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in aggressiveSnapshottingNotEnabled if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in developerModeNotEnabledForApplication if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data_overhead, top=5.0)

    # 7. Error Configuration (Slide 8)
    slide = root.slides[8]

    successPercentageOfWorstTransaction = getValuesInColumn(wb["ErrorConfigurationAPM"], "successPercentageOfWorstTransaction")
    numberOfCustomRules = getValuesInColumn(wb["ErrorConfigurationAPM"], "numberOfCustomRules")

    data_error = [
        [
            "Controller",
            "% Apps w/BTs at 100% Error Rate",
            "% Apps without Custom Error Detection Rules",
        ],
        [
            folder,
            str(format((len([x for x in successPercentageOfWorstTransaction if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in numberOfCustomRules if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data_error, top=5.0)

    # 8. Health Rules and Alerting (Slide 9)
    slide = root.slides[9]

    numberOfHealthRuleViolations = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfHealthRuleViolations")
    numberOfDefaultHealthRulesModified = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfDefaultHealthRulesModified")
    numberOfActionsBoundToEnabledPolicies = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfActionsBoundToEnabledPolicies")
    numberOfCustomHealthRules = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfCustomHealthRules")

    data_hr = [
        [
            "Controller",
            "% Apps with > 50 Health Rule Violations/Day",
            "% Apps without Default Health Rules Modified",
            "% Apps not Sending Alerts Anywhere",
            "% Apps without Custom Health Rules",
        ],
        [
            folder,
            str(format((len([x for x in numberOfHealthRuleViolations if x >= 50]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in numberOfDefaultHealthRulesModified if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in numberOfActionsBoundToEnabledPolicies if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
            str(format((len([x for x in numberOfCustomHealthRules if x == 0]) / totalApplications if totalApplications > 0 else 0) * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data_hr, top=5.0)

    # 9. Recommendations (Slide 10 - Section Header)
    # slide = root.slides[10]
    # updateTitle(slide, "Recommendations")

    # 10. Low-Hanging Fruit List (Slide 11)
    slide = root.slides[11]
    updateTitle(slide, "Low-Hanging Fruit")
    text = {
        "Remove Overhead inducing settings": [
            "Can be done on controller without app team involvement",
            "Total time investment: 1-2 hours per controller",
        ],
        "Create Error Detection rules to reduce error rates": [
            "Requires app team involvement",
            "Total time investment: 1-2 hours per application",
        ],
        "Disable Health Rules which violate too frequently": [
            "Requires app team involvement",
            "Total time investment: 1-2 hours per application",
        ],
        "Add Policies and Actions to send off events": [
            "Requires app team involvement",
            "Total time investment: 1-2 hours per application",
        ],
    }
    addNestedBulletedText(slide, text)

    # 11. Raise Gold Apps (Slide 12)
    slide = root.slides[12]
    updateTitle(slide, "Raise Gold Apps to Platinum Status")
    goldApps = getAppsWithScore(wb["Analysis"], "gold")
    text_gold = {
        f"These apps are currently in Gold status. See {folder}-MaturityAssessment-apm.xlsx Analysis sheet for a full set of applications.": [],
        "We recommend working with them to raise them to Platinum status:": goldApps[:10],
    }
    addNestedBulletedText(slide, text_gold)

    # 12. Appendix (Slide 13 - Section Header)
    # slide = root.slides[13]
    # updateTitle(slide, "Appendix")

    # Saving file
    output_path = os.path.join(job_dir, f"{folder}-cx-presentation.pptx")
    logging.info(f"Saving presentation to {output_path}")
    root.save(output_path)
