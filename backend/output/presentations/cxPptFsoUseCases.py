import json
import logging
import os
import re
from enum import Enum
from typing import Dict, Optional

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt, Inches


from enum import Enum

class SlideId(Enum):
    INTRO = ("INTRO", 0)
    RACETRACK = ("RACETRACK", 1)
    ONBOARD = ("onboard", 2)
    ONBOARD_REMEDIATION_PRIMARY = ("onboard_remediation_primary", 3)
    ONBOARD_REMEDIATION_SECONDARY = ("onboard_remediation_secondary", 4)
    IMPLEMENT = ("implement", 5)
    IMPLEMENT_REMEDIATION_PRIMARY = ("implement_remediation_primary", 6)
    IMPLEMENT_REMEDIATION_SECONDARY = ("implement_remediation_secondary", 7)
    USE = ("use", 8)
    USE_REMEDIATION_PRIMARY = ("use_remediation_primary", 9)
    USE_REMEDIATION_SECONDARY = ("use_remediation_secondary", 10)
    ENGAGE = ("engage", 11)
    ENGAGE_REMEDIATION_PRIMARY = ("engage_remediation_primary", 12)
    ENGAGE_REMEDIATION_SECONDARY = ("engage_remediation_secondary", 13)
    ADOPT = ("adopt", 14)
    ADOPT_REMEDIATION_PRIMARY = ("adopt_remediation_primary", 15)
    ADOPT_REMEDIATION_SECONDARY = ("adopt_remediation_secondary", 16)
    OPTIMIZE = ("optimize", 17)
    OPTIMIZE_REMEDIATION_PRIMARY = ("optimize_remediation_primary", 18)
    OPTIMIZE_REMEDIATION_SECONDARY = ("optimize_remediation_secondary", 19)
    FINAL = ("FINAL", 20)

    @property
    def id(self) -> str:
        return self.value[0]

    @property
    def number(self) -> int:
        return self.value[1]

    def isSecondary(slide):
        if slide in [
            SlideId.ONBOARD_REMEDIATION_SECONDARY,
            SlideId.IMPLEMENT_REMEDIATION_SECONDARY,
            SlideId.USE_REMEDIATION_SECONDARY,
            SlideId.ENGAGE_REMEDIATION_SECONDARY,
            SlideId.ADOPT_REMEDIATION_SECONDARY,
            SlideId.OPTIMIZE_REMEDIATION_SECONDARY
        ]:
            return True


class ExcelSheets(object):
    def __init__(self, directory: str, filenames: tuple):
        self.workbooks = {}
        for filename in filenames:
            try:
                full_path = os.path.join(directory, filename)
                self.workbooks[filename] = pd.read_excel(full_path, sheet_name=None)
            except Exception as e:
                logging.warning(f"Not able to load workbook {filename}")
                logging.warning(f"Was it generated to begin with? Ignoring. "
                                f"Continuing with the rest of the workbooks.")

    def findSheetByHeader(self, header_name):
        result = []
        # iterate over all loaded workbooks
        for filename, workbook in self.workbooks.items():
            # iterate over all sheets in the workbook
            for sheet_name, df in workbook.items():
                # check if the given header exists in the sheet
                if header_name in df.columns:
                    return (filename, sheet_name)
        return None

    def getHeaders(self, filename, sheetname):
        if filename not in self.workbooks:
            logging.error(f"No workbook found with filename: {filename}")
            return

        if sheetname not in self.workbooks[filename]:
            logging.error(f"No sheet found with name: {sheetname} in workbook: {filename}")
            return

        df = self.workbooks[filename][sheetname]
        return list(df.columns)


    def getWorkBooks(self):
        return self.workbooks

    def validColumn(self, filename, sheet_name, column_name):

        if filename not in self.workbooks:
            logging.error(f"No workbook found with filename: {filename}")
            return False

        if sheet_name not in self.workbooks[filename]:
            logging.error(f"No sheet found with name: {sheet_name} in workbook: {filename}")
            return False

        return True

    def getColumnTotal(self, filename, sheet_name, column_name):
        df = None
        if self.validColumn(filename, sheet_name, column_name) \
                and column_name in self.workbooks[filename][sheet_name].columns:
            df = self.workbooks[filename][sheet_name]
        else:
            logging.error(f"No column found with name: {column_name}")
            return
        return df[column_name].sum()

    def getColumnAverage(self, filename, sheet_name, column_name):
        df = None
        if self.validColumn(filename, sheet_name, column_name) \
                and column_name in self.workbooks[filename][sheet_name].columns:
            df = self.workbooks[filename][sheet_name]
        else:
            logging.error(f"No column found with name: {column_name}")
            return
        return df[column_name].mean()

    def getRowCountForColumnValue(self, filename, sheet_name, column_name, operator, value):
        if filename not in self.workbooks:
            logging.error(f"No workbook found with filename: {filename}")
            return

        if sheet_name not in self.workbooks[filename]:
            logging.error(f"No sheet found with name: {sheet_name}")
            return

        df = self.workbooks[filename][sheet_name]
        if column_name not in df.columns:
            logging.error(f"No column found with name: {column_name}")
            return

        # Map string operators to Pandas methods
        operator_mapping = {
            '==': df[column_name].eq,
            '!=': df[column_name].ne,
            '<': df[column_name].lt,
            '<=': df[column_name].le,
            '>': df[column_name].gt,
            '>=': df[column_name].ge,
        }

        # Get the appropriate method and apply it
        comparison_method = operator_mapping.get(operator)
        if comparison_method:
            return comparison_method(value).sum()
        else:
            logging.error(f"Invalid operator: {operator}")
            return


    def parseExpression(self, expression):
        match = re.search(r"<eval>\s*count\((.*?) (>=|<=|>|<|==|!=) (.*?)\)\s*</eval>", expression)
        if not match:
            logging.debug(f"no evaluation found, ignoring expression: {expression}")
            return None, None, None, None, False

        column_name = match.group(1)
        operator = match.group(2)
        value_str = match.group(3)

        try:
            value = int(value_str)
        except ValueError:
            if value_str.lower() == 'true':
                value = True
            elif value_str.lower() == 'false':
                value = False
            else:
                value = value_str

        return f"count({column_name} {operator} {str(value)})", column_name, operator, value, True

    def getValue(self, filename, sheet_name, column_name, operator, value):
        return self.getRowCountForColumnValue(filename, sheet_name, column_name, operator, value)

    def substituteExpression(self, expression, sub_exp, value):
        return expression.replace("<eval>" + sub_exp + "</eval>", str(value))



class UseCase(tuple):
    def __init__(self, file: str):

        self.task_id_to_slide = {}
        self.task_id_to_health_check_status = {}

        with open(file, 'r') as file:
            data = json.load(file)
            self.pitstop_data = data['pitstop']
        self._initSlideMapping()

    def pitstop_data(self) -> Dict:
        return self.pitstop_data

    def __str__(self) -> str:
        return self.pitstop_data.values()

    def _initSlideMapping(self):
        for slide in SlideId:
            self.setSlideId(slide.id, slide.number)

    def getAllSlideIndexes(self):
        return list(self.task_id_to_slide.values())

    def getAllHealthCheckValues(self):
        return list(self.task_id_to_health_check_status.values())

    def isFailed(self, pitstop: str) -> bool:
        tasks = self.getPitstopTasks(pitstop)
        return any('fail' in self.task_id_to_health_check_status[task].lower() for task in tasks if task in self.task_id_to_health_check_status)

    def pitStopContainsFailureOrManualCheck(self, pitstop: str) -> bool:
        tasks = self.getPitstopTasks(pitstop)
        for task in tasks:
            if task in self.task_id_to_health_check_status \
                    and \
                    'fail' in self.task_id_to_health_check_status[task].lower() \
                    or \
                    'manual check' in self.task_id_to_health_check_status[task].lower():
                return True

        return False

    def _get_task_data(self, task_id: str) -> Dict[str, Optional[str]]:
        for pitstop_stage in self.pitstop_data:
            for sequence_number, task_data in self.pitstop_data[pitstop_stage]['checklist_sequence'].items():
                if task_data['task_id'] == task_id:
                    return {
                        'checklist_item': task_data['checklist_item'],
                        'tooltip': task_data.get('tooltip', None),
                        'exit_criteria_logic': task_data.get('exit_criteria_logic', None),
                        'min_pass_threshold': task_data.get('min_pass_threshold', None),
                        'remediation_items': task_data.get('remediation_steps', None),
                        'remediation_items_secondary': task_data.get('remediation_steps_secondary', None)
                    }
        return {}

    def getAllTaskIds(self):
        task_ids = []
        for pitstop_stage in self.pitstop_data:
            for sequence_number, task_data in self.pitstop_data[pitstop_stage]['checklist_sequence'].items():
                task_ids.append(task_data['task_id'])
        return task_ids

    def getPitstopTasks(self, pitstop: str) -> Dict[str, str]:
        task_ids = []
        for pitstop_stage in self.pitstop_data[pitstop].values():
            for k in pitstop_stage.values():
                task_ids.append(k['task_id'])
        return task_ids if task_ids else None

    def getChecklistItem(self, task_id: str) -> Optional[str]:
        task_data = self._get_task_data(task_id)
        return task_data['checklist_item'] if task_data else None

    def getToolTip(self, task_id: str) -> Optional[str]:
        task_data = self._get_task_data(task_id)
        return task_data.get('tooltip', None) if task_data else None

    def getExitCriteriaLogic(self, task_id: str) -> Optional[str]:
        task_data = self._get_task_data(task_id)
        return task_data.get('exit_criteria_logic', None) if task_data else None

    def getMinPassThreshold(self, task_id: str) -> Optional[str]:
        task_data = self._get_task_data(task_id)
        return task_data.get('min_pass_threshold', None) if task_data else None

    def getRemediationList(self, task_id: str) -> Optional[list]:
        task_data = self._get_task_data(task_id)
        return task_data.get('remediation_items', None) if task_data else None

    def getRemediationSecondaryList(self, task_id: str) -> Optional[list]:
        task_data = self._get_task_data(task_id)
        return task_data.get('remediation_items_secondary', None) if task_data else None

    def setSlideId(self, slide_id: str, slide_number: int):
        self.task_id_to_slide[slide_id] = slide_number

    def getSlideId(self, slide_id: str) -> Optional[int]:
        # this avoids a potential KeyError
        return self.task_id_to_slide.get(slide_id, None)

    def setHealthCheckStatus(self, task_id: str, hc: str):
        self.task_id_to_health_check_status[task_id] = hc

    def getHealthCheckStatus(self, task_id: str) -> Optional[str]:
        # this avoids a potential KeyError
        return self.task_id_to_health_check_status.get(task_id, None)


class Color(Enum):
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    RED = RGBColor(255, 0, 0)
    GREEN = RGBColor(0, 255, 0)
    BLUE = RGBColor(0, 0, 255)


def addTable(slide, data, color: Color = Color.BLACK, fontSize: int = 16, left: int = 0.25, top: int = 3.5, width: int = 9.5, height: int = 1.5):
    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table

    # Set fixed row height
    for row_index, row in enumerate(table.rows):
        if row_index != 0:  # Skip the title row
            row.height = Inches(1)

    pass_mark = search("checkmark.png", "../")
    fail_mark = search("xmark.png", "../")

    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            cell_obj = table.cell(i, j)
            cell_obj.text = str(cell)

            image_left = Inches(left+2) + table.columns[j].width  * j
            image_top = Inches(top) + table.rows[i].height * i
            # add marker image
            if "pass" in str(cell).lower() or "manual check" in str(cell).lower():
                slide.shapes.add_picture(pass_mark, image_left, image_top, width=Inches(0.2), height=Inches(0.2))
            if "fail" in str(cell).lower():
                slide.shapes.add_picture(fail_mark, image_left, image_top, width=Inches(0.2), height=Inches(0.2))

            for paragraph in cell_obj.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fontSize)
                    run.font.color.rgb = color.value

def addCell(cell, text, color=None, fontSize=16):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(fontSize)
            if color:
                run.font.color.rgb = color.value


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def makeParaBulletPointed(para):
    """Bullets are set to Arial,
        actual text can be a different font"""
    pPr = para._p.get_or_add_pPr()
    ## Set marL and indent attributes
    pPr.set('marL', '171450')
    pPr.set('indent', '171450')
    ## Add buFont
    _ = SubElement(parent=pPr,
                   tagname="a:buFont",
                   typeface="Arial",
                   panose="020B0604020202020204",
                   pitchFamily="34",
                   charset="0"
                   )
    ## Add buChar
    _ = SubElement(parent=pPr,
                   tagname='a:buChar',
                   char="•")


def addHyperLinkCell(run, cell_text):
    hyperlink_pattern = re.compile(r'<a href=[\'"](.*?)[\'"]>(.*?)<\/a>')
    parts = hyperlink_pattern.split(cell_text)
    for idx, part in enumerate(parts):
        if idx % 3 == 0:
            run.text += part
        elif idx % 3 == 1:
            hyperlink_url = part
        else:
            run.text += part
            hlink = run.hyperlink
            hlink.address = hyperlink_url


def addRemediationTable(slide, data, color=None, fontSize=16, left=1.5, top=3.5, width=9.5, height=1.5):
    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[2].width = Inches(2)
    table.columns[2].width = Inches(7)
    hyperlink_pattern = re.compile(r'<a href=[\'"](.*?)[\'"]>(.*?)<\/a>')

    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            cell_text = str(cell)
            text_frame = table.cell(i, j).text_frame
            # Clear any existing paragraphs
            text_frame.clear()
            if "\n" not in cell_text:
                addCell(table.cell(i, j), cell_text, color, fontSize)
                continue

            for idx, line in enumerate(cell_text.split("\n")):
                segments = re.split(hyperlink_pattern, line)

                if idx == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()

                # paragraph = text_frame.add_paragraph()
                for i, segment in enumerate(segments):
                    if i % 3 == 0:
                        paragraph.add_run().text = segment
                    elif i % 3 == 1:  # url
                        url = segment
                    elif i % 3 == 2:  # hyperlink text
                        hyperlink_run = paragraph.add_run()
                        hyperlink_run.hyperlink.address = url
                        hyperlink_run.text = segment

                makeParaBulletPointed(paragraph)
                paragraph.font.size = Pt(10)


def getValuesInColumn(sheet, col1_value):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == col1_value:
            j = 0
            for data in column_cell[1:]:
                values.append(data.value)
            break
    return values


def getValuesInColumnForController(sheet, col1_value, controller):
    values = []
    col1_index = -1
    col2_index = -1

    col2_value = "controller"

    # Find the indexes of col1_value and col2_value in the header row
    for i, cell in enumerate(sheet[1]):
        if cell.value == col1_value:
            col1_index = i
        if cell.value == col2_value:
            col2_index = i
        if col1_index != -1 and col2_index != -1:
            break

    # Check if both columns were found
    if col1_index == -1 or col2_index == -1:
        return values

    # Iterate through the rows in the worksheet, starting from row 2 to skip the header row
    for row in sheet.iter_rows(min_row=2):
        if row[col2_index].value == controller:
            values.append(row[col1_index].value)

    return values


def getRowCountForController(sheet, controller: str) -> int:
    CONTROLLER_NAME_COLUMN_INDEX = 1
    matching_rows = sum(row[0] == controller for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, min_col=CONTROLLER_NAME_COLUMN_INDEX, max_col=CONTROLLER_NAME_COLUMN_INDEX, values_only=True))
    return matching_rows


def getAppsWithScore(sheet, assessmentScore):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == "OverallAssessment":
            j = 0
            for idx, data in enumerate(column_cell[1:]):
                if data.value == assessmentScore:
                    values.append(sheet[f"C{idx + 2}"].value)  # +2 because of the header
            break
    return values

def add_image(slide, image_path, left, top, width, height):
    return slide.shapes.add_picture(image_path, left, top, width, height)


def search(filename, search_path="."):
    for root, _, files in os.walk(search_path):
        if filename in files:
            return os.path.join(root, filename)
    return None

def markRaceTrackFailures(root, uc: UseCase):
    # racetrack slide is the second slide
    slide = root.slides[1]
    pass_mark = search("checkmark.png", "../")
    fail_mark = search("xmark.png", "../")
    image_width = Inches(0.3)
    image_height = Inches(0.3)

    # IMPORTANT - EXTERNAL DEPENDENCY each pitstop has a named shape in the
    # PPTX template and is used for the iteration below else pitstop shapes won't be found
    # and will not be visually marked for fail/pass checkmarks
    pitstops = ("onboard","use","implement","optimize","adopt","engage")

    shapes_dict = {shape.name: shape for shape in slide.shapes}
    for pitstop in pitstops:
        try:
            if uc.isFailed(pitstop):
                _ = add_image(slide, fail_mark , shapes_dict[pitstop].left, shapes_dict[pitstop].top, image_width, image_height)
            else:
                _ = add_image(slide, pass_mark , shapes_dict[pitstop].left, shapes_dict[pitstop].top, image_width, image_height)
        except KeyError:
            logging.error(f"Shape with name '{pitstop}' not found in the slide. This prevents properly marking the race track with visual markers. ")


def createCxHamUseCasePpt(folder: str):
    logging.info(f"Creating CX HAM Use Case PPT for {folder}")
    directory = f"output/{folder}"
    file_prefix = f"{folder}"
    apm_wb = load_workbook(f"{directory}/{file_prefix}-MaturityAssessment-apm.xlsx")
    db_wb = load_workbook(f"{directory}/{file_prefix}-AgentMatrix.xlsx")

    excels = ExcelSheets(directory,
                         (
                             f"{file_prefix}-AgentMatrix.xlsx",
                             f"{file_prefix}-CustomMetrics.xlsx",
                             f"{file_prefix}-Dashboards.xlsx",
                             f"{file_prefix}-License.xlsx",
                             f"{file_prefix}-MaturityAssessment-apm.xlsx",
                             f"{file_prefix}-MaturityAssessment-brum.xlsx",
                             f"{file_prefix}-MaturityAssessment-mrum.xlsx",
                             f"{file_prefix}-MaturityAssessmentRaw-apm.xlsx",
                             f"{file_prefix}-MaturityAssessmentRaw-brum.xlsx",
                             f"{file_prefix}-MaturityAssessmentRaw-mrum.xlsx",
                             f"{file_prefix}-Synthetics.xlsx"
                         ))

    assert len(excels.getWorkBooks()) >= 10

    # currently only 1st controller in the job file is examined.
    controller = getValuesInColumn(apm_wb["Analysis"], "controller")[0]

    json_file = search("HybridApplicationMonitoringUseCase.json", "../")
    uc = UseCase(json_file)
    _ = calculate_kpis(apm_wb, db_wb, uc)

    template_file = search("HybridApplicationMonitoringUseCase_template.pptx", "../")
    root = Presentation(template_file)

    ############################# Onboard ###########################
    generatePitstopHealthCheckTable(controller, root, uc, "onboard")
    generateRemediationSlides(controller, root, uc, excels, "onboard", SlideId.ONBOARD_REMEDIATION_PRIMARY)
    generateRemediationSlides(controller, root, uc, excels, "onboard", SlideId.ONBOARD_REMEDIATION_SECONDARY)
    ############################# Implement ###########################
    generatePitstopHealthCheckTable(controller, root, uc, "implement")
    generateRemediationSlides(controller, root, uc, excels, "implement", SlideId.IMPLEMENT_REMEDIATION_PRIMARY)
    generateRemediationSlides(controller, root, uc, excels, "implement", SlideId.IMPLEMENT_REMEDIATION_SECONDARY)
    ############################ Use ##################################
    generatePitstopHealthCheckTable(controller, root, uc, "use")
    generateRemediationSlides(controller, root, uc, excels, "use", SlideId.USE_REMEDIATION_PRIMARY)
    generateRemediationSlides(controller, root, uc, excels, "use", SlideId.USE_REMEDIATION_SECONDARY)
    ############################ Engage ###############################
    generatePitstopHealthCheckTable(controller, root, uc, "engage")
    generateRemediationSlides(controller, root, uc, excels, "engage", SlideId.ENGAGE_REMEDIATION_PRIMARY)
    generateRemediationSlides(controller, root, uc, excels, "engage", SlideId.ENGAGE_REMEDIATION_SECONDARY)
    ############################ Adopt ###############################
    generatePitstopHealthCheckTable(controller, root, uc, "adopt")
    generateRemediationSlides(controller, root, uc, excels, "adopt", SlideId.ADOPT_REMEDIATION_PRIMARY)
    generateRemediationSlides(controller, root, uc, excels, "adopt", SlideId.ADOPT_REMEDIATION_SECONDARY)
    ########################### Optimize ##############################
    generatePitstopHealthCheckTable(controller, root, uc, "optimize")
    generateRemediationSlides(controller, root, uc, excels, "optimize", SlideId.OPTIMIZE_REMEDIATION_PRIMARY)
    generateRemediationSlides(controller, root, uc, excels, "optimize", SlideId.OPTIMIZE_REMEDIATION_SECONDARY)

    # for now, we will not clean up slides. Primary and secondary feedback slides will
    # remain as further resource for the customer
    # cleanup_slides(root, uc)
    markRaceTrackFailures(root, uc)

    root.save(f"output/{folder}/{folder}-cx-HybridApplicationMonitoringUseCaseMaturityAssessment-presentation.pptx")


def generatePitstopHealthCheckTable(folder, root, uc, pitstop):
    slide = root.slides[uc.getSlideId(pitstop)]
    data = [["Controller", "Checklist Item", "Tooltips ***", "Exit Criteria Logic"]]
    for task in uc.getPitstopTasks(pitstop):
        data.append([folder, f"{uc.getChecklistItem(task)}", f"{uc.getToolTip(task)}", uc.getHealthCheckStatus(task)])
    addTable(slide, data, fontSize=10, top=2, left=1.5)


def parseHelper(expression, excels: ExcelSheets):
    sub_exp, column_name, operator, value, is_eval_expr = excels.parseExpression(expression)
    if is_eval_expr:
        wb,sh = excels.findSheetByHeader(column_name)
        count = excels.getValue(wb, sh, column_name, operator, value)
        return excels.substituteExpression(expression, sub_exp, count)
    else:
        return expression


def generateRemediationSlides(folder: str, root: Presentation, uc: UseCase, excels: ExcelSheets, pitstop: str, remediation_slide: SlideId):

    # take out the check for now. We want to generate feedback pages regardless of pass/fail/manual status
    # if uc.pitStopContainsFailureOrManualCheck(pitstop):
    slide = root.slides[uc.getSlideId(remediation_slide.id)]
    data = []
    for task in uc.getPitstopTasks(pitstop):
        if SlideId.isSecondary(remediation_slide):
            values_to_use = uc.getRemediationSecondaryList(task)
        else:
            values_to_use = uc.getRemediationList(task)
        if not values_to_use:
            continue
        else:
            data.append(
                [folder,
                 f"{uc.getChecklistItem(task)}",
                 '\n'.join(parseHelper(value["remediation_item"],excels) for value in values_to_use.values())
                 ]
            )

    if len(data) > 0:
        data.insert(0,["Controller", "Checklist Item", "Recommendation"] )
        addRemediationTable(slide, data, fontSize=10, top=2, left=.5)

def filter_slides(keep_slide_indexes, prs: Presentation):
    total_slides = len(prs.slides)
    slides_to_remove = {i for i in range(total_slides) if i not in keep_slide_indexes}
    for slide_index in sorted(slides_to_remove, reverse=True):
        slide = prs.slides[slide_index]
        id_to_index_mapping = {slide.id: (i, slide.rId) for i, slide in enumerate(prs.slides._sldIdLst)}
        slide_id = slide.slide_id
        slide_rId = id_to_index_mapping[slide_id][1]
        prs.part.drop_rel(slide_rId)
        del prs.slides._sldIdLst[id_to_index_mapping[slide_id][0]]
    return prs


def cleanup_slides(root: Presentation, uc: UseCase):
    slides_to_keep = uc.getAllSlideIndexes()
    if not uc.pitStopContainsFailureOrManualCheck("onboard"):
        slides_to_keep.remove(uc.getSlideId(SlideId.ONBOARD_REMEDIATION_PRIMARY.id))
    if not uc.pitStopContainsFailureOrManualCheck("implement"):
        slides_to_keep.remove(uc.getSlideId(SlideId.IMPLEMENT_REMEDIATION_PRIMARY.id))
    if not uc.pitStopContainsFailureOrManualCheck("use"):
        slides_to_keep.remove(uc.getSlideId(SlideId.USE_REMEDIATION_PRIMARY.id))
    if not uc.pitStopContainsFailureOrManualCheck("engage"):
        slides_to_keep.remove(uc.getSlideId(SlideId.ENGAGE_REMEDIATION_PRIMARY.id))
    if not uc.pitStopContainsFailureOrManualCheck("adopt"):
        slides_to_keep.remove(uc.getSlideId(SlideId.ADOPT_REMEDIATION_PRIMARY.id))
    if not uc.pitStopContainsFailureOrManualCheck("optimize"):
        slides_to_keep.remove(uc.getSlideId(SlideId.OPTIMIZE_REMEDIATION_PRIMARY.id))

    filter_slides(slides_to_keep, root)


def calculate_kpis(apm_wb, agent_wb, uc: UseCase):
    # currently only supports one controller report out of the workbook
    controller = getValuesInColumn(apm_wb["Analysis"], "controller")[0]
    logging.info(f"processing report for 1st controller only as multiple controllers are not supported yet: {controller}")

    totalApplications = getRowCountForController(apm_wb["Analysis"], controller)
    percentAgentsReportingData = getValuesInColumnForController(apm_wb["AppAgentsAPM"], "percentAgentsReportingData", controller)
    countOfpercentAgentsReportingData = len([x for x in percentAgentsReportingData if x >= 1])
    numberOfBTsList = getValuesInColumnForController(apm_wb["BusinessTransactionsAPM"], "numberOfBTs", controller)
    customMatchRulesList = getValuesInColumnForController(apm_wb["BusinessTransactionsAPM"], "numberCustomMatchRules", controller)
    numberOfDataCollectorsConfigured = getValuesInColumnForController(apm_wb["DataCollectorsAPM"], "numberOfDataCollectorFieldsConfigured", controller)
    countOfNumberOfDataCollectorsConfigured = len([x for x in numberOfDataCollectorsConfigured if x >= 1])
    numberOfCustomHealthRules = getValuesInColumnForController(apm_wb["HealthRulesAndAlertingAPM"], "numberOfCustomHealthRules", controller)
    numberOfHealthRuleViolations = getValuesInColumnForController(apm_wb["HealthRulesAndAlertingAPM"], "numberOfHealthRuleViolations", controller)
    numberOfDefaultHealthRulesModified = getValuesInColumnForController(apm_wb["HealthRulesAndAlertingAPM"], "numberOfDefaultHealthRulesModified", controller)
    numberOfActionsBoundToEnabledPoliciesList = getValuesInColumnForController(apm_wb["HealthRulesAndAlertingAPM"], "numberOfActionsBoundToEnabledPolicies", controller)
    dashboardsList = getValuesInColumnForController(apm_wb["DashboardsAPM"], "numberOfDashboards", controller)

    dbAgentList = getValuesInColumnForController(agent_wb["Individual - dbAgents"], "status", controller)
    dbAgentsActiveCount = len([x for x in dbAgentList if x == "ACTIVE"])

    ### ONB
    FSO_HAM_ONB_1 = f'manual check'
    FSO_HAM_ONB_2 = f'manual check'
    FSO_HAM_ONB_3 = f'manual check'
    FSO_HAM_ONB_4 = f'manual check'

    ### IMP
    FSO_HAM_IMP_1 = f'manual check'
    FSO_HAM_IMP_2 = f'Pass ({totalApplications} applications on-boarded)' if totalApplications >= 1 else f'Fail. At least one application needs to be instrumented and reporting metric data into AppDynamics controller'
    FSO_HAM_IMP_3 = f'manual check'

    # FSO_HAM_IMP_4 = f'Pass ({customMatchRulesList})' if all(countOfCustomRule >= 2 for countOfCustomRule in customMatchRulesList) else f'Fail (Not all applications have at least 2 custom BT match rules). Only {len([x for x in customMatchRulesList if x >= 2])} have at least 2 custom match rules out of {totalApplications}.'
    FSO_HAM_IMP_4 = f'Pass' if countOfpercentAgentsReportingData == totalApplications else f'Fail (Not all application agents are reporting data). Only  {countOfpercentAgentsReportingData} reporting data out of {totalApplications}.'

    ### USE
    # FSO_HAM_USE_1 = f'Pass' if countOfpercentAgentsReportingData == totalApplications else f'Fail (Not all application agents are reporting data). Only  {countOfpercentAgentsReportingData} reporting data out of {totalApplications}.'
    FSO_HAM_USE_1 = f'Pass ({customMatchRulesList})' if all(countOfCustomRule >= 2 for countOfCustomRule in customMatchRulesList) else f'Fail (Not all applications have at least 2 custom BT match rules). Only {len([x for x in customMatchRulesList if x >= 2])} have at least 2 custom match rules out of {totalApplications}.'

    FSO_HAM_USE_2 = f'Pass' if all(count >= 5 for count in numberOfBTsList) else f'Fail (Not all applications have at least 5 Business Transactions). Only {len([x for x in numberOfBTsList if x >= 5])} have at least 5 Business Transactions out of {totalApplications}.'

    FSO_HAM_USE_3 = f'Pass' if all(count >= 2 for count in numberOfCustomHealthRules) else f'Fail (At least 2 custom health rules must be configured per application). Only {len([count for count in numberOfCustomHealthRules if count >= 2])} applications have at least 2 custom health rules configured out of {totalApplications}.'
    FSO_HAM_USE_4 = f'Pass ({countOfNumberOfDataCollectorsConfigured} data collectors found)' if countOfNumberOfDataCollectorsConfigured >= 2 else 'Fail'

    ### ENG
    FSO_HAM_ENG_1 = f'Pass' if all(count >= 5 for count in numberOfActionsBoundToEnabledPoliciesList) else f'Fail (Not all applications have at least 2 policies with actionable alerts). Only {len([x for x in numberOfActionsBoundToEnabledPoliciesList if x >= 2])} applications have at least 2 actionable alerts out of {totalApplications} applications.'

    FSO_HAM_ENG_2 = f'Pass' if all(count >= 2 for count in dashboardsList) else f'Fail (there should be at least two dashboards configured per application). Only {len([x for x in dashboardsList if x >= 2])} applications have at least two dashboards configured out of {totalApplications} applications. '
    FSO_HAM_ENG_3 = f'manual check'
    FSO_HAM_ENG_4 = f'manual check'

    ### ADO
    FSO_HAM_ADO_1 = f'manual check'
    FSO_HAM_ADO_2 = f'Pass' if dbAgentsActiveCount > 0 else f'Fail (there should be at least one Active Database Agent configured). Currently {dbAgentsActiveCount} Active Database agents are configured.'

    ### OPT
    FSO_HAM_OPT_1 = f'manual check'
    FSO_HAM_OPT_2 = f'manual check'
    FSO_HAM_OPT_3 = f'Pass' if all(count >= 6 for count in numberOfCustomHealthRules) else f'Fail (At least 6 custom health rules must be configured per applications). Only {len([count for count in numberOfCustomHealthRules if count >= 6])} application has 6 custom health rules configured.'
    FSO_HAM_OPT_4 = f'Pass' if all(count >= 10 for count in numberOfBTsList) else f'Fail (there should be at least 10 Business Translations detected per application). Only {len([count for count in numberOfBTsList if count >= 10])} applications have at least 10 Business Transaction'

    uc.setHealthCheckStatus('FSO_HAM_ONB_1', FSO_HAM_ONB_1)
    uc.setHealthCheckStatus('FSO_HAM_ONB_2', FSO_HAM_ONB_2)
    uc.setHealthCheckStatus('FSO_HAM_ONB_3', FSO_HAM_ONB_3)
    uc.setHealthCheckStatus('FSO_HAM_ONB_4', FSO_HAM_ONB_4)
    uc.setHealthCheckStatus('FSO_HAM_USE_1', FSO_HAM_USE_1)
    uc.setHealthCheckStatus('FSO_HAM_USE_2', FSO_HAM_USE_2)
    uc.setHealthCheckStatus('FSO_HAM_USE_3', FSO_HAM_USE_3)
    uc.setHealthCheckStatus('FSO_HAM_USE_4', FSO_HAM_USE_4)
    uc.setHealthCheckStatus('FSO_HAM_IMP_1', FSO_HAM_IMP_1)
    uc.setHealthCheckStatus('FSO_HAM_IMP_2', FSO_HAM_IMP_2)
    uc.setHealthCheckStatus('FSO_HAM_IMP_3', FSO_HAM_IMP_3)
    uc.setHealthCheckStatus('FSO_HAM_IMP_4', FSO_HAM_IMP_4)
    uc.setHealthCheckStatus('FSO_HAM_ENG_1', FSO_HAM_ENG_1)
    uc.setHealthCheckStatus('FSO_HAM_ENG_2', FSO_HAM_ENG_2)
    uc.setHealthCheckStatus('FSO_HAM_ENG_3', FSO_HAM_ENG_3)
    uc.setHealthCheckStatus('FSO_HAM_ENG_4', FSO_HAM_ENG_4)
    uc.setHealthCheckStatus('FSO_HAM_ADO_1', FSO_HAM_ADO_1)
    uc.setHealthCheckStatus('FSO_HAM_ADO_2', FSO_HAM_ADO_2)
    uc.setHealthCheckStatus('FSO_HAM_OPT_1', FSO_HAM_OPT_1)
    uc.setHealthCheckStatus('FSO_HAM_OPT_2', FSO_HAM_OPT_2)
    uc.setHealthCheckStatus('FSO_HAM_OPT_3', FSO_HAM_OPT_3)
    uc.setHealthCheckStatus('FSO_HAM_OPT_4', FSO_HAM_OPT_4)

    kpi_dictionary = {
        'FSO_HAM_ONB_1': FSO_HAM_ONB_1,
        'FSO_HAM_ONB_2': FSO_HAM_ONB_2,
        'FSO_HAM_ONB_3': FSO_HAM_ONB_3,
        'FSO_HAM_ONB_4': FSO_HAM_ONB_4,
        'FSO_HAM_USE_1': FSO_HAM_USE_1,
        'FSO_HAM_USE_2': FSO_HAM_USE_2,
        'FSO_HAM_USE_3': FSO_HAM_USE_3,
        'FSO_HAM_USE_4': FSO_HAM_USE_4,
        'FSO_HAM_IMP_1': FSO_HAM_IMP_1,
        'FSO_HAM_IMP_2': FSO_HAM_IMP_2,
        'FSO_HAM_IMP_3': FSO_HAM_IMP_3,
        'FSO_HAM_IMP_4': FSO_HAM_IMP_4,
        'FSO_HAM_ENG_1': FSO_HAM_ENG_1,
        'FSO_HAM_ENG_2': FSO_HAM_ENG_2,
        'FSO_HAM_ENG_3': FSO_HAM_ENG_3,
        'FSO_HAM_ENG_4': FSO_HAM_ENG_4,
        'FSO_HAM_ADO_1': FSO_HAM_ADO_1,
        'FSO_HAM_ADO_2': FSO_HAM_ADO_2,
        'FSO_HAM_OPT_1': FSO_HAM_OPT_1,
        'FSO_HAM_OPT_2': FSO_HAM_OPT_2,
        'FSO_HAM_OPT_3': FSO_HAM_OPT_3,
        'FSO_HAM_OPT_4': FSO_HAM_OPT_4
    }

    return kpi_dictionary
