import json
import logging
import os
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import List

from openpyxl import load_workbook
from tzlocal import get_localzone
from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx.slide import Slide
from pptx.util import Inches, Pt

"""
- Add in some links

- KPIs
Agent versions
BTs
HRs
Errors
Policies and Actions

- Maturity Level 3
5% Platinum
10% Gold
25% Silver

- Maturity Level 4
10% Platinum
25% Gold
50% Silver
"""


class Color(Enum):
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    RED = RGBColor(255, 0, 0)
    GREEN = RGBColor(0, 255, 0)
    BLUE = RGBColor(0, 0, 255)


def setBackgroundImage(root: Presentation, slide: Slide, image_path: str):
    left = top = Inches(0)
    img_path = image_path
    pic = slide.shapes.add_picture(img_path, left, top, width=root.slide_width, height=root.slide_height)

    # Move it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def setTitle(slide: Slide, text: str, color: Color = Color.BLACK, fontSize: int = 32, top: float = 0.75, left: float = 0.5, width: float = 9):
    title = slide.shapes.title
    title.top = Inches(top)
    title.left = Inches(left)
    title.width = Inches(width)
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(fontSize)
    title.text_frame.paragraphs[0].font.color.rgb = color.value


def addBulletedText(slide: Slide, text: List[str], color: Color = Color.BLACK, fontSize: int = 12):
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    if len(text) == 0:
        return

    tf = body_shape.text_frame
    tf.text = text[0]

    for i in range(1, len(text)):
        p = tf.add_paragraph()
        p.text = text[i]

    for paragraph in body_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(fontSize)
            run.font.color.rgb = color.value


def addTable(slide, data, color: Color = Color.BLACK, fontSize: int = 16, left: int = 0.25, top: int = 3.5, width: int = 9.5, height: int = 1.5):
    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table

    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            for paragraph in table.cell(i, j).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fontSize)
                    run.font.color.rgb = color.value


def getValuesInColumn(sheet, param):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == param:
            j = 0
            for data in column_cell[1:]:
                values.append(data.value)
            break
    return values


def addNestedBulletedText(slide, text, color: Color = Color.BLACK, headerFontSize: int = 20, subheaderFontSize: int = 16):
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    if len(text) == 0:
        return

    tf = body_shape.text_frame
    tf.text = next(iter(text))

    firstIter = True
    for header in text.keys():
        if not firstIter:
            p = tf.add_paragraph()
            p.text = header
        else:
            firstIter = False
        for subheader in text[header]:
            p = tf.add_paragraph()
            p.text = subheader
            p.level = 1

    for paragraph in body_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(headerFontSize)
            run.font.color.rgb = color.value


def getAppsWithScore(sheet, assessmentScore):
    values = []
    for column_cell in sheet.iter_cols(1, sheet.max_column):
        if column_cell[0].value == "OverallAssessment":
            j = 0
            for idx, data in enumerate(column_cell[1:]):
                if data.value == assessmentScore:
                    values.append(sheet[f"C{idx + 2}"].value)  # +2 because of the header
            break
    return values


def createCxPpt(folder: str):
    logging.info(f"Creating presentation from output folder: {folder}")

    root = Presentation()

    """ Ref for slide types: 
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only 
    6 ->  Blank 
    7 ->  Content with caption
    8 ->  Pic with caption
    """

    # Title Slide
    slide = root.slides.add_slide(root.slide_layouts[0])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background.jpg")
    setTitle(slide, f"{folder} Configuration Assessment Highlights", Color.WHITE, fontSize=48, top=2.5)
    info = json.loads(open(f"output/{folder}/info.json").read())
    slide.placeholders[1].text = f'Data As Of: {datetime.fromtimestamp(info["lastRun"], get_localzone()).strftime("%m-%d-%Y at %H:%M:%S")}'

    # Current State Transition Slide
    slide = root.slides.add_slide(root.slide_layouts[5])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Current State", fontSize=48, top=3.5)

    # S/G/P Criteria & Scoring Slide
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"B/S/G/P Criteria & Scoring")
    text = [
        "Platinum - 50% of criteria for an application must be Platinum, all remaining criteria must be at least Gold",
        "Gold - 50% of criteria for an application must be Gold or higher, all remaining criteria must be at least Silver",
        "Silver - 80% of criteria for an application must be Silver or higher",
        "Bronze - All remaining applications",
    ]
    addBulletedText(slide, text)
    wb = load_workbook(filename=f"output/{folder}/{folder}-MaturityAssessment-apm.xlsx")
    totalApplications = wb["Analysis"].max_row - 1
    sheet = wb["Analysis"]
    scores = getValuesInColumn(sheet, "OverallAssessment")
    data = [
        ["Controller", "Apps", "Bronze% (#)", "Silver% (#)", "Gold% (#)", "Platinum% (#)"],
        [
            folder,
            str(totalApplications),
            f"{format(scores.count('bronze') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('bronze')})",
            f"{format(scores.count('silver') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('silver')})",
            f"{format(scores.count('gold') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('gold')})",
            f"{format(scores.count('platinum') / (wb['Analysis'].max_row - 1) * 100, '.0f')}% ({scores.count('platinum')})",
        ],
    ]
    addTable(slide, data)

    # App & Machine Agents
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"App & Machine Agents")
    text = [
        "AppD Agents are supported for 1 year after release",
        "*Machine Agents Reporting No Data includes nodes with an App Agent, but no Machine Agent installed",
    ]
    addBulletedText(slide, text)
    percentAgentsLessThan1YearOld = getValuesInColumn(wb["AppAgentsAPM"], "percentAgentsLessThan1YearOld")
    percentAgentsReportingData = getValuesInColumn(wb["AppAgentsAPM"], "percentAgentsReportingData")
    percentMachineAgentsLessThan1YearOld = getValuesInColumn(wb["MachineAgentsAPM"], "percentAgentsLessThan1YearOld")
    percentMachineAgentsReportingData = getValuesInColumn(wb["MachineAgentsAPM"], "percentAgentsReportingData")
    data = [
        [
            "Controller",
            "100% App Agents Older Than 1yr",
            "% App Agents Reporting No Data",
            "100% Machine Agents Older Than 1yr",
            "% Machine Agents Reporting No Data*",
        ],
        [
            folder,
            str(format(len([x for x in percentAgentsLessThan1YearOld if x != 100]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentAgentsReportingData if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentMachineAgentsLessThan1YearOld if x != 100]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in percentMachineAgentsReportingData if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Low Hanging Fruit Slide
    slide = root.slides.add_slide(root.slide_layouts[5])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Low Hanging Fruit", fontSize=48, top=3.5)

    # Overhead
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Overhead")
    text = [
        "Within the AppDynamics UI, we can enable the extraction of additional data",
        "While okay in lower environments, it is never recommended to use these features in a production environment",
        "Leaving some of these settings enabled can cause significant overhead",
    ]
    addBulletedText(slide, text)
    developerModeNotEnabledForAnyBT = getValuesInColumn(wb["OverheadAPM"], "developerModeNotEnabledForAnyBT")
    findEntryPointsNotEnabled = getValuesInColumn(wb["OverheadAPM"], "findEntryPointsNotEnabled")
    aggressiveSnapshottingNotEnabled = getValuesInColumn(wb["OverheadAPM"], "aggressiveSnapshottingNotEnabled")
    developerModeNotEnabledForApplication = getValuesInColumn(wb["OverheadAPM"], "developerModeNotEnabledForApplication")
    data = [
        [
            "Controller",
            "% Apps with BT Developer Mode Enabled",
            "% Apps with Find Entry Points Enabled",
            "% Apps with Aggressive Snapshotting Enabled",
            "% Apps with Developer Mode Enabled",
        ],
        [
            folder,
            str(format(len([x for x in developerModeNotEnabledForAnyBT if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in findEntryPointsNotEnabled if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in aggressiveSnapshottingNotEnabled if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in developerModeNotEnabledForApplication if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Error Configuration
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Error Configuration")
    text = [
        "It is critical to effectively configure Error detection settings ",
        "Failure to do so will lead to an inflated application error rate",
        "Ideally, there should be no BT with 100% error rate, and at least one custom error detection rule",
    ]
    addBulletedText(slide, text)
    successPercentageOfWorstTransaction = getValuesInColumn(wb["ErrorConfigurationAPM"], "successPercentageOfWorstTransaction")
    numberOfCustomRules = getValuesInColumn(wb["ErrorConfigurationAPM"], "numberOfCustomRules")
    data = [
        [
            "Controller",
            "% Apps w/BTs at 100% Error Rate",
            "% Apps without Custom Error Detection Rules",
        ],
        [
            folder,
            str(format(len([x for x in successPercentageOfWorstTransaction if x == 100]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in numberOfCustomRules if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Health Rules and Alerting
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Health Rules and Alerting")
    text = [
        "Too many Health Rule violations can lead to a 'boy who cried wolf' scenario. Consider tuning Health rules so they do not violate too frequently and become background noise,"
        "AppDynamics default Health Rules should always be modified to suit the needs of individual applications",
        "It is recommend to configure alerts to be sent to external tools",
        "Having no custom Health Rules is a telltale sign of poor application configuration health",
    ]
    addBulletedText(slide, text)
    numberOfHealthRuleViolationsLast24Hours = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfHealthRuleViolationsLast24Hours")
    numberOfDefaultHealthRulesModified = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfDefaultHealthRulesModified")
    numberOfActionsBoundToEnabledPolicies = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfActionsBoundToEnabledPolicies")
    numberOfCustomHealthRules = getValuesInColumn(wb["HealthRulesAndAlertingAPM"], "numberOfCustomHealthRules")
    data = [
        [
            "Controller",
            "% Apps with > 50 Health Rule Violations/Day",
            "% Apps without Default Health Rules Modified",
            "% Apps not Sending Alerts Anywhere",
            "% Apps without Custom Health Rules",
        ],
        [
            folder,
            str(format(len([x for x in numberOfHealthRuleViolationsLast24Hours if x >= 50]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in numberOfDefaultHealthRulesModified if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in numberOfActionsBoundToEnabledPolicies if x == 0]) / totalApplications * 100, ".0f")) + "%",
            str(format(len([x for x in numberOfCustomHealthRules if x == 0]) / totalApplications * 100, ".0f")) + "%",
        ],
    ]
    addTable(slide, data)

    # Recommendations Slide
    slide = root.slides.add_slide(root.slide_layouts[5])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Recommendations", fontSize=48, top=3.5)

    # Low-Hanging Fruit Slide with List
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Low-Hanging Fruit")
    text = {
        "Remove Overhead inducing settings": [
            "Can be done on controller without app team involvement",
            "Total time investment: 1-2 hours per controller",
        ],
        "Create Error Detection rules to reduce error rates": [
            "Requires app team involvement",
            "Total time investment: 1-2 hours per application",
        ],
        "Disable Health Rules which violate too frequently": [
            "Requires app team involvement",
            "Total time investment: 1-2 hours per application",
        ],
        "Add Policies and Actions to send off events": [
            "Requires app team involvement",
            "Total time investment: 1-2 hours per application",
        ],
    }
    addNestedBulletedText(slide, text)

    # Low-Hanging Fruit Slide with List
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Raise Gold Apps to Platinum Status")
    goldApps = getAppsWithScore(wb["Analysis"], "gold")
    text = {
        f"These apps are currently in Gold status. See {folder}-MaturityAssessment-apm.xlsx Analysis sheet for a full set of applications.": [],
        "We recommend working with them to raise them to Platinum status:": goldApps[:10],
    }
    addNestedBulletedText(slide, text)

    # Recommendations Slide
    slide = root.slides.add_slide(root.slide_layouts[5])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Appendix", fontSize=48, top=3.5)

    # Criteria Slide
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Configuration Assessment Tool Criteria")
    slide.shapes.add_picture("backend/resources/pptAssets/criteria.png", Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(5))

    # Criteria ctd... Slide
    slide = root.slides.add_slide(root.slide_layouts[1])
    setBackgroundImage(root, slide, "backend/resources/pptAssets/background_2.jpg")
    setTitle(slide, f"Configuration Assessment Tool Criteria ctd...")
    slide.shapes.add_picture("backend/resources/pptAssets/criteria2.png", Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(4))

    # Saving file
    root.save(f"output/{folder}/{folder}-cx-presentation.pptx")
