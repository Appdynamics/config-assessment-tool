import os
import json
import sys
import logging
import webbrowser
from flask import Flask, request, send_file, render_template
import openpyxl
from openpyxl.styles import PatternFill, Font
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl import load_workbook
import pandas as pd
from copy import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import xlwings as xw

app = Flask(__name__)

# Load configuration from JSON file
def load_config():
    # Determine the current script directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Construct the full path to the config.json file in the compare_results directory
    config_path = os.path.join(script_dir, 'config.json')  # Adjusted to directly access config.json

    if not os.path.exists(config_path):
        raise FileNotFoundError(f"Config file not found: {config_path}")

    with open(config_path) as f:
        config = json.load(f)
    
    # Convert relative paths to absolute paths
    config['upload_folder'] = os.path.join(script_dir, config.get('upload_folder', 'uploads'))
    config['result_folder'] = os.path.join(script_dir, config.get('result_folder', 'results'))
    
    return config

# Load the configuration
config = load_config()

# Configure upload and result folders from config
UPLOAD_FOLDER = config['upload_folder']
RESULT_FOLDER = config['result_folder']

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['RESULT_FOLDER'] = RESULT_FOLDER
logging.basicConfig(level=logging.DEBUG)

# Ensure folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['RESULT_FOLDER'], exist_ok=True)

# Add the controller check after the files are loaded and before processing
def check_controllers_match(previous_file_path, current_file_path):
    # Load the previous and current workbooks
    previous_workbook = pd.read_excel(previous_file_path, sheet_name='Analysis')
    current_workbook = pd.read_excel(current_file_path, sheet_name='Analysis')

    # Extract the 'controller' column from both workbooks
    previous_controllers = previous_workbook['controller'].unique()
    current_controllers = current_workbook['controller'].unique()

    # Check if the controllers match
    if not (len(previous_controllers) == 1 and len(current_controllers) == 1 and previous_controllers[0] == current_controllers[0]):
        print(f"Error: The controllers in the provided workbooks do not match.\n"
              f"Previous workbook controller: {previous_controllers[0]}\n"
              f"Current workbook controller: {current_controllers[0]}")
        sys.exit(1)  # Exit the script

# Function to generate PowerPoint slide
def generate_powerpoint_from_analysis(comparison_result_path, powerpoint_output_path):
    logging.debug("Generating PowerPoint presentation...")
    
    try:
        # Load the Analysis sheet from the Excel file
        df = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        logging.debug("Loaded Analysis sheet successfully.")
        logging.debug(f"DataFrame head:\n{df.head()}")  # Print the first few rows for inspection
        
        # Initialize counts
        results = {}
        columns = [
            'AppAgentsAPM', 'MachineAgentsAPM', 'BusinessTransactionsAPM',
            'BackendsAPM', 'OverheadAPM', 'ServiceEndpointsAPM',
            'ErrorConfigurationAPM', 'HealthRulesAndAlertingAPM',
            'DataCollectorsAPM', 'DashboardsAPM', 'OverallAssessment'
        ]
        
        total_applications = len(df)  # Total number of applications
        
        for col in columns:
            # Extract upgrade and downgrade information
            df[col] = df[col].astype(str)  # Ensure all data is treated as strings
            
            upgraded_count = df[col].str.contains('upgraded', case=False, na=False).sum()
            downgraded_count = df[col].str.contains('downgraded', case=False, na=False).sum()
            
            # Store results
            results[col] = {
                'upgraded': upgraded_count,
                'downgraded': downgraded_count
            }
            logging.debug(f"Column: {col}, Upgraded: {upgraded_count}, Downgraded: {downgraded_count}")
        
        # Create a PowerPoint presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        
        # Add title and subheading
        title_shape = shapes.title
        title_shape.text = "Assessment Comparison - Status"
        
        # Add a subheading for the total number of applications
        textbox = shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8), Inches(0.5))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f"Total number of Applications Compared: {total_applications}"
        p.font.size = Pt(14)  # Set the font size for the subheading
        
        # Center align the text in the textbox
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        textbox.left = Inches(0.5)  # Center the textbox horizontally
        textbox.width = Inches(8)  # Ensure it spans across the slide width
        
        # Add table to the slide
        rows = len(columns) + 1
        cols = 3  # Total, Upgraded, Downgraded columns
        table = shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(8), Inches(4)).table
        
        # Set table headers
        table.cell(0, 0).text = 'Metric'
        table.cell(0, 1).text = 'Maturity Increase'
        table.cell(0, 2).text = 'Maturity Decrease'
        
        # Set font size for headers
        for cell in table.rows[0].cells:
            cell.text_frame.paragraphs[0].font.size = Pt(14)
        
        # Populate the table
        for i, col in enumerate(columns, start=1):
            table.cell(i, 0).text = col
            table.cell(i, 1).text = f"{results[col]['upgraded'] / total_applications * 100:.2f}%"
            table.cell(i, 2).text = f"{results[col]['downgraded'] / total_applications * 100:.2f}%"
            
            # Set font size for data cells
            for cell in table.rows[i].cells:
                cell.text_frame.paragraphs[0].font.size = Pt(14)
        
        # Save the presentation
        prs.save(powerpoint_output_path)
        logging.debug(f"PowerPoint saved to {powerpoint_output_path}.")
    
    except Exception as e:
        logging.error(f"Error generating PowerPoint: {e}", exc_info=True)
        raise

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    logging.debug("Request files: %s", request.files)
    
    if 'previous_file' not in request.files or 'current_file' not in request.files:
        logging.error("No file part")
        return render_template('index.html', message="Error: No file part was uploaded."), 400
    
    previous_file = request.files.get('previous_file')
    current_file = request.files.get('current_file')
    
    if not previous_file or previous_file.filename == '':
        logging.error("No selected file for previous")
        return render_template('index.html', message="Error: No previous file was selected."), 400
    
    if not current_file or current_file.filename == '':
        logging.error("No selected file for current")
        return render_template('index.html', message="Error: No current file was selected."), 400
    
    previous_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous.xlsx')
    current_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'current.xlsx')
    output_file_path = os.path.join(app.config['RESULT_FOLDER'], 'comparison_result.xlsx')
    previous_sum_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous_sum.xlsx')
    current_sum_path = os.path.join(app.config['UPLOAD_FOLDER'], 'current_sum.xlsx')
    comparison_sum_path = os.path.join(app.config['RESULT_FOLDER'], 'comparison_sum.xlsx')
    powerpoint_output_path = os.path.join(app.config['RESULT_FOLDER'], 'Analysis_Summary.pptx')

    try:
        # Save the uploaded files
        previous_file.save(previous_file_path)
        current_file.save(current_file_path)

        # Load previous and current workbooks
        previous_workbook = pd.read_excel(previous_file_path, sheet_name='Analysis')
        current_workbook = pd.read_excel(current_file_path, sheet_name='Analysis')

        # Check if controllers are the same
        previous_controller = previous_workbook['controller'].unique()
        current_controller = current_workbook['controller'].unique()

        if len(previous_controller) != 1 or len(current_controller) != 1 or previous_controller[0] != current_controller[0]:
            logging.error("Controllers do not match.")
            return render_template('index.html', message="Error: The controllers in the two files do not match. Please upload files from the same controller."), 400

        # Create previous_sum and current_sum workbooks for Summary sheet
        create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path)

        # Perform the comparison only for Summary using previous_sum and current_sum
        compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path)

        # Perform comparison for other sheets
        compare_files_other_sheets(previous_file_path, current_file_path, output_file_path)

        # Copy the Summary sheet from comparison_sum to comparison_result
        copy_summary_to_result(comparison_sum_path, output_file_path)

        # Generate PowerPoint presentation from the Analysis sheet
        logging.debug("Generating PowerPoint presentation...")
        generate_powerpoint_from_analysis(output_file_path, powerpoint_output_path)

        # After successful processing, return a success message with a link to the result file
        return render_template(
            'index.html', 
            message=(
                f"Comparison completed successfully. "
                f"You can download the results from your browser by clicking "
                f"<a href='/download/{os.path.basename(output_file_path)}' style='color: #32CD32;'>here</a> "
                f"and the PowerPoint summary by clicking "
                f"<a href='/download/{os.path.basename(powerpoint_output_path)}'style='color: #32CD32;'>here</a>.<br>"
                f"<br>Click <a href='http://127.0.0.1:5000/' style='color: #32CD32;'>here</a> to return to the home page" 
            )
        )
    
    except Exception as e:
        logging.error(f"Error during file upload or comparison: {e}", exc_info=True)
        return render_template('index.html', message="Error during file upload or comparison"), 500


@app.route('/download/<filename>')
def download_file(filename):
    # Provide a download link for the output file
    return send_file(os.path.join(app.config['RESULT_FOLDER'], filename), as_attachment=True)


# Define color fills for Excel cells
red_fill = PatternFill(start_color='FF0000', end_color='FF0000', fill_type='solid')
green_fill = PatternFill(start_color='00FF00', end_color='00FF00', fill_type='solid')
added_fill = PatternFill(start_color='ADD8E6', end_color='ADD8E6', fill_type='solid')

# Helper function to find column index by name
def get_key_column(sheet, key_name):
    for i, cell in enumerate(sheet[1], 1):
        if cell.value == key_name:
            return i
    return None

# Function to create summary workbooks
def create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path):
    try:
        wb_previous = load_workbook(previous_file_path, data_only=True)
        wb_current = load_workbook(current_file_path, data_only=True)

        if 'Summary' not in wb_previous.sheetnames or 'Summary' not in wb_current.sheetnames:
            logging.error("'Summary' sheet is missing in one of the files.")
            return

        ws_previous = wb_previous['Summary']
        ws_current = wb_current['Summary']

        # Create new workbooks for the summaries
        wb_previous_sum = openpyxl.Workbook()
        wb_current_sum = openpyxl.Workbook()

        ws_previous_sum = wb_previous_sum.active
        ws_current_sum = wb_current_sum.active

        ws_previous_sum.title = 'Summary'
        ws_current_sum.title = 'Summary'

        # Copy data from original workbooks to summary workbooks as values only
        for row in ws_previous.iter_rows(values_only=True):
            ws_previous_sum.append(row)
        for row in ws_current.iter_rows(values_only=True):
            ws_current_sum.append(row)

        # Save the cleaned-up summary workbooks
        wb_previous_sum.save(previous_sum_path)
        wb_current_sum.save(current_sum_path)

    except Exception as e:
        logging.error(f"Error in create_summary_workbooks: {e}", exc_info=True)
        raise


# Function to compare 'Summary' sheet and save to a new workbook
def compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path):
    try:
        # Load the previous_sum and current_sum Excel files
        wb_previous = load_workbook(previous_sum_path, data_only=True)
        wb_current = load_workbook(current_sum_path, data_only=True)
        wb_output = openpyxl.Workbook()

        ws_previous = wb_previous['Summary']
        ws_current = wb_current['Summary']
        ws_output = wb_output.active
        ws_output.title = 'Summary'

        # Debugging: Print sheet names and some data from the 'Summary' sheet
        print("Previous Workbook Sheets:", wb_previous.sheetnames)
        print("Current Workbook Sheets:", wb_current.sheetnames)
        print("Example data from 'Summary' sheet in Previous Workbook (1,1):", ws_previous.cell(row=1, column=1).value)
        print("Example data from 'Summary' sheet in Current Workbook (1,1):", ws_current.cell(row=1, column=1).value)
        
        logging.debug(f"Processing sheet: 'Summary'")
        
        compare_summary(ws_previous, ws_current, ws_output)

        # Save the workbook after all modifications have been completed
        wb_output.save(comparison_sum_path)
        logging.debug(f"Summary comparison saved to: {comparison_sum_path}")

    except Exception as e:
        logging.error(f"Error in compare_files_summary: {e}", exc_info=True)
        raise


# Function to compare summaries of both sheets
def compare_summary(ws_previous, ws_current, ws_output):
    from openpyxl.styles import PatternFill

    # Define fill styles
    red_fill = PatternFill(start_color="FF0000", end_color="FF0000", fill_type="solid")
    green_fill = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")

    for row in ws_previous.iter_rows(min_row=1, min_col=1, max_col=ws_previous.max_column, max_row=ws_previous.max_row):
        for cell in row:
            prev_cell = ws_previous.cell(row=cell.row, column=cell.column)
            curr_cell = ws_current.cell(row=cell.row, column=cell.column)
            output_cell = ws_output.cell(row=cell.row, column=cell.column)

            prev_value = prev_cell.value
            curr_value = curr_cell.value

            if prev_value is None:
                prev_value = ''
            if curr_value is None:
                curr_value = ''

            logging.debug(f"Comparing cell ({cell.row},{cell.column}): Previous Value: {prev_value}, Current Value: {curr_value}")

            if prev_value != curr_value:
                if isinstance(prev_value, (int, float)) and isinstance(curr_value, (int, float)):
                    if curr_value > prev_value:
                        output_cell.fill = green_fill
                    else:
                        output_cell.fill = red_fill
                    output_cell.value = f"{prev_value} → {curr_value}"
                else:
                    output_cell.fill = red_fill
                    output_cell.value = f"{prev_value} → {curr_value}"
            else:
                output_cell.value = prev_value

            logging.debug(f"Cell ({cell.row},{cell.column}) updated to: {output_cell.value}")


def eval_formula(ws, cell):
    """ Helper function to evaluate a cell's formula and return the result. """
    try:
        if cell.data_type == 'f':
            # Create a temporary workbook to evaluate the formula
            temp_wb = openpyxl.Workbook()
            temp_ws = temp_wb.active
            temp_ws.cell(row=1, column=1).value = cell.value

            # Copy relevant data from the original sheet
            for row in range(1, ws.max_row + 1):
                for col in range(1, ws.max_column + 1):
                    temp_ws.cell(row=row + 1, column=col).value = ws.cell(row=row, column=col).value

            # Save and reopen to evaluate formulas
            temp_file = "temp.xlsx"
            temp_wb.save(temp_file)
            temp_wb.close()

            temp_wb = openpyxl.load_workbook(temp_file)
            temp_ws = temp_wb.active
            eval_result = temp_ws['A1'].value
            temp_wb.close()

            return eval_result
    except Exception as e:
        logging.error(f"Error evaluating formula in cell {cell.coordinate}: {e}", exc_info=True)
        return None

# Function to copy the Summary sheet from comparison_sum to comparison_result
def copy_summary_to_result(comparison_sum_path, output_file_path):
    try:
        # Load the comparison_sum and output workbooks
        wb_comparison_sum = load_workbook(comparison_sum_path)
        wb_output = load_workbook(output_file_path)

        # Get the Summary sheet from comparison_sum
        ws_comparison_sum = wb_comparison_sum['Summary']

        # If the Summary sheet already exists in output, delete it
        if 'Summary' in wb_output.sheetnames:
            del wb_output['Summary']

        # Create a new Summary sheet in the output workbook
        ws_output = wb_output.create_sheet('Summary', 0)  # Insert Summary as the first sheet

        # Copy data and formatting from the comparison_sum Summary sheet to the output Summary sheet
        for row in ws_comparison_sum.iter_rows():
            for cell in row:
                new_cell = ws_output.cell(row=cell.row, column=cell.column, value=cell.value)
                
                # Copy individual style attributes
                if cell.has_style:
                    new_cell.font = copy(cell.font)
                    new_cell.border = copy(cell.border)
                    new_cell.fill = copy(cell.fill)
                    new_cell.number_format = cell.number_format
                    new_cell.protection = copy(cell.protection)
                    new_cell.alignment = copy(cell.alignment)

        # Define header colors for Bronze, Silver, Gold, and Platinum
        header_colors = {
            'B1': 'cd7f32',  # Bronze
            'C1': 'c0c0c0',  # Silver
            'D1': 'ffd700',  # Gold
            'E1': 'e5e4e2'   # Platinum
        }

        # Apply colors to the headers in the first row
        for col in header_colors:
            cell = ws_output[col]
            cell.fill = PatternFill(start_color=header_colors[col], end_color=header_colors[col], fill_type="solid")
            cell.font = Font(bold=True, color="000000")

        # Save the output workbook
        wb_output.save(output_file_path)
        logging.debug("Summary sheet copied to the final comparison result and placed as the first sheet with highlighted headers.")

    except Exception as e:
        logging.error(f"Error in copy_summary_to_result: {e}", exc_info=True)
        raise

# Function to compare files for all other sheets
def compare_files_other_sheets(previous_file_path, current_file_path, output_file_path):
    try:
        wb_previous = load_workbook(previous_file_path)
        wb_current = load_workbook(current_file_path)

        for sheet_name in wb_current.sheetnames:
            if sheet_name in wb_previous.sheetnames:
                ws_previous = wb_previous[sheet_name]
                ws_current = wb_current[sheet_name]

                logging.debug(f"Processing sheet: {sheet_name}")

                if sheet_name == 'Analysis':
                    compare_analysis(ws_previous, ws_current)
                elif  sheet_name == 'AppAgentsAPM':
                    compare_appagentsapm(ws_previous, ws_current) 
                elif  sheet_name == 'MachineAgentsAPM':
                    compare_machineagentsapm(ws_previous, ws_current)
                elif sheet_name == 'BusinessTransactionsAPM':
                    compare_businesstransactionsapm(ws_previous, ws_current)
                elif  sheet_name == 'BackendsAPM':
                    compare_backendsapm(ws_previous, ws_current) 
                elif  sheet_name == 'OverheadAPM':
                    compare_overheadapm(ws_previous, ws_current)
                elif sheet_name == 'ServiceEndpointsAPM':
                    compare_serviceendpointsapm(ws_previous, ws_current)
                elif sheet_name == 'ErrorConfigurationAPM':
                    compare_errorconfigurationapm(ws_previous, ws_current)
                elif  sheet_name == 'HealthRulesAndAlertingAPM':
                    compare_healthrulesandalertingapm(ws_previous, ws_current) 
                elif  sheet_name == 'DataCollectorsAPM':
                    compare_datacollectorsapm(ws_previous, ws_current)
                elif sheet_name == 'DashboardsAPM':
                    compare_dashboardsapm(ws_previous, ws_current)
                elif sheet_name == 'OverallAssessmentAPM':
                    compare_overallassessmentapm(ws_previous, ws_current)
                elif sheet_name == 'Summary':
                    continue
                else:
                    logging.warning(f"No comparison function defined for sheet: {sheet_name}")

        wb_current.save(output_file_path)
        logging.debug(f"Comparison results saved to: {output_file_path}")

    except Exception as e:
        logging.error(f"Error in compare_files_other_sheets: {e}", exc_info=True)
        raise

def compare_appagentsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'metricLimitNotHit': None,
            'percentAgentsLessThan1YearOld': None,
            'percentAgentsLessThan2YearsOld': None,
            'percentAgentsReportingData': None,
            'percentAgentsRunningSameVersion': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    if column == 'metricLimitNotHit':
                        # Handle boolean logic for metricLimitNotHit
                        if current_value == 'TRUE' and previous_value == 'FALSE':
                            cell_output.fill = green_fill
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif current_value == 'FALSE' and previous_value == 'TRUE':
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"
                    else:
                        # Handle numeric logic for other columns
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_appagentsapm: {e}", exc_info=True)
        raise

def compare_machineagentsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'percentAgentsLessThan1YearOld': None,
            'percentAgentsLessThan2YearsOld': None,
            'percentAgentsReportingData': None,
            'percentAgentsRunningSameVersion': None,
            'percentAgentsInstalledAlongsideAppAgents': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    # Handle numeric logic for percentage columns
                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_machineagentsapm: {e}", exc_info=True)
        raise

def compare_datacollectorsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'numberOfDataCollectorFieldsConfigured': None,
            'numberOfDataCollectorFieldsCollectedInSnapshots': None,
            'numberOfDataCollectorFieldsCollectedInAnalytics': None,
            'biqEnabled': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                  #  logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        cell_output.value = previous_value  # Ensure the value is set to the previous value
                        continue  # Skip any further formatting or changes

                    if column == 'biqEnabled':
                        # Handle boolean logic for biqEnabled
                        if current_value == 'TRUE' and previous_value == 'FALSE':
                            cell_output.fill = green_fill
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif current_value == 'FALSE' and previous_value == 'TRUE':
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"
                    else:
                        # Handle numeric logic for other columns
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_datacollectorsapm: {e}", exc_info=True)
        raise

def compare_backendsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'percentBackendsWithLoad': None,
            'backendLimitNotHit': None,
            'numberOfCustomBackendRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        cell_output.value = previous_value  # Ensure the value is set to the previous value
                        continue  # Skip any further formatting or changes

                    if column == 'backendLimitNotHit':
                        # Handle boolean logic for backendLimitNotHit
                        if current_value == 'TRUE' and previous_value == 'FALSE':
                            cell_output.fill = green_fill
                            cell_output.value = f"{previous_value} → {current_value} (Increased)"
                        elif current_value == 'FALSE' and previous_value == 'TRUE':
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Decreased)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"
                    else:
                        # Handle numeric logic for other columns
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_backendsapm: {e}", exc_info=True)
        raise

def compare_overheadapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'developerModeNotEnabledForAnyBT': None,
            'findEntryPointsNotEnabled': None,
            'aggressiveSnapshottingNotEnabled': None,
            'developerModeNotEnabledForApplication': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                  #  logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    # Handle boolean logic for specified columns
                    if current_value == 'TRUE' and previous_value == 'FALSE':
                        cell_output.fill = green_fill
                        cell_output.value = f"{previous_value} → {current_value} (Improved)"
                    elif current_value == 'FALSE' and previous_value == 'TRUE':
                        cell_output.fill = red_fill
                        cell_output.value = f"{previous_value} → {current_value} (Declined)"
                    else:
                        cell_output.fill = red_fill
                        cell_output.value = f"{previous_value} → {current_value} (Changed)"

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_overheadapm: {e}", exc_info=True)
        raise

def compare_healthrulesandalertingapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'numberOfHealthRuleViolations': None,
            'numberOfDefaultHealthRulesModified': None,
            'numberOfActionsBoundToEnabledPolicies': None,
            'numberOfCustomHealthRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if column == 'numberOfHealthRuleViolations':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        else:
                            # For other columns
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_healthrulesandalertingapm: {e}", exc_info=True)
        raise

def compare_errorconfigurationapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'successPercentageOfWorstTransaction': None,
            'numberOfCustomRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if column == 'successPercentageOfWorstTransaction':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        elif column == 'numberOfCustomRules':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_errorconfigurationapm: {e}", exc_info=True)
        raise

def compare_serviceendpointsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfCustomServiceEndpointRules': None,
            'serviceEndpointLimitNotHit': None,
            'percentServiceEndpointsWithLoadOrDisabled': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        if column == 'numberOfCustomServiceEndpointRules':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                        elif column == 'serviceEndpointLimitNotHit':
                            if previous_value == "FALSE" and current_value == "TRUE":
                                cell_output.fill = green_fill
                                cell_output.value = "FALSE → TRUE"
                            elif previous_value == "TRUE" and current_value == "FALSE":
                                cell_output.fill = red_fill
                                cell_output.value = "TRUE → FALSE"

                        elif column == 'percentServiceEndpointsWithLoadOrDisabled':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_serviceendpointsapm: {e}", exc_info=True)
        raise

def compare_dashboardsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfDashboards': None,
            'percentageOfDashboardsModifiedLast6Months': None,
            'numberOfDashboardsUsingBiQ': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                        elif curr_value_num < prev_value_num:
                            cell_output.fill = red_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_dashboardsapm: {e}", exc_info=True)
        raise

def compare_overallassessmentapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'percentageTotalPlatinum': None,
            'percentageTotalGoldOrBetter': None,
            'percentageTotalSilverOrBetter': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                        elif curr_value_num < prev_value_num:
                            cell_output.fill = red_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_overallassessmentapm: {e}", exc_info=True)
        raise

# Function to compare 'Business Transactions' sheet
def compare_businesstransactionsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfBTs': None,
            'percentBTsWithLoad': None,
            'btLockdownEnabled': None,
            'numberCustomMatchRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                  #  logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        if column == 'numberOfBTs':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if 201 <= prev_value_num <= 600:
                                if curr_value_num < prev_value_num:
                                    cell_output.fill = green_fill
                                    cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                                elif curr_value_num > prev_value_num:
                                    cell_output.fill = red_fill
                                    cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"

                        elif column == 'percentBTsWithLoad':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                        elif column == 'btLockdownEnabled':
                            if previous_value == "FALSE" and current_value == "TRUE":
                                cell_output.fill = green_fill
                                cell_output.value = "FALSE → TRUE"
                            elif previous_value == "TRUE" and current_value == "FALSE":
                                cell_output.fill = red_fill
                                cell_output.value = "TRUE → FALSE"

                        elif column == 'numberCustomMatchRules':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_businesstransactionsapm: {e}", exc_info=True)
        raise 

def compare_analysis(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'AppAgentsAPM': None,
            'MachineAgentsAPM': None,
            'BusinessTransactionsAPM': None,
            'BackendsAPM': None,
            'OverheadAPM': None,
            'ServiceEndpointsAPM': None,
            'ErrorConfigurationAPM': None,
            'HealthRulesAndAlertingAPM': None,
            'DataCollectorsAPM': None,
            'DashboardsAPM': None,
            'OverallAssessment': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        name_col_prev = get_key_column(ws_previous, 'name')
        name_col_curr = get_key_column(ws_current, 'name')

        if name_col_prev is None or name_col_curr is None:
            logging.error("The 'name' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            name_value = row[name_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (name_value, ctrl_value)
            if name_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            name_value = row[name_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (name_value, ctrl_value)
            if name_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Comparison logic based on ranking
                        ranking = {'bronze': 1, 'silver': 2, 'gold': 3, 'platinum': 4}
                        prev_rank = ranking.get(previous_value.lower(), 0)
                        curr_rank = ranking.get(current_value.lower(), 0)

                        if curr_rank > prev_rank:
                            cell_output.fill = green_fill
                            cell_output.value = f"{previous_value} → {current_value} (Upgraded)"
                        elif curr_rank < prev_rank:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Downgraded)"
                    except ValueError:
                        logging.error(f"Invalid ranking value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_analysis: {e}", exc_info=True)
        raise

if __name__ == '__main__':
    # Automatically open the default web browser to the correct URL
    url = "http://127.0.0.1:5000"  # This is the default URL where Flask serves
    webbrowser.open(url)

    # Run the Flask app
    app.run(debug=True)
