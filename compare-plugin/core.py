

import os
import json
import sys
import logging
import webbrowser
from flask import Flask, request, send_file, render_template
import openpyxl
from openpyxl.styles import PatternFill, Font
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl import load_workbook
import pandas as pd
from copy import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.placeholder import TablePlaceholder
import xlwings as xw

# Set up logging configuration
def setup_logging():
    # Set log level to INFO directly
    log_level = logging.INFO  # INFO, WARNING, ERROR, CRITICAL
    logging.basicConfig(
        level=log_level,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    logging.info("Logging is set up!")

# Call the setup_logging function to initialize logging
setup_logging()

app = Flask(__name__)

def save_workbook(filepath):
    """Open and save the workbook to ensure formulas are recalculated."""
    # Open the workbook
    app = xw.App(visible=False)  # Set visible=False to avoid showing the Excel window
    wb = app.books.open(filepath)
    
    # Save the workbook
    wb.save()
    
    # Close the workbook
    wb.close()
    
    # Quit the application
    app.quit()

def process_files(previous_file_path, current_file_path):
    # Save the workbooks to trigger formula recalculation
    save_workbook(previous_file_path)
    save_workbook(current_file_path)
    
    # Continue with your comparison logic
    compare_files_summary(previous_file_path, current_file_path, 'comparison_summary.xlsx')

# Load configuration from JSON file
def load_config():
    # Determine the current script directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Construct the full path to the config.json file
    config_path = os.path.join(script_dir, 'config.json')

    if not os.path.exists(config_path):
        raise FileNotFoundError(f"Config file not found: {config_path}")

    with open(config_path) as f:
        config = json.load(f)
    
    # Convert relative paths to absolute paths
    config['upload_folder'] = os.path.join(script_dir, config.get('upload_folder', 'uploads'))
    config['result_folder'] = os.path.join(script_dir, config.get('result_folder', 'results'))
    config['TEMPLATE_FOLDER'] = os.path.join(script_dir, config.get('TEMPLATE_FOLDER', 'templates'))  # Adding this line
    
    return config

# Load the configuration
config = load_config()

# Configure upload and result folders from config
UPLOAD_FOLDER = config['upload_folder']
RESULT_FOLDER = config['result_folder']

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['RESULT_FOLDER'] = RESULT_FOLDER
logging.basicConfig(level=logging.DEBUG)

# Ensure folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['RESULT_FOLDER'], exist_ok=True)

def check_controllers_match(previous_file_path, current_file_path):
    # Load the previous and current workbooks
    previous_workbook = pd.read_excel(previous_file_path, sheet_name='Analysis')
    current_workbook = pd.read_excel(current_file_path, sheet_name='Analysis')

    # Extract the 'controller' column from both workbooks and strip whitespaces
    previous_controllers = previous_workbook['controller'].dropna().str.strip().unique()
    current_controllers = current_workbook['controller'].dropna().str.strip().unique()

    # Log the exact controller values for debugging
    logging.debug(f"Previous controller(s): {previous_controllers}")
    logging.debug(f"Current controller(s): {current_controllers}")

    # Check if the controllers match
    if not (len(previous_controllers) == 1 and len(current_controllers) == 1 and previous_controllers[0] == current_controllers[0]):
        logging.error(f"Controllers do not match. Previous controller: {previous_controllers}, Current controller: {current_controllers}")
        return False

    return True

def generate_powerpoint_from_analysis(comparison_result_path, powerpoint_output_path, current_file_path, previous_file_path):
    logging.debug("Generating PowerPoint presentation...")

    try:
        # prs = Presentation(template_path)  # Open the template

        # Define the relative path for the template using the TEMPLATE_FOLDER
        template_folder = config.get('TEMPLATE_FOLDER', 'templates')  # 'templates' is the default folder name
        template_path = os.path.join(template_folder, 'template.pptx')

        # Load the 'Analysis' sheet from the current workbook (uploaded by the user)
        df_current_analysis = pd.read_excel(current_file_path, sheet_name='Analysis')

        # Count the number of valid applications by counting rows where 'name' column is not NaN or empty
        number_of_apps = df_current_analysis['name'].dropna().str.strip().ne('').sum()

        # Log the number of valid applications
        logging.info(f"Number of applications in the current 'Analysis' sheet: {number_of_apps}")

        # Check if the template exists, otherwise, ask the user for input or use environment variables
        if not os.path.exists(template_path):
            template_path = os.getenv('TEMPLATE_PATH', template_path)  # Allow user to set this via an environment variable
            if not os.path.exists(template_path):
                template_path = input("Template not found! Please provide the full path to the template: ")

        # Load the template
        prs = Presentation(template_path)
        logging.debug(f"Template loaded from: {template_path}")

        # Load the Summary sheet
        summary_df = pd.read_excel(comparison_result_path, sheet_name='Summary')
        logging.debug("Loaded Summary sheet successfully.")
        logging.debug(f"Summary DataFrame head:\n{summary_df.head()}")

        # Load the Analysis sheet
        df_analysis = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        # Load the 'AppAgentsAPM' sheet from the Excel file
        df_app_agents = pd.read_excel(comparison_result_path, sheet_name='AppAgentsAPM')
        # Load the 'MachineAgentsAPM' sheet from the Excel file
        df_machine_agents = pd.read_excel(comparison_result_path, sheet_name='MachineAgentsAPM')
        # Load the 'BusinessTransactionsAPM' sheet from the Excel file
        df_BTs = pd.read_excel(comparison_result_path, sheet_name='BusinessTransactionsAPM')
        # Load the 'BackendsAPM' sheet from the Excel file
        df_Backends = pd.read_excel(comparison_result_path, sheet_name='BackendsAPM')
        # Load the 'OverheadAPM' sheet from the Excel file
        df_Overhead = pd.read_excel(comparison_result_path, sheet_name='OverheadAPM')
        # Load the 'ServiceEndpointsAPM' sheet from the Excel file
        df_ServiceEndpoints = pd.read_excel(comparison_result_path, sheet_name='ServiceEndpointsAPM')
        # Load the 'ErrorConfigurationAPM' sheet from the Excel file
        df_ErrorConfiguration = pd.read_excel(comparison_result_path, sheet_name='ErrorConfigurationAPM')
        # Load the 'HealthRulesAndAlertingAPM' sheet from the Excel file
        df_HealthRulesAndAlerting = pd.read_excel(comparison_result_path, sheet_name='HealthRulesAndAlertingAPM')
        # Load the 'DataCollectorsAPM' sheet from the Excel file
        df_DataCollectors = pd.read_excel(comparison_result_path, sheet_name='DataCollectorsAPM')
        # Load the 'DashboardsAPM' sheet from the Excel file
        df_Dashboards = pd.read_excel(comparison_result_path, sheet_name='DashboardsAPM')

        # Function to find table placeholders by name
        def find_table_placeholder_by_name(slide, name):
            for shape in slide.shapes:
                if shape.is_placeholder and shape.name == name:
                    return shape
            return None  # Return None if not found

        def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
            """Insert a table at the position of a placeholder."""
            placeholder = find_table_placeholder_by_name(slide, placeholder_name)
            
            if not placeholder:
                logging.error(f"Placeholder '{placeholder_name}' not found on the slide.")
                return None

            # Get placeholder dimensions
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            logging.debug(f"Inserting table at placeholder position: left={left}, top={top}, width={width}, height={height}")

            # Insert table at the placeholder's position
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            return table_shape.table  # Return the inserted table
        
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Now handle Slide 3 table with "Upgraded" applications **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        slide = prs.slides[2]  # Slide 3 (index 2)
        upgraded_apps = df_analysis[df_analysis['OverallAssessment'].str.contains('upgraded', case=False, na=False)]['name'].tolist()

        # Count the number of applications in the current "Analysis" sheet
        current_analysis_df = pd.read_excel(current_file_path, sheet_name='Analysis')  # Load the current "Analysis" sheet
        number_of_apps = len(current_analysis_df)

        # Insert the count into TextBox 7
        textbox_7 = None
        for shape in slide.shapes:
            if shape.name == "TextBox 7":
                textbox_7 = shape
                break

        if textbox_7:
            textbox_7.text = f"{number_of_apps}"  # Set the text with the count
        else:
            logging.warning("TextBox 8 not found on Slide 3.")

        # Insert Upgraded Applications Table onto Slide 3 (Slide index 2) - using Table Placeholder 1
        upgraded_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # We are now using the same placeholder
        if upgraded_placeholder:
            logging.debug("Found Upgraded Applications table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(upgraded_apps) + 1, 1)
        else:
            logging.warning("Upgraded Applications table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(upgraded_apps) + 1, 1, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Add header for the new table
        table.cell(0, 0).text = "Applications with Upgraded Metrics"
        table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with upgraded applications
        for idx, app in enumerate(upgraded_apps):
            table.cell(idx + 1, 0).text = app
            table.cell(idx + 1, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table onto Slide 4 (Slide index 3) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[3]  # Slide 4 (index 3)
        summary_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder for Summary Table

        if summary_placeholder:
            logging.debug("Found Summary table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(summary_df) + 1, len(summary_df.columns))
        else:
            logging.warning("Summary table placeholder not found. Adding manually.")
            # Explicitly add a new table with defined dimensions for Slide 4
            table = slide.shapes.add_table(len(summary_df) + 1, len(summary_df.columns), Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Set column headers for the Summary table
        for col_idx, column in enumerate(summary_df.columns):
            table.cell(0, col_idx).text = str(column)
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate table with Summary data
        for row_idx, row in summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)
                table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Add the title for Slide 4 (Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Comparison Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white


        # Load the Analysis sheet
        df = pd.read_excel(comparison_result_path, sheet_name='Analysis')

        columns = [
            'AppAgentsAPM', 'MachineAgentsAPM', 'BusinessTransactionsAPM',
            'BackendsAPM', 'OverheadAPM', 'ServiceEndpointsAPM',
            'ErrorConfigurationAPM', 'HealthRulesAndAlertingAPM',
            'DataCollectorsAPM', 'DashboardsAPM', 'OverallAssessment'
        ]

        results = {}
        total_applications = len(df)
        
        for col in columns:
            df[col] = df[col].astype(str)
            upgraded_count = df[col].str.contains('upgraded', case=False, na=False).sum()
            downgraded_count = df[col].str.contains('downgraded', case=False, na=False).sum()

            # Total applications is the length of the column
            total_applications = len(df[col])

            overall_result = "Increase" if upgraded_count > downgraded_count else "Decrease" if downgraded_count > upgraded_count else "Even"
            percentage_value = 0 if overall_result == "Even" else round((upgraded_count / total_applications) * 100)

            # Log the results for each column
            # logging.debug(f"Column: {col}")
            # logging.debug(f"Upgraded Count: {upgraded_count}")
            # logging.debug(f"Total Applications: {total_applications}")
            # logging.debug(f"Overall Result: {overall_result}")
            # logging.debug(f"Percentage: {percentage_value}%")

            results[col] = {
                'upgraded': upgraded_count,
                'downgraded': downgraded_count,
                'overall_result': overall_result,
                'percentage': percentage_value
            }

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table from Previous Workbook onto Slide 4 (Table Placeholder 4) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        slide = prs.slides[3]  # Slide 4 (index 3)

        # Load the previous summary data
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Add to Table Placeholder 4 (for previous summary)
        summary_placeholder_previous = find_table_placeholder_by_name(slide, "Table Placeholder 4")  # Placeholder for Previous Summary Table
        if summary_placeholder_previous:
            logging.debug("Found Table Placeholder 4. Inserting table for previous summary.")
            table_previous = insert_table_at_placeholder(slide, "Table Placeholder 4", len(previous_summary_df) + 1, len(previous_summary_df.columns))
        else:
            logging.warning("Table Placeholder 4 not found. Adding manually.")
            table_previous = slide.shapes.add_table(len(previous_summary_df) + 1, len(previous_summary_df.columns), Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Set column headers for the previous summary table
        for col_idx, column in enumerate(previous_summary_df.columns):
            table_previous.cell(0, col_idx).text = str(column)
            table_previous.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with previous summary data
        for row_idx, row in previous_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_previous.cell(row_idx + 1, col_idx).text = str(value)
                table_previous.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table from Current Workbook onto Slide 4 (Table Placeholder 3) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # Load the current summary data
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')

        # Add to Table Placeholder 3 (for current summary)
        summary_placeholder_current = find_table_placeholder_by_name(slide, "Table Placeholder 3")  # Placeholder for Current Summary Table
        if summary_placeholder_current:
            logging.debug("Found Table Placeholder 3. Inserting table for current summary.")
            table_current = insert_table_at_placeholder(slide, "Table Placeholder 3", len(current_summary_df) + 1, len(current_summary_df.columns))
        else:
            logging.warning("Table Placeholder 3 not found. Adding manually.")
            table_current = slide.shapes.add_table(len(current_summary_df) + 1, len(current_summary_df.columns), Inches(0.5), Inches(6), Inches(9), Inches(4)).table  

        # Set column headers for the current summary table
        for col_idx, column in enumerate(current_summary_df.columns):
            table_current.cell(0, col_idx).text = str(column)
            table_current.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with current summary data
        for row_idx, row in current_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_current.cell(row_idx + 1, col_idx).text = str(value)
                table_current.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Add the title for Slide 4 (Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Comparison Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white


        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        # Insert Overall Assessment Table onto Slide 5 (Slide index 4)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[4]  # Slide 5 (index 4)
        overall_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder name

        if overall_placeholder:
            # logging.debug("Found Overall Assessment table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", 2, 5)
        else:
            # logging.warning("Overall Assessment table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(2, 5, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).table  

        headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        overall_assessment = results['OverallAssessment']
        table.cell(1, 0).text = 'OverallAssessment'
        table.cell(1, 1).text = str(overall_assessment['upgraded'])
        table.cell(1, 2).text = str(overall_assessment['downgraded'])
        table.cell(1, 3).text = overall_assessment['overall_result']
        table.cell(1, 4).text = f"{overall_assessment['percentage']}%"

        if overall_assessment['overall_result'] == "Increase":
            table.cell(1, 4).fill.solid()
            table.cell(1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Add the title for Slide 4
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Overall Assessment Result"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        # Insert Status Table onto Slide 6 (Slide index 5)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[5]  # Slide 6 (index 5)
        status_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder name

        num_rows = len(columns)  # Should match the expected row count
        num_cols = 5  

        if status_placeholder:
            # logging.debug("Found Status table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", num_rows, num_cols)
        else:
            # logging.warning("Status table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        for i, col in enumerate(columns[:-1]):  
            table.cell(i + 1, 0).text = col
            table.cell(i + 1, 1).text = str(results[col]['upgraded'])
            table.cell(i + 1, 2).text = str(results[col]['downgraded'])
            table.cell(i + 1, 3).text = results[col]['overall_result']
            table.cell(i + 1, 4).text = f"{results[col]['percentage']}%"

            if results[col]['overall_result'] == "Increase":
                table.cell(i + 1, 4).fill.solid()
                table.cell(i + 1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Add the title for Slide 5
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "APM Maturity Assessment Result"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert APM AGENT Downgrade Table onto Slide 10 (Slide index 9) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[9]  # Slide 10 (index 9)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the AppAgentsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['AppAgentsAPM']
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 10 (Slide index 9)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")

            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 10 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "APM Agent - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentAgentsLessThan1YearOld': 'Rectangle 11',
            'metricLimitNotHit': 'Rectangle 10',
            'percentAgentsLessThan2YearsOld': 'Rectangle 12',
            'percentAgentsReportingData': 'Rectangle 13',
            'percentAgentsRunningSameVersion': 'Rectangle 14'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_app_agents.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'declined' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 10
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 10.") 


        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert MACHINE AGENT Downgrade Table onto Slide 11 (Slide index 10) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[10]  # Slide 11 (index 10)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the MachineAgentsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['MachineAgentsAPM']  # Use MachineAgentsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 11 (Slide index 10)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 11 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Machine Agent - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentAgentsLessThan1YearOld': 'Rectangle 8',
            'percentAgentsLessThan2YearsOld': 'Rectangle 9',
            'percentAgentsReportingData': 'Rectangle 10',
            'percentAgentsRunningSameVersion': 'Rectangle 11',
            'percentAgentsInstalledAlongsideAppAgents': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_machine_agents.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'declined' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 10
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 11.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert BT Downgrade Table onto Slide 12 (Slide index 11) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[11]  # Slide 12 (index 11)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the BusinessTransactionsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['BusinessTransactionsAPM']  # Use BusinessTransactionsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 12 (Slide index 11)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 12 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Business Transactions - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfBTs': 'Rectangle 17',
            'percentBTsWithLoad': 'Rectangle 18',
            'btLockdownEnabled': 'Rectangle 19',
            'numberCustomMatchRules': 'Rectangle 20'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_BTs.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 10
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 11.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Backend Downgrade Table onto Slide 13 (Slide index 12) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[12]  # Slide 13 (index 12)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the BackendsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['BackendsAPM']  # Use BackendsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 13 (Slide index 12)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            # Convert all items in the applications list to strings
            applications_str = ', '.join(str(app) for app in downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            
            # Log the grade and applications
            logging.debug(f"Grade: {grade}, Applications: {applications_str}")
            
            # Populate the table with the data
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = applications_str  # Display the application names
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 13 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Backends - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentBackendsWithLoad': 'Rectangle 10',
            'backendLimitNotHit': 'Rectangle 11',
            'numberOfCustomBackendRules': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Backends.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 10
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 11.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert SEP Downgrade Table onto Slide 14 (Slide index 13) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[13]  # Slide 14 (index 13)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the ServiceEndpointsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['ServiceEndpointsAPM']  # Use ServiceEndpointsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 14 (Slide index 13)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 14 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Service Endpoints - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfCustomServiceEndpointRules': 'Rectangle 10',
            'serviceEndpointLimitNotHit': 'Rectangle 11',
            'percentServiceEndpointsWithLoadOrDisabled': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_ServiceEndpoints.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 14
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 14.") 

        
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert ERROR CONFIG Downgrade Table onto Slide 15 (Slide index 14) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[14]  # Slide 15 (index 14)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the ErrorConfigurationAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['ErrorConfigurationAPM']  # Use ErrorConfigurationAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 15 (Slide index 14)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 15 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Error Configuration - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'successPercentageOfWorstTransaction': 'Rectangle 10',
            'numberOfCustomRules': 'Rectangle 11'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_ErrorConfiguration.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 15
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert HR & ALERTS Downgrade Table onto Slide 16 (Slide index 15) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[15]  # Slide 16 (index 15)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the HealthRulesAndAlertingAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['HealthRulesAndAlertingAPM']  # Use HealthRulesAndAlertingAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 16 (Slide index 15)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 16 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Health Rules & Alerting - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfHealthRuleViolations': 'Rectangle 10',
            'numberOfDefaultHealthRulesModified': 'Rectangle 11',
            'numberOfActionsBoundToEnabledPolicies': 'Rectangle 12',
            'numberOfCustomHealthRules': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_HealthRulesAndAlerting.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 16
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 16.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert DATA COLLECTORS Downgrade Table onto Slide 17 (Slide index 16) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[16]  # Slide 17 (index 16)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the DataCollectorsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['DataCollectorsAPM']  # Use DataCollectorsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 17 (Slide index 16)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 17 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Data Collectors - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfDataCollectorFieldsConfigured': 'Rectangle 10',
            'numberOfDataCollectorFieldsCollectedInSnapshots': 'Rectangle 11',
            'numberOfDataCollectorFieldsCollectedInAnalytics': 'Rectangle 12',
            'biqEnabled': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_DataCollectors.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 17
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 17.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert DASHBOARDS Downgrade Table onto Slide 18 (Slide index 17) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[17]  # Slide 18 (index 17)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the DashboardsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['DashboardsAPM']  # Use DashboardsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 18 (Slide index 17)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 18 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Dashboards - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

                # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfDashboards': 'Rectangle 10',
            'percentageOfDashboardsModifiedLast6Months': 'Rectangle 11',
            'numberOfDashboardsUsingBiQ': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Dashboards.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 18
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 18.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert OVERHEAD Downgrade Table onto Slide 19 (Slide index 18) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[18]  # Slide 19 (index 18)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the OverheadAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['OverheadAPM']  # Use OverheadAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 19 (Slide index 18)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 19 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Overhead - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'developerModeNotEnabledForAnyBT': 'Rectangle 10',
            'findEntryPointsNotEnabled': 'Rectangle 11',
            'aggressiveSnapshottingNotEnabled': 'Rectangle 12',
            'developerModeNotEnabledForApplication': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Overhead.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'changed' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 19
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 19.")

        # Save the PowerPoint
        prs.save(powerpoint_output_path)
        logging.debug(f"PowerPoint saved to {powerpoint_output_path}.")

    except Exception as e:
        logging.error(f"Error generating PowerPoint: {e}", exc_info=True)
        raise




 

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    logging.debug("Request files: %s", request.files)
    
    if 'previous_file' not in request.files or 'current_file' not in request.files:
        logging.error("No file part")
        return render_template('index.html', message="Error: No file part was uploaded."), 400
    
    previous_file = request.files.get('previous_file')
    current_file = request.files.get('current_file')
    
    if not previous_file or previous_file.filename == '':
        logging.error("No selected file for previous")
        return render_template('index.html', message="Error: No previous file was selected."), 400
    
    if not current_file or current_file.filename == '':
        logging.error("No selected file for current")
        return render_template('index.html', message="Error: No current file was selected."), 400
    
    previous_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous.xlsx')
    current_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'current.xlsx')
    output_file_path = os.path.join(app.config['RESULT_FOLDER'], 'comparison_result.xlsx')
    previous_sum_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous_sum.xlsx')
    current_sum_path = os.path.join(app.config['UPLOAD_FOLDER'], 'current_sum.xlsx')
    comparison_sum_path = os.path.join(app.config['RESULT_FOLDER'], 'comparison_sum.xlsx')
    powerpoint_output_path = os.path.join(app.config['RESULT_FOLDER'], 'Analysis_Summary.pptx')


    try:
        # Save the uploaded files
        previous_file.save(previous_file_path)
        current_file.save(current_file_path)

        # Automatically save the workbooks to recalculate formulas
        save_workbook(previous_file_path)
        save_workbook(current_file_path)

        # Check if controllers match
        if not check_controllers_match(previous_file_path, current_file_path):
            logging.error("Controllers do not match.")
            return render_template('index.html', message="Error: The controllers in the two files do not match. Please upload files from the same controller."), 400

        # Proceed with comparison...
        create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path)
        compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path)
        compare_files_other_sheets(previous_file_path, current_file_path, output_file_path)
        copy_summary_to_result(comparison_sum_path, output_file_path)
        generate_powerpoint_from_analysis(output_file_path, powerpoint_output_path, current_file_path, previous_file_path)

        return render_template(
            'index.html', 
            message=(
                f"Comparison completed successfully. "
                f"You can download the results from your browser by clicking "
                f"<a href='/download/{os.path.basename(output_file_path)}' style='color: #32CD32;'>here</a> "
                f"and the PowerPoint summary by clicking "
                f"<a href='/download/{os.path.basename(powerpoint_output_path)}'style='color: #32CD32;'>here</a>.<br>"
                f"<br>Click <a href='http://127.0.0.1:5000/' style='color: #32CD32;'>here</a> to return to the home page"
            )
        )
    
    except Exception as e:
        logging.error(f"Error during file upload or comparison: {e}", exc_info=True)
        return render_template('index.html', message="Error during file upload or comparison"), 500

@app.route('/download/<filename>')
def download_file(filename):
    # Provide a download link for the output file
    return send_file(os.path.join(app.config['RESULT_FOLDER'], filename), as_attachment=True)


# Define color fills for Excel cells
red_fill = PatternFill(start_color='FF0000', end_color='FF0000', fill_type='solid')
green_fill = PatternFill(start_color='00FF00', end_color='00FF00', fill_type='solid')
added_fill = PatternFill(start_color='ADD8E6', end_color='ADD8E6', fill_type='solid')

# Helper function to find column index by name
def get_key_column(sheet, key_name):
    for i, cell in enumerate(sheet[1], 1):
        if cell.value == key_name:
            return i
    return None

# Function to create summary workbooks
def create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path):
    try:
        wb_previous = load_workbook(previous_file_path, data_only=True)
        wb_current = load_workbook(current_file_path, data_only=True)

        if 'Summary' not in wb_previous.sheetnames or 'Summary' not in wb_current.sheetnames:
            logging.error("'Summary' sheet is missing in one of the files.")
            return

        ws_previous = wb_previous['Summary']
        ws_current = wb_current['Summary']

        # Create new workbooks for the summaries
        wb_previous_sum = openpyxl.Workbook()
        wb_current_sum = openpyxl.Workbook()

        ws_previous_sum = wb_previous_sum.active
        ws_current_sum = wb_current_sum.active

        ws_previous_sum.title = 'Summary'
        ws_current_sum.title = 'Summary'

        # Copy data from original workbooks to summary workbooks as values only
        for row in ws_previous.iter_rows(values_only=True):
            ws_previous_sum.append(row)
        for row in ws_current.iter_rows(values_only=True):
            ws_current_sum.append(row)

        # Save the cleaned-up summary workbooks
        wb_previous_sum.save(previous_sum_path)
        wb_current_sum.save(current_sum_path)

    except Exception as e:
        logging.error(f"Error in create_summary_workbooks: {e}", exc_info=True)
        raise


# Function to compare 'Summary' sheet and save to a new workbook
def compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path):
    try:
        # Load the previous_sum and current_sum Excel files
        wb_previous = load_workbook(previous_sum_path, data_only=True)
        wb_current = load_workbook(current_sum_path, data_only=True)
        wb_output = openpyxl.Workbook()

        ws_previous = wb_previous['Summary']
        ws_current = wb_current['Summary']
        ws_output = wb_output.active
        ws_output.title = 'Summary'

        # Debugging: Print sheet names and some data from the 'Summary' sheet
        print("Previous Workbook Sheets:", wb_previous.sheetnames)
        print("Current Workbook Sheets:", wb_current.sheetnames)
        print("Example data from 'Summary' sheet in Previous Workbook (1,1):", ws_previous.cell(row=1, column=1).value)
        print("Example data from 'Summary' sheet in Current Workbook (1,1):", ws_current.cell(row=1, column=1).value)
        
        logging.debug(f"Processing sheet: 'Summary'")
        
        compare_summary(ws_previous, ws_current, ws_output)

        # Save the workbook after all modifications have been completed
        wb_output.save(comparison_sum_path)
        logging.debug(f"Summary comparison saved to: {comparison_sum_path}")

    except Exception as e:
        logging.error(f"Error in compare_files_summary: {e}", exc_info=True)
        raise


# Function to compare summaries of both sheets
def compare_summary(ws_previous, ws_current, ws_output):
    from openpyxl.styles import PatternFill

    # Define fill styles
    red_fill = PatternFill(start_color="FF0000", end_color="FF0000", fill_type="solid")
    green_fill = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")

    for row in ws_previous.iter_rows(min_row=1, min_col=1, max_col=ws_previous.max_column, max_row=ws_previous.max_row):
        for cell in row:
            prev_cell = ws_previous.cell(row=cell.row, column=cell.column)
            curr_cell = ws_current.cell(row=cell.row, column=cell.column)
            output_cell = ws_output.cell(row=cell.row, column=cell.column)

            prev_value = prev_cell.value
            curr_value = curr_cell.value

            if prev_value is None:
                prev_value = ''
            if curr_value is None:
                curr_value = ''

            logging.debug(f"Comparing cell ({cell.row},{cell.column}): Previous Value: {prev_value}, Current Value: {curr_value}")

            if prev_value != curr_value:
                if isinstance(prev_value, (int, float)) and isinstance(curr_value, (int, float)):
                    if curr_value > prev_value:
                        output_cell.fill = green_fill
                    else:
                        output_cell.fill = red_fill
                    output_cell.value = f"{prev_value} → {curr_value}"
                else:
                    output_cell.fill = red_fill
                    output_cell.value = f"{prev_value} → {curr_value}"
            else:
                output_cell.value = prev_value

            logging.debug(f"Cell ({cell.row},{cell.column}) updated to: {output_cell.value}")


def eval_formula(ws, cell):
    """ Helper function to evaluate a cell's formula and return the result. """
    try:
        if cell.data_type == 'f':
            # Create a temporary workbook to evaluate the formula
            temp_wb = openpyxl.Workbook()
            temp_ws = temp_wb.active
            temp_ws.cell(row=1, column=1).value = cell.value

            # Copy relevant data from the original sheet
            for row in range(1, ws.max_row + 1):
                for col in range(1, ws.max_column + 1):
                    temp_ws.cell(row=row + 1, column=col).value = ws.cell(row=row, column=col).value

            # Save and reopen to evaluate formulas
            temp_file = "temp.xlsx"
            temp_wb.save(temp_file)
            temp_wb.close()

            temp_wb = openpyxl.load_workbook(temp_file)
            temp_ws = temp_wb.active
            eval_result = temp_ws['A1'].value
            temp_wb.close()

            return eval_result
    except Exception as e:
        logging.error(f"Error evaluating formula in cell {cell.coordinate}: {e}", exc_info=True)
        return None

# Function to copy the Summary sheet from comparison_sum to comparison_result
def copy_summary_to_result(comparison_sum_path, output_file_path):
    try:
        # Load the comparison_sum and output workbooks
        wb_comparison_sum = load_workbook(comparison_sum_path)
        wb_output = load_workbook(output_file_path)

        # Get the Summary sheet from comparison_sum
        ws_comparison_sum = wb_comparison_sum['Summary']

        # If the Summary sheet already exists in output, delete it
        if 'Summary' in wb_output.sheetnames:
            del wb_output['Summary']

        # Create a new Summary sheet in the output workbook
        ws_output = wb_output.create_sheet('Summary', 0)  # Insert Summary as the first sheet

        # Copy data and formatting from the comparison_sum Summary sheet to the output Summary sheet
        for row in ws_comparison_sum.iter_rows():
            for cell in row:
                new_cell = ws_output.cell(row=cell.row, column=cell.column, value=cell.value)
                
                # Copy individual style attributes
                if cell.has_style:
                    new_cell.font = copy(cell.font)
                    new_cell.border = copy(cell.border)
                    new_cell.fill = copy(cell.fill)
                    new_cell.number_format = cell.number_format
                    new_cell.protection = copy(cell.protection)
                    new_cell.alignment = copy(cell.alignment)

        # Define header colors for Bronze, Silver, Gold, and Platinum
        header_colors = {
            'B1': 'cd7f32',  # Bronze
            'C1': 'c0c0c0',  # Silver
            'D1': 'ffd700',  # Gold
            'E1': 'e5e4e2'   # Platinum
        }

        # Apply colors to the headers in the first row
        for col in header_colors:
            cell = ws_output[col]
            cell.fill = PatternFill(start_color=header_colors[col], end_color=header_colors[col], fill_type="solid")
            cell.font = Font(bold=True, color="000000")

        # Save the output workbook
        wb_output.save(output_file_path)
        logging.debug("Summary sheet copied to the final comparison result and placed as the first sheet with highlighted headers.")

    except Exception as e:
        logging.error(f"Error in copy_summary_to_result: {e}", exc_info=True)
        raise

# Function to compare files for all other sheets
def compare_files_other_sheets(previous_file_path, current_file_path, output_file_path):
    try:
        wb_previous = load_workbook(previous_file_path)
        wb_current = load_workbook(current_file_path)

        for sheet_name in wb_current.sheetnames:
            if sheet_name in wb_previous.sheetnames:
                ws_previous = wb_previous[sheet_name]
                ws_current = wb_current[sheet_name]

                logging.debug(f"Processing sheet: {sheet_name}")

                if sheet_name == 'Analysis':
                    compare_analysis(ws_previous, ws_current)
                elif  sheet_name == 'AppAgentsAPM':
                    compare_appagentsapm(ws_previous, ws_current) 
                elif  sheet_name == 'MachineAgentsAPM':
                    compare_machineagentsapm(ws_previous, ws_current)
                elif sheet_name == 'BusinessTransactionsAPM':
                    compare_businesstransactionsapm(ws_previous, ws_current)
                elif  sheet_name == 'BackendsAPM':
                    compare_backendsapm(ws_previous, ws_current) 
                elif  sheet_name == 'OverheadAPM':
                    compare_overheadapm(ws_previous, ws_current)
                elif sheet_name == 'ServiceEndpointsAPM':
                    compare_serviceendpointsapm(ws_previous, ws_current)
                elif sheet_name == 'ErrorConfigurationAPM':
                    compare_errorconfigurationapm(ws_previous, ws_current)
                elif  sheet_name == 'HealthRulesAndAlertingAPM':
                    compare_healthrulesandalertingapm(ws_previous, ws_current) 
                elif  sheet_name == 'DataCollectorsAPM':
                    compare_datacollectorsapm(ws_previous, ws_current)
                elif sheet_name == 'DashboardsAPM':
                    compare_dashboardsapm(ws_previous, ws_current)
                elif sheet_name == 'OverallAssessmentAPM':
                    compare_overallassessmentapm(ws_previous, ws_current)
                elif sheet_name == 'Summary':
                    continue
                else:
                    logging.warning(f"No comparison function defined for sheet: {sheet_name}")

        wb_current.save(output_file_path)
        logging.info(f"Comparison results saved to: {output_file_path}")

    except Exception as e:
        logging.error(f"Error in compare_files_other_sheets: {e}", exc_info=True)
        raise

def compare_appagentsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'metricLimitNotHit': None,
            'percentAgentsLessThan1YearOld': None,
            'percentAgentsLessThan2YearsOld': None,
            'percentAgentsReportingData': None,
            'percentAgentsRunningSameVersion': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        continue

                    if column == 'metricLimitNotHit':
                        # Handle boolean logic for metricLimitNotHit (True/False)
                        prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                        curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                        logging.info(f"Comparing {column}: Previous={prev_value_str}, Current={curr_value_str}")

                        # Compare "TRUE" vs "FALSE"
                        if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                            cell_output.fill = green_fill  # Green for improvement (False → True)
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                            cell_output.fill = red_fill  # Red for decline (True → False)
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"

                    # Handle numeric logic for other columns
                    else:
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_appagentsapm: {e}", exc_info=True)
        raise
 

def compare_machineagentsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'percentAgentsLessThan1YearOld': None,
            'percentAgentsLessThan2YearsOld': None,
            'percentAgentsReportingData': None,
            'percentAgentsRunningSameVersion': None,
            'percentAgentsInstalledAlongsideAppAgents': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    # Handle numeric logic for percentage columns
                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_machineagentsapm: {e}", exc_info=True)
        raise

def compare_datacollectorsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'numberOfDataCollectorFieldsConfigured': None,
            'numberOfDataCollectorFieldsCollectedInSnapshots': None,
            'numberOfDataCollectorFieldsCollectedInAnalytics': None,
            'biqEnabled': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        cell_output.value = previous_value  # Ensure the value is set to the previous value
                        continue  # Skip any further formatting or changes

                    if column == 'biqEnabled':
                        # Handle boolean logic for biqEnabled
                        prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                        curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                        # Log the comparison for debugging
                        logging.info(f"Comparing {column}: Previous={prev_value_str}, Current={curr_value_str}")

                        # Compare "TRUE" vs "FALSE"
                        if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                            cell_output.fill = green_fill  # Green for improvement (False → True)
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                            cell_output.fill = red_fill  # Red for decline (True → False)
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            # For other cases, we just mark it as changed
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_datacollectorsapm: {e}", exc_info=True)
        raise


# Function to compare 'Backends' sheet 
def compare_backendsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'percentBackendsWithLoad': None,
            'backendLimitNotHit': None,  # Column to compare
            'numberOfCustomBackendRules': None
        }

        # Retrieve column indices for the columns to be compared
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices (application and controller columns)
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        cell_output.value = previous_value  # Ensure the value is set to the previous value
                        continue  # Skip any further formatting or changes

                    if column == 'backendLimitNotHit':
                        # Handle boolean logic for backendLimitNotHit
                        prev_value = str(previous_value).strip().upper()  # Convert to string and handle case insensitivity
                        curr_value = str(current_value).strip().upper()  # Same for the current value

                        # Log the comparison for debugging
                        logging.info(f"Comparing backendLimitNotHit: Previous={prev_value}, Current={curr_value}")

                        # Compare "FALSE" vs "TRUE"
                        if prev_value == "FALSE" and curr_value == "TRUE":
                            cell_output.fill = green_fill  # Correct color for FALSE → TRUE
                            cell_output.value = "FALSE → TRUE"
                        elif prev_value == "TRUE" and curr_value == "FALSE":
                            cell_output.fill = red_fill  # Correct color for TRUE → FALSE
                            cell_output.value = "TRUE → FALSE"
                        else:
                            cell_output.fill = red_fill  # Default for unexpected values
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"
                    else:
                        # Handle numeric logic for other columns
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_backendsapm: {e}", exc_info=True)
        raise

def compare_overheadapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'developerModeNotEnabledForAnyBT': None,
            'findEntryPointsNotEnabled': None,
            'aggressiveSnapshottingNotEnabled': None,
            'developerModeNotEnabledForApplication': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    # Handle boolean logic for specified columns
                    if column == 'developerModeNotEnabledForAnyBT' or column == 'findEntryPointsNotEnabled' or column == 'aggressiveSnapshottingNotEnabled' or column == 'developerModeNotEnabledForApplication':
                        prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                        curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                        # Log the comparison for debugging
                        logging.info(f"Comparing {column}: Previous={prev_value_str}, Current={curr_value_str}")

                        # Compare "TRUE" vs "FALSE"
                        if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                            cell_output.fill = green_fill  # Green for True (Improvement)
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                            cell_output.fill = red_fill  # Red for False (Decline)
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            # For other cases, we just mark it as changed
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_overheadapm: {e}", exc_info=True)
        raise


def compare_healthrulesandalertingapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'numberOfHealthRuleViolations': None,
            'numberOfDefaultHealthRulesModified': None,
            'numberOfActionsBoundToEnabledPolicies': None,
            'numberOfCustomHealthRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if column == 'numberOfHealthRuleViolations':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        else:
                            # For other columns
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_healthrulesandalertingapm: {e}", exc_info=True)
        raise

def compare_errorconfigurationapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'successPercentageOfWorstTransaction': None,
            'numberOfCustomRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if column == 'successPercentageOfWorstTransaction':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        elif column == 'numberOfCustomRules':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_errorconfigurationapm: {e}", exc_info=True)
        raise

# Function to compare 'Service Endpoints' sheet
def compare_serviceendpointsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfCustomServiceEndpointRules': None,
            'serviceEndpointLimitNotHit': None,
            'percentServiceEndpointsWithLoadOrDisabled': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        if column == 'numberOfCustomServiceEndpointRules':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                        elif column == 'serviceEndpointLimitNotHit':
                            # Explicitly handle TRUE/FALSE as strings and booleans
                            prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                            curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                            logging.info(f"Comparing serviceEndpointLimitNotHit: Previous={prev_value_str}, Current={curr_value_str}")

                            # Compare "TRUE" vs "FALSE"
                            if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                                cell_output.fill = green_fill
                                cell_output.value = "FALSE → TRUE"
                            elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                                cell_output.fill = red_fill
                                cell_output.value = "TRUE → FALSE"
                            else:
                                # Log if we encounter an unexpected value
                                logging.info(f"Unexpected values for serviceEndpointLimitNotHit: Previous={previous_value}, Current={current_value}")

                        elif column == 'percentServiceEndpointsWithLoadOrDisabled':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_serviceendpointsapm: {e}", exc_info=True)
        raise


def compare_dashboardsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfDashboards': None,
            'percentageOfDashboardsModifiedLast6Months': None,
            'numberOfDashboardsUsingBiQ': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                        elif curr_value_num < prev_value_num:
                            cell_output.fill = red_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_dashboardsapm: {e}", exc_info=True)
        raise

def compare_overallassessmentapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'percentageTotalPlatinum': None,
            'percentageTotalGoldOrBetter': None,
            'percentageTotalSilverOrBetter': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                        elif curr_value_num < prev_value_num:
                            cell_output.fill = red_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_overallassessmentapm: {e}", exc_info=True)
        raise

# Function to compare 'Business Transactions' sheet
def compare_businesstransactionsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfBTs': None,
            'percentBTsWithLoad': None,
            'btLockdownEnabled': None,
            'numberCustomMatchRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                print(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            print("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            print("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        print(f"Column: {column}, Previous Value: {previous_value}, Current Value: {current_value}")

                        # Handle each column's specific comparison logic
                        if column == 'numberOfBTs':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if 201 <= prev_value_num <= 600:
                                if curr_value_num < prev_value_num:
                                    cell_output.fill = green_fill
                                    cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                                elif curr_value_num > prev_value_num:
                                    cell_output.fill = red_fill
                                    cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"

                        elif column == 'percentBTsWithLoad':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                        elif column == 'btLockdownEnabled':
                            # Normalize TRUE/FALSE as strings
                            prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                            curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                            # Log the comparison for debugging
                            print(f"Comparing btLockdownEnabled for app {key}: Previous={prev_value_str}, Current={curr_value_str}")

                            # Compare "TRUE" vs "FALSE"
                            if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                                print(f"Update: FALSE → TRUE for app {key}")
                                cell_output.fill = green_fill
                                cell_output.value = "FALSE → TRUE"
                            elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                                print(f"Update: TRUE → FALSE for app {key}")
                                cell_output.fill = red_fill
                                cell_output.value = "TRUE → FALSE"
                            else:
                                print(f"Unexpected values for btLockdownEnabled: Previous={previous_value}, Current={current_value}")

                        elif column == 'numberCustomMatchRules':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        print(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        print(f"Error in compare_businesstransactionsapm: {e}")
        raise


def compare_analysis(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'AppAgentsAPM': None,
            'MachineAgentsAPM': None,
            'BusinessTransactionsAPM': None,
            'BackendsAPM': None,
            'OverheadAPM': None,
            'ServiceEndpointsAPM': None,
            'ErrorConfigurationAPM': None,
            'HealthRulesAndAlertingAPM': None,
            'DataCollectorsAPM': None,
            'DashboardsAPM': None,
            'OverallAssessment': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        name_col_prev = get_key_column(ws_previous, 'name')
        name_col_curr = get_key_column(ws_current, 'name')

        if name_col_prev is None or name_col_curr is None:
            logging.error("The 'name' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            name_value = row[name_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (name_value, ctrl_value)
            if name_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            name_value = row[name_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (name_value, ctrl_value)
            if name_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Comparison logic based on ranking
                        ranking = {'bronze': 1, 'silver': 2, 'gold': 3, 'platinum': 4}
                        prev_rank = ranking.get(previous_value.lower(), 0)
                        curr_rank = ranking.get(current_value.lower(), 0)

                        if curr_rank > prev_rank:
                            cell_output.fill = green_fill
                            cell_output.value = f"{previous_value} → {current_value} (Upgraded)"
                        elif curr_rank < prev_rank:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Downgraded)"
                    except ValueError:
                        logging.error(f"Invalid ranking value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_analysis: {e}", exc_info=True)
        raise

if __name__ == '__main__':
    # Automatically open the default web browser to the correct URL
    url = "http://127.0.0.1:5000"  # This is the default URL where Flask serves
    webbrowser.open(url)

    # Run the Flask app
    app.run(debug=True)
