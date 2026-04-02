# compare_tool/powerpoint/__init__.py

import logging
import os
from typing import Any, Dict, Optional, List

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)

# NEW: import the specialised generators
try:
    from .brum import generate_powerpoint_from_brum
except ImportError:
    generate_powerpoint_from_brum = None  # type: ignore

try:
    from .mrum import generate_powerpoint_from_mrum
except ImportError:
    generate_powerpoint_from_mrum = None  # type: ignore

# Optional: if you have a dedicated APM generator in apm.py
try:
    from .apm import generate_powerpoint_from_apm
except ImportError:
    generate_powerpoint_from_apm = None  # type: ignore



# ---------------------------------------------------------------------------
# Helper: load the Analysis sheet
# ---------------------------------------------------------------------------
def _load_analysis(comparison_result_path: str) -> pd.DataFrame:
    """
    Load the 'Analysis' sheet from the comparison_result workbook.
    """
    log.info("Loading Analysis sheet from %s", comparison_result_path)
    return pd.read_excel(comparison_result_path, sheet_name="Analysis")


# ---------------------------------------------------------------------------
# Helper: overall upgrade/downgrade counts from OverallAssessment
# ---------------------------------------------------------------------------
def _count_overall(df: pd.DataFrame) -> Dict[str, Any]:
    """
    Compute how many rows are Upgraded / Downgraded based on OverallAssessment.
    """
    if "OverallAssessment" not in df.columns:
        return {"improved": 0, "degraded": 0, "result": "Even", "percentage": 0}

    s = df["OverallAssessment"].astype(str)
    improved = int(s.str.contains("Upgraded", case=False, na=False).sum())
    degraded = int(s.str.contains("Downgraded", case=False, na=False).sum())

    if improved > degraded:
        result = "Increase"
    elif degraded > improved:
        result = "Decrease"
    else:
        result = "Even"

    pct = 0 if result == "Even" else round(
        (improved / max(1, improved + degraded)) * 100
    )

    return {
        "improved": improved,
        "degraded": degraded,
        "result": result,
        "percentage": pct,
    }


# ---------------------------------------------------------------------------
# Helper: pick template path based on config + domain
# ---------------------------------------------------------------------------
def _pick_template_path(config: Optional[Dict[str, Any]], domain: str) -> Optional[str]:
    """
    Very simple template resolution:
      1) config["TEMPLATE_APM"/"TEMPLATE_BRUM"/"TEMPLATE_MRUM"] if present
      2) <TEMPLATE_FOLDER>/template_brum.pptx, template_mrum.pptx, template.pptx
      3) <TEMPLATE_FOLDER>/template.pptx as a generic fallback
    """
    cfg = config or {}
    domain = (domain or "APM").upper()

    key_by_domain = {
        "APM": "TEMPLATE_APM",
        "BRUM": "TEMPLATE_BRUM",
        "MRUM": "TEMPLATE_MRUM",
    }
    dom_key = key_by_domain.get(domain)
    if dom_key and cfg.get(dom_key):
        path = cfg[dom_key]
        if os.path.exists(path):
            return path

    template_folder = cfg.get("TEMPLATE_FOLDER", "templates")

    domain_name = {
        "APM": "template.pptx",          # legacy APM template name
        "BRUM": "template_brum.pptx",
        "MRUM": "template_mrum.pptx",
    }.get(domain, "template.pptx")

    candidate = os.path.join(template_folder, domain_name)
    if os.path.exists(candidate):
        return candidate

    generic = os.path.join(template_folder, "template.pptx")
    return generic if os.path.exists(generic) else None


# ---------------------------------------------------------------------------
# Helper: find a reasonable "name" column for apps
# ---------------------------------------------------------------------------
def _find_name_column(df: pd.DataFrame) -> Optional[str]:
    """
    Try to detect which column holds the application name.
    Works with 'name', 'Name', 'application', 'Application', etc.
    """
    candidates = ["name", "Name", "application", "Application", "ApplicationName"]
    for col in candidates:
        if col in df.columns:
            return col
    return None


# ---------------------------------------------------------------------------
# Helper: build rows for the detail slide
# ---------------------------------------------------------------------------
def _build_detail_rows(df: pd.DataFrame) -> List[str]:
    """
    Build human-readable lines for each application for the detail slide.
    The idea is to make BRUM/MRUM obviously show real data.
    """
    if df.empty:
        return []

    name_col = _find_name_column(df)
    if not name_col:
        # Fallback: just show row index and OverallAssessment
        name_series = df.index.astype(str)
    else:
        name_series = df[name_col].astype(str)

    overall_series = (
        df["OverallAssessment"].astype(str)
        if "OverallAssessment" in df.columns
        else pd.Series([""] * len(df))
    )

    # Mark which rows have any "change"
    has_change = overall_series.str.contains(
        "Upgraded|Downgraded", case=False, na=False
    )

    # We want changed apps first, then the rest
    df_detail = pd.DataFrame(
        {
            "name": name_series,
            "overall": overall_series,
            "changed": has_change,
        }
    )

    # Sort: changed first, then alphabetical by name
    df_detail = df_detail.sort_values(
        by=["changed", "name"], ascending=[False, True]
    )

    lines: List[str] = []
    for _, row in df_detail.iterrows():
        name = row["name"].strip() or "(Unnamed App)"
        overall = row["overall"].strip()

        if overall:
            line = f"{name} – {overall}"
        else:
            line = name

        lines.append(line)

    # Don't go totally wild – cap to, say, 40 entries
    return lines[:40]


# ---------------------------------------------------------------------------
# Helper: ensure at least one slide exists and return slide 0
# ---------------------------------------------------------------------------
def _ensure_first_slide(prs: Presentation) -> Any:
    """
    Make sure there is at least one slide in the deck and return it.
    """
    if not prs.slides:
        layout = prs.slide_layouts[0]
        return prs.slides.add_slide(layout)
    return prs.slides[0]


# ---------------------------------------------------------------------------
# Helper: add or update a "detail" slide (2nd slide) with app list
# ---------------------------------------------------------------------------
def _add_detail_slide(prs: Presentation, df: pd.DataFrame, domain: str) -> None:
    """
    Create / update a second slide with a list of applications and their
    OverallAssessment so that BRUM/MRUM decks clearly show live data.
    """
    domain = (domain or "APM").upper()

    # Use existing 2nd slide if present, otherwise add a new one
    if len(prs.slides) > 1:
        slide = prs.slides[1]
    else:
        # Prefer "Title + Content" layout if available
        layout_index = 1 if len(prs.slide_layouts) > 1 else 0
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title is not None:
        slide.shapes.title.text = f"{domain} – Application Detail"
    else:
        # Add a title textbox if needed
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_tf = title_box.text_frame
        title_tf.text = f"{domain} – Application Detail"
        title_tf.paragraphs[0].font.size = Pt(28)

    # Main body textbox for the list
    left = Inches(0.6)
    top = Inches(1.3)
    width = Inches(9.0)
    height = Inches(5.0)

    # Remove any existing text boxes in the main content area? We'll keep it simple:
    body_box = slide.shapes.add_textbox(left, top, width, height)
    tf = body_box.text_frame
    tf.clear()

    lines = _build_detail_rows(df)
    if not lines:
        p = tf.paragraphs[0]
        p.text = "No application rows found in Analysis sheet."
        p.font.size = Pt(18)
        return

    # First line
    first = tf.paragraphs[0]
    first.text = lines[0]
    first.font.size = Pt(18)

    # Remaining lines as bullets
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)


# ---------------------------------------------------------------------------
# Main entry point: generate_powerpoint_from_analysis
# ---------------------------------------------------------------------------
def generate_powerpoint_from_analysis(*args, **kwargs) -> str:
    """
    Domain-aware, but generic PowerPoint generator.

    It:
      * Reads the Analysis sheet from comparison_result_path.
      * Computes total apps and upgrade/downgrade stats.
      * Populates slide 1 with a clear summary.
      * Populates slide 2 with a per-application detail list
        (very useful to make BRUM/MRUM obviously show data).

    It is forgiving about how it's called (positional vs keyword args) so it
    works both with older and newer service.run_comparison implementations.
    """

    # ---- Normalise arguments (supports both old & new call styles) ---------
    comparison_result_path = kwargs.pop("comparison_result_path", None)
    powerpoint_output_path = kwargs.pop("powerpoint_output_path", None)
    current_file_path = kwargs.pop("current_file_path", None)    # not used yet
    previous_file_path = kwargs.pop("previous_file_path", None)  # not used yet
    domain = kwargs.pop("domain", "APM")
    config = kwargs.pop("config", None)

    # Allow positional usage: (comparison_result_path, powerpoint_output_path, ...)
    if comparison_result_path is None and len(args) > 0:
        comparison_result_path = args[0]
    if powerpoint_output_path is None and len(args) > 1:
        powerpoint_output_path = args[1]
    if current_file_path is None and len(args) > 2:
        current_file_path = args[2]
    if previous_file_path is None and len(args) > 3:
        previous_file_path = args[3]
    if len(args) > 4 and isinstance(args[4], str):
        domain = args[4]
    if len(args) > 5 and isinstance(args[5], dict):
        config = args[5]

    if comparison_result_path is None or powerpoint_output_path is None:
        raise ValueError(
            "comparison_result_path and powerpoint_output_path are required"
        )

    domain = (domain or "APM").upper()

    log.info(
        "Generating PPT (generic): domain=%s, comparison_result=%s, output=%s",
        domain,
        comparison_result_path,
        powerpoint_output_path,
    )

    # ---- Load analysis and calculate numbers -------------------------------
    df_analysis = _load_analysis(comparison_result_path)

    # Count apps (try to be robust to column naming)
    name_col = _find_name_column(df_analysis)
    if name_col:
        total_apps = int(
            df_analysis[name_col]
            .dropna()
            .astype(str)
            .str.strip()
            .ne("")
            .sum()
        )
    else:
        total_apps = len(df_analysis)

    overall = _count_overall(df_analysis)

    log.info(
        "PPT stats: total_apps=%s, improved=%s, degraded=%s, result=%s, pct=%s",
        total_apps,
        overall["improved"],
        overall["degraded"],
        overall["result"],
        overall["percentage"],
    )

    # ---- Load template (or blank deck) ------------------------------------
    template_path = _pick_template_path(config, domain)
    if template_path and os.path.exists(template_path):
        log.info("Using PPT template: %s", template_path)
        prs = Presentation(template_path)
    else:
        log.warning(
            "No PPT template found (looked for %s). Using a blank presentation.",
            template_path,
        )
        prs = Presentation()

    # ---- Slide 1: summary -------------------------------------------------
    slide = _ensure_first_slide(prs)

    # Update title if possible
    if slide.shapes.title is not None:
        slide.shapes.title.text = f"{domain} Maturity Comparison – Summary"

    # Big textbox with the key stats
    left = Inches(0.75)
    top = Inches(2.0)
    width = Inches(9)
    height = Inches(4)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.text = f"Total applications analysed: {total_apps}"
    p0.font.size = Pt(24)

    p1 = tf.add_paragraph()
    p1.text = (
        f"Improved (OverallAssessment contains 'Upgraded'): "
        f"{overall['improved']}"
    )
    p1.level = 0
    p1.font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = (
        f"Degraded (OverallAssessment contains 'Downgraded'): "
        f"{overall['degraded']}"
    )
    p2.level = 0
    p2.font.size = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = (
        f"Net result: {overall['result']} "
        f"({overall['percentage']}% of changed apps)"
    )
    p3.level = 0
    p3.font.size = Pt(20)

    # ---- Slide 2: application detail list ---------------------------------
    _add_detail_slide(prs, df_analysis, domain)

    # ---- Save and return path ---------------------------------------------
    os.makedirs(os.path.dirname(powerpoint_output_path), exist_ok=True)
    prs.save(powerpoint_output_path)
    log.info("PPT saved to %s", powerpoint_output_path)

    return powerpoint_output_path

# ---------------------------------------------------------------------------
# Unified entry point: generate_powerpoint
# ---------------------------------------------------------------------------
def generate_powerpoint(*args, **kwargs) -> str:
    """
    Unified entry point for PPT generation.

    Dispatch rules:
      - domain == 'BRUM' -> generate_powerpoint_from_brum(...)
      - domain == 'MRUM' -> generate_powerpoint_from_mrum(...)
      - domain == 'APM'  -> generate_powerpoint_from_apm(...) if available,
                            otherwise falls back to generate_powerpoint_from_analysis(...)
      - anything else     -> generate_powerpoint_from_analysis(...)
    
    It supports both positional and keyword arguments in the same way
    as generate_powerpoint_from_analysis.
    """

    # Extract / normalise domain from kwargs or 5th positional argument
    domain = kwargs.get("domain")
    if domain is None and len(args) > 4 and isinstance(args[4], str):
        domain = args[4]

    domain_norm = (domain or "APM").upper()

    # BRUM
    if domain_norm == "BRUM" and generate_powerpoint_from_brum is not None:
        log.info("Dispatching to BRUM generator")
        return generate_powerpoint_from_brum(*args, **kwargs)

    # MRUM
    if domain_norm == "MRUM" and generate_powerpoint_from_mrum is not None:
        log.info("Dispatching to MRUM generator")
        return generate_powerpoint_from_mrum(*args, **kwargs)

    # APM – if you have a dedicated APM generator, prefer that
    if domain_norm == "APM" and generate_powerpoint_from_apm is not None:
        log.info("Dispatching to APM generator")
        return generate_powerpoint_from_apm(*args, **kwargs)

    # Fallback: your existing generic implementation
    log.info(
        "Dispatching to generic Analysis-based PPT generator (domain=%s)",
        domain_norm,
    )
    # Ensure the domain kwarg matches the normalised value for consistent behaviour
    kwargs["domain"] = domain_norm
    return generate_powerpoint_from_analysis(*args, **kwargs)

__all__ = [
    "generate_powerpoint",
    "generate_powerpoint_from_analysis",
    "generate_powerpoint_from_brum",
    "generate_powerpoint_from_mrum",
]
