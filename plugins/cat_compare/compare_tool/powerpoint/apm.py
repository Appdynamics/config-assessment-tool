"""
apm.py
------
This module handles the generation of PowerPoint presentations for APM (Application Performance Management) comparisons.

Purpose:
- Creates PowerPoint slides specific to APM comparison results.
- Uses templates to dynamically populate slides with APM-related data.

Key Features:
- Reads APM comparison data and generates slides with metrics, charts, and summaries.
- Saves the generated PowerPoint presentation to the specified output path.
"""

import logging
import os
import re

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

log = logging.getLogger(__name__)
log.info("[APM] apm.py imported")

# Nite: Optional config import – if your app has a central config, use it.
try:
    # Adjust to your real config location if different
    from ..config import config  # type: ignore
except Exception:  # pragma: no cover
    config = {}


# Nite: Import shared PPTX utilities from a common base module.
# These helpers are reused across APM / BRUM / MRUM.
try:
    # Adjust path if your base module lives somewhere else
    from .base import (
        autosize_col_to_header,
        set_arrow_cell,
        PINK,
        overall_maturity_from_df,
        color_oval_for_maturity,
    )
except Exception as e:  # pragma: no cover
    logging.warning(
        "Could not import shared PPT helpers from powerpoint.base. "
        "Make sure autosize_col_to_header, set_arrow_cell, PINK, "
        "overall_maturity_from_df, color_oval_for_maturity are available. "
        f"Error: {e}"
    )

    # Optional: define super-simple fallbacks so module still imports.
    def autosize_col_to_header(*args, **kwargs):
        pass

    def set_arrow_cell(cell, value, color=None, size_pt=36):
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(size_pt)

    PINK = RGBColor(255, 105, 180)

    def overall_maturity_from_df(df, grade_func=None):
        return None

    def color_oval_for_maturity(slide, shape_name, tier, update_text=False):
        pass


# ---------------------------------------------------------------------------
# Nite: Generic placeholder helpers used across many slides
# ---------------------------------------------------------------------------

def find_table_placeholder_by_name(slide, name):
    """
    Nite: Find a placeholder shape on the given slide by its name.

    Returns the shape if found, otherwise None.
    """
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and shape.name == name:
            return shape
    return None


def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
    """
    Nite: Insert a table at the position/dimensions of a named placeholder.

    Returns: the pptx.table.Table instance, or None if placeholder not found.
    """
    placeholder = find_table_placeholder_by_name(slide, placeholder_name)

    if not placeholder:
        logging.error(f"Placeholder '{placeholder_name}' not found on the slide.")
        return None

    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    logging.debug(
        f"Inserting table at placeholder position: "
        f"left={left}, top={top}, width={width}, height={height}"
    )

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


# ---------------------------------------------------------------------------
# Nite: Main entry point for APM PPT generation
# ---------------------------------------------------------------------------

def generate_powerpoint_from_apm(
    comparison_result_path,
    powerpoint_output_path,
    current_file_path,
    previous_file_path,
    template_path=None,
    domain="APM",
    config=None,
):
    logging.debug("Generating PowerPoint presentation (APM)...")
    logging.warning(">>> ENTERED generate_powerpoint_from_apm() <<<")

    try:
        # -------------------------------------------------------------------
        # Resolve template path
        # -------------------------------------------------------------------
        # Prefer explicit template_path if provided by caller
        if template_path:
            effective_template_path = template_path
        else:
            # fall back to config (argument) > module-level config
            cfg = config or globals().get("config", {}) or {}
            template_folder = cfg.get("TEMPLATE_FOLDER", "templates")
            effective_template_path = os.path.join(template_folder, "template.pptx")

        if not os.path.exists(effective_template_path):
            # last-chance: env var or manual prompt (legacy behaviour)
            env_path = os.getenv("TEMPLATE_PATH")
            if env_path and os.path.exists(env_path):
                effective_template_path = env_path
            else:
                effective_template_path = input(
                    "Template not found! Please provide the full path to the template: "
                )

        prs = Presentation(effective_template_path)
        logging.debug(f"Template loaded from: {effective_template_path}")

        # Deprecated
        # def generate_powerpoint_from_analysis(
        #     comparison_result_path,
        #     powerpoint_output_path,
        #     current_file_path,
        #     previous_file_path,
        # ):
        #     """
        #     Backwards-compatible wrapper around generate_powerpoint_from_apm.
        #     Older code still calls this name.
        #     """
        #     return generate_powerpoint_from_apm(
        #         comparison_result_path=comparison_result_path,
        #         powerpoint_output_path=powerpoint_output_path,
        #         current_file_path=current_file_path,
        #         previous_file_path=previous_file_path,
        #     )


        # -------------------------------------------------------------------
        # Nite: Load Analysis from CURRENT workbook to drive counts & maturity
        # -------------------------------------------------------------------
        df_current_analysis = pd.read_excel(current_file_path, sheet_name="Analysis")

        # Count valid applications (non-empty 'name')
        number_of_apps = (
            df_current_analysis["name"].dropna().str.strip().ne("").sum()
        )
        logging.info(
            f"Number of applications in the current 'Analysis' sheet: {number_of_apps}"
        )

        # -------------------------------------------------------------------
        # Nite: Load summary sheets (current, previous, and comparison summary)
        # -------------------------------------------------------------------
        current_summary_df = pd.read_excel(current_file_path, sheet_name="Summary")
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name="Summary")

        summary_df = pd.read_excel(comparison_result_path, sheet_name="Summary")
        logging.debug("Loaded Summary sheet successfully.")
        logging.debug(f"Summary DataFrame head:\n{summary_df.head()}")

        # -------------------------------------------------------------------
        # Nite: Load comparison_result APM sheets for domain-specific slides
        # -------------------------------------------------------------------
        df_analysis = pd.read_excel(comparison_result_path, sheet_name="Analysis")
        df_app_agents = pd.read_excel(comparison_result_path, sheet_name="AppAgentsAPM")
        df_machine_agents = pd.read_excel(
            comparison_result_path, sheet_name="MachineAgentsAPM"
        )
        df_BTs = pd.read_excel(
            comparison_result_path, sheet_name="BusinessTransactionsAPM"
        )
        df_Backends = pd.read_excel(
            comparison_result_path, sheet_name="BackendsAPM"
        )
        df_Overhead = pd.read_excel(
            comparison_result_path, sheet_name="OverheadAPM"
        )
        df_ServiceEndpoints = pd.read_excel(
            comparison_result_path, sheet_name="ServiceEndpointsAPM"
        )
        df_ErrorConfiguration = pd.read_excel(
            comparison_result_path, sheet_name="ErrorConfigurationAPM"
        )
        df_HealthRulesAndAlerting = pd.read_excel(
            comparison_result_path, sheet_name="HealthRulesAndAlertingAPM"
        )
        df_DataCollectors = pd.read_excel(
            comparison_result_path, sheet_name="DataCollectorsAPM"
        )
        df_Dashboards = pd.read_excel(
            comparison_result_path, sheet_name="DashboardsAPM"
        )

        # -------------------------------------------------------------------
        # Nite: Local helpers for Slide 2 – Key Callouts + coverage
        # -------------------------------------------------------------------
        def _parse_percent_to_float(val):
            if pd.isna(val):
                return None
            if isinstance(val, (int, float)):
                return float(val)
            s = str(val).strip()
            if "→" in s:
                s = s.split("→")[-1].strip()
            s = s.replace("%", "")
            try:
                return float(s)
            except Exception:
                return None

        def _get_tier_percent(df, tier):
            name_map = {c.lower(): c for c in df.columns}
            candidates = [
                f"{tier.lower()} %",
                f"{tier.lower()}%",
                f"percentage{tier.lower()}",
                f"{tier.lower()}percentage",
            ]
            for cand in candidates:
                if cand in name_map:
                    return _parse_percent_to_float(df[name_map[cand]].iloc[0])

            needed = ["bronze", "silver", "gold", "platinum"]
            if all(k in name_map for k in needed):
                try:
                    total = 0.0
                    counts = {}
                    for k in needed:
                        val = pd.to_numeric(
                            df[name_map[k]].iloc[0], errors="coerce"
                        )
                        counts[k] = 0.0 if pd.isna(val) else float(val)
                        total += counts[k]
                    if total > 0:
                        return (counts[tier.lower()] / total) * 100.0
                except Exception:
                    return None
            return None

        def _arrow(curr, prev):
            if curr is None or prev is None:
                return "→"
            if curr > prev:
                return "↑"
            if curr < prev:
                return "↓"
            return "→"

        def _trend_word(curr, prev):
            if curr is None or prev is None:
                return "held steady"
            if curr > prev:
                return "increased"
            if curr < prev:
                return "decreased"
            return "held steady"

        def _delta_pp(curr, prev):
            if curr is None or prev is None:
                return None
            return round(curr - prev, 1)

        def _grade_token(s: str):
            if not s:
                return None
            m = re.search(r"(platinum|gold|silver|bronze)", str(s), re.I)
            return m.group(1).lower() if m else None

        def _apps_coverage(path):
            try:
                df = pd.read_excel(path, sheet_name="Analysis")
                total = int(
                    df["name"]
                    .dropna()
                    .astype(str)
                    .str.strip()
                    .ne("")
                    .sum()
                )
                if total == 0:
                    return (0, 0, 0.0)
                rated = int(
                    df["OverallAssessment"].apply(_grade_token).notna().sum()
                )
                pct = (rated / total) * 100.0
                return (total, rated, pct)
            except Exception:
                return (0, 0, None)

        def _arrow_threshold(curr, prev, threshold_pp=0.5):
            if curr is None or prev is None:
                return "→"
            delta = curr - prev
            if delta >= threshold_pp:
                return "↑"
            if delta <= -threshold_pp:
                return "↓"
            return "→"

        def _fmt_pp_delta(prev, curr):
            if prev is None or curr is None:
                return None
            d = curr - prev
            sign = "+" if d > 0 else "−" if d < 0 else "±"
            return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(d):.1f} pp)."

        def _fmt_outcome(prev, curr, delta):
            if prev is None or curr is None:
                return "Data not available."
            sign = (
                "+"
                if delta is not None and delta > 0
                else "−"
                if delta is not None and delta < 0
                else "±"
            )
            if delta is None:
                return f"{prev:.1f}%→{curr:.1f}%"
            return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(delta):.1f} pp)."

        # -------------------------------------------------------------------
        # Nite: Slide 2 (index 1) — Assessment Result - Key Callouts
        # -------------------------------------------------------------------
        slide = prs.slides[1]

        curr_gold = _get_tier_percent(current_summary_df, "Gold")
        prev_gold = _get_tier_percent(previous_summary_df, "Gold")
        curr_plat = _get_tier_percent(current_summary_df, "Platinum")
        prev_plat = _get_tier_percent(previous_summary_df, "Platinum")

        total_prev, rated_prev, cov_prev = _apps_coverage(previous_file_path)
        total_curr, rated_curr, cov_curr = _apps_coverage(current_file_path)
        cov_arrow = _arrow_threshold(cov_curr, cov_prev)

        cov_outcome = (
            f"Coverage: {cov_curr:.1f}% of apps rated ({rated_curr}/{total_curr})."
            if cov_curr is not None and total_curr > 0
            else "Coverage data not available."
        )
        cov_prev_curr = _fmt_pp_delta(cov_prev, cov_curr)
        if cov_prev_curr:
            cov_outcome = f"{cov_outcome} {cov_prev_curr}"

        try:
            df_cmp = pd.read_excel(comparison_result_path, sheet_name="Analysis")
        except Exception:
            df_cmp = None

        def _count_changes(df, col):
            if df is None or col not in df.columns:
                return 0, 0
            s = df[col].astype(str)
            up = s.str.contains("Upgraded", case=False, na=False).sum()
            down = s.str.contains("Downgraded", case=False, na=False).sum()
            return int(up), int(down)

        up_overall, down_overall = _count_changes(df_cmp, "OverallAssessment")
        overall_result_text = (
            "Increase"
            if up_overall > down_overall
            else "Decrease"
            if down_overall > up_overall
            else "Even"
        )

        area_cols = [
            "AppAgentsAPM",
            "MachineAgentsAPM",
            "BusinessTransactionsAPM",
            "BackendsAPM",
            "OverheadAPM",
            "ServiceEndpointsAPM",
            "ErrorConfigurationAPM",
            "HealthRulesAndAlertingAPM",
            "DataCollectorsAPM",
            "DashboardsAPM",
        ]
        pretty = {
            "AppAgentsAPM": "App Agents",
            "MachineAgentsAPM": "Machine Agents",
            "BusinessTransactionsAPM": "Business Transactions",
            "BackendsAPM": "Backends",
            "OverheadAPM": "Overhead",
            "ServiceEndpointsAPM": "Service Endpoints",
            "ErrorConfigurationAPM": "Error Configuration",
            "HealthRulesAndAlertingAPM": "Health Rules & Alerting",
            "DataCollectorsAPM": "Data Collectors",
            "DashboardsAPM": "Dashboards",
        }
        downgraded_counts = []
        if df_cmp is not None:
            for col in area_cols:
                if col in df_cmp.columns:
                    s = df_cmp[col].astype(str)
                    cnt = s.str.contains("Downgraded", case=False, na=False).sum()
                    downgraded_counts.append((col, int(cnt)))
        downgraded_counts.sort(key=lambda x: x[1], reverse=True)
        focus_list = [pretty[c] for c, n in downgraded_counts if n > 0][:2]
        next_focus_text = (
            ", ".join(focus_list) if focus_list else "Maintain current progress"
        )

        delta_gold = _delta_pp(curr_gold, prev_gold)
        delta_plat = _delta_pp(curr_plat, prev_plat)

        headers = [
            "AppD Maturity Progression & Engagement",
            "Commentary",
            "Outcomes",
            "Change/Status Since Last",
        ]

        rows = [
            [
                "B/S/G/P Model Adoption & Maturity Status",
                f"B/S/G/P model applied to APM; assessment covered {int(total_curr)} apps.",
                cov_outcome,
                cov_arrow,
            ],
            [
                "Gold Status Apps",
                f"Gold-or-better coverage {_trend_word(curr_gold, prev_gold)} across the portfolio.",
                _fmt_outcome(prev_gold, curr_gold, delta_gold),
                _arrow(curr_gold, prev_gold),
            ],
            [
                "Platinum Status Apps",
                f"Platinum presence {_trend_word(curr_plat, prev_plat)}; teams progressing on prerequisites.",
                _fmt_outcome(prev_plat, curr_plat, delta_plat),
                _arrow(curr_plat, prev_plat),
            ],
            [
                "Maturity Partnership",
                "Working cadence in place; recommendations implemented during this period.",
                f"Overall result: {overall_result_text}. Next focus: {next_focus_text}.",
                "↑"
                if overall_result_text == "Increase"
                else "↓"
                if overall_result_text == "Decrease"
                else "→",
            ],
        ]

        key_callouts_ph = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if key_callouts_ph:
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(rows) + 1, len(headers)
            )
        else:
            table = (
                slide.shapes.add_table(
                    len(rows) + 1,
                    len(headers),
                    Inches(0.6),
                    Inches(2.1),
                    Inches(9.0),
                    Inches(4.0),
                ).table
            )

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            cell.text_frame.word_wrap = False

        autosize_col_to_header(
            table, 3, header_pt=12, padding_in=0.6, avg_char_em=0.55
        )

        for r_idx, row in enumerate(rows, start=1):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                if c_idx == 3 and value in ("↑", "↓", "→"):
                    set_arrow_cell(cell, value, color=PINK, size_pt=36)
                else:
                    cell.text = str(value)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)

        # -------------------------------------------------------------------
        # Nite: Color the overall maturity oval & write slide notes rationale
        # -------------------------------------------------------------------
        overall_tier = overall_maturity_from_df(
            df_current_analysis, grade_func=_grade_token
        )
        if overall_tier:
            color_oval_for_maturity(
                slide, shape_name="Oval 10", tier=overall_tier, update_text=False
            )

        def _tier_counts(df):
            counts = {"bronze": 0, "silver": 0, "gold": 0, "platinum": 0}
            col = "OverallAssessment"
            if df is None or col not in df.columns:
                return counts, 0
            for v in df[col]:
                t = _grade_token(v)
                if t in counts:
                    counts[t] += 1
            total = sum(counts.values())
            return counts, total

        def _pct(n, d):
            return (n / d) * 100.0 if d else 0.0

        tier_counts, tier_total = _tier_counts(df_current_analysis)
        b, s, g, p = (
            tier_counts["bronze"],
            tier_counts["silver"],
            tier_counts["gold"],
            tier_counts["platinum"],
        )
        pb, ps, pg, pp_ = (
            _pct(b, tier_total),
            _pct(s, tier_total),
            _pct(g, tier_total),
            _pct(p, tier_total),
        )

        rationale = (
            f"Status is {overall_tier} because it has the largest share of rated apps this run. "
            f"Distribution — Platinum {pp_:.1f}% ({p}), Gold {pg:.1f}% ({g}), "
            f"Silver {ps:.1f}% ({s}), Bronze {pb:.1f}% ({b})."
        )

        coverage_note = (
            f"Rated coverage: {cov_curr:.1f}% ({rated_curr}/{total_curr})."
            if cov_curr is not None and total_curr > 0
            else "Rated coverage: n/a."
        )
        next_focus_note = f"Next focus: {next_focus_text}."

        notes = slide.notes_slide
        tf = notes.notes_text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = (
            "Overall tier selection: majority of app ratings in Analysis; "
            "ties prefer the higher tier."
        )
        p2 = tf.add_paragraph()
        p2.text = rationale
        p3 = tf.add_paragraph()
        p3.text = f"{coverage_note} {next_focus_note}"

        # -------------------------------------------------------------------
        # Nite: Slide 4 – Upgraded applications list + count in TextBox 7
        # -------------------------------------------------------------------
        slide = prs.slides[3]
        upgraded_apps = df_analysis[
            df_analysis["OverallAssessment"].str.contains(
                "upgraded", case=False, na=False
            )
        ]["name"].tolist()

        current_analysis_df = pd.read_excel(
            current_file_path, sheet_name="Analysis"
        )
        number_of_apps = len(current_analysis_df)

        textbox_7 = None
        for shape in slide.shapes:
            if shape.name == "TextBox 7":
                textbox_7 = shape
                break

        if textbox_7:
            textbox_7.text = f"{number_of_apps}"
        else:
            logging.warning("TextBox 8 not found on Slide 3.")

        upgraded_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if upgraded_placeholder:
            logging.debug(
                "Found Upgraded Applications table placeholder. Inserting table."
            )
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(upgraded_apps) + 1, 1
            )
        else:
            logging.warning(
                "Upgraded Applications table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(upgraded_apps) + 1,
                1,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Applications with Upgraded Metrics"
        table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(12)

        for idx, app in enumerate(upgraded_apps):
            table.cell(idx + 1, 0).text = app
            table.cell(idx + 1, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # -------------------------------------------------------------------
        # Nite: Slide 5 – Comparison summary + previous vs current summary tables
        # -------------------------------------------------------------------
        slide = prs.slides[4]
        summary_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )

        # Nite: main comparison Result Summary table
        if summary_placeholder:
            logging.debug("Found Summary table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(summary_df) + 1, len(summary_df.columns)
            )
        else:
            logging.warning(
                "Summary table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(summary_df) + 1,
                len(summary_df.columns),
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        for col_idx, column in enumerate(summary_df.columns):
            table.cell(0, col_idx).text = str(column)
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        for row_idx, row in summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)
                table.cell(row_idx + 1, col_idx).text_frame.paragraphs[
                    0
                ].font.size = Pt(12)

        # Nite: previous summary mini-table (Table Placeholder 4)
        summary_placeholder_previous = find_table_placeholder_by_name(
            slide, "Table Placeholder 4"
        )
        if summary_placeholder_previous:
            logging.debug("Found Table Placeholder 4. Inserting table for previous summary.")
            table_previous = insert_table_at_placeholder(
                slide,
                "Table Placeholder 4",
                len(previous_summary_df) + 1,
                len(previous_summary_df.columns),
            )
        else:
            logging.warning("Table Placeholder 4 not found. Adding manually.")
            table_previous = slide.shapes.add_table(
                len(previous_summary_df) + 1,
                len(previous_summary_df.columns),
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        for col_idx, column in enumerate(previous_summary_df.columns):
            table_previous.cell(0, col_idx).text = str(column)
            table_previous.cell(0, col_idx).text_frame.paragraphs[
                0
            ].font.size = Pt(12)

        for row_idx, row in previous_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_previous.cell(row_idx + 1, col_idx).text = str(value)
                table_previous.cell(row_idx + 1, col_idx).text_frame.paragraphs[
                    0
                ].font.size = Pt(12)

        # Nite: current summary mini-table (Table Placeholder 3)
        summary_placeholder_current = find_table_placeholder_by_name(
            slide, "Table Placeholder 3"
        )
        if summary_placeholder_current:
            logging.debug("Found Table Placeholder 3. Inserting table for current summary.")
            table_current = insert_table_at_placeholder(
                slide,
                "Table Placeholder 3",
                len(current_summary_df) + 1,
                len(current_summary_df.columns),
            )
        else:
            logging.warning("Table Placeholder 3 not found. Adding manually.")
            table_current = slide.shapes.add_table(
                len(current_summary_df) + 1,
                len(current_summary_df.columns),
                Inches(0.5),
                Inches(6),
                Inches(9),
                Inches(4),
            ).table

        for col_idx, column in enumerate(current_summary_df.columns):
            table_current.cell(0, col_idx).text = str(column)
            table_current.cell(0, col_idx).text_frame.paragraphs[
                0
            ].font.size = Pt(12)

        for row_idx, row in current_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_current.cell(row_idx + 1, col_idx).text = str(value)
                table_current.cell(row_idx + 1, col_idx).text_frame.paragraphs[
                    0
                ].font.size = Pt(12)

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Comparison Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        # -------------------------------------------------------------------
        # Nite: Overall / per-area upgraded vs downgraded counts for Slides 7 & 8
        # -------------------------------------------------------------------
        df = pd.read_excel(comparison_result_path, sheet_name="Analysis")
        columns = [
            "AppAgentsAPM",
            "MachineAgentsAPM",
            "BusinessTransactionsAPM",
            "BackendsAPM",
            "OverheadAPM",
            "ServiceEndpointsAPM",
            "ErrorConfigurationAPM",
            "HealthRulesAndAlertingAPM",
            "DataCollectorsAPM",
            "DashboardsAPM",
            "OverallAssessment",
        ]

        results = {}
        total_applications = len(df)

        for col in columns:
            df[col] = df[col].astype(str)
            upgraded_count = df[col].str.contains(
                "upgraded", case=False, na=False
            ).sum()
            downgraded_count = df[col].str.contains(
                "downgraded", case=False, na=False
            ).sum()

            total_applications = len(df[col])
            overall_result = (
                "Increase"
                if upgraded_count > downgraded_count
                else "Decrease"
                if downgraded_count > upgraded_count
                else "Even"
            )
            percentage_value = (
                0
                if overall_result == "Even"
                else round((upgraded_count / total_applications) * 100)
            )

            results[col] = {
                "upgraded": upgraded_count,
                "downgraded": downgraded_count,
                "overall_result": overall_result,
                "percentage": percentage_value,
            }

        # -------------------------------------------------------------------
        # Nite: Slide 7 – Overall Assessment Result table
        # -------------------------------------------------------------------
        slide = prs.slides[6]
        overall_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )

        if overall_placeholder:
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", 2, 5)
        else:
            table = slide.shapes.add_table(
                2, 5, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
            ).table

        headers = [
            "Metric",
            "# of Apps Improved",
            "# Apps Degraded",
            "Overall Result",
            "Percentage Value",
        ]
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        overall_assessment = results["OverallAssessment"]
        table.cell(1, 0).text = "OverallAssessment"
        table.cell(1, 1).text = str(overall_assessment["upgraded"])
        table.cell(1, 2).text = str(overall_assessment["downgraded"])
        table.cell(1, 3).text = overall_assessment["overall_result"]
        table.cell(1, 4).text = f"{overall_assessment['percentage']}%"

        if overall_assessment["overall_result"] == "Increase":
            table.cell(1, 4).fill.solid()
            table.cell(1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Overall Assessment Result"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        # -------------------------------------------------------------------
        # Nite: Slide 8 – Status table per APM domain
        # -------------------------------------------------------------------
        slide = prs.slides[7]
        status_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )

        num_rows = len(columns)
        num_cols = 5

        if status_placeholder:
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", num_rows, num_cols
            )
        else:
            table = slide.shapes.add_table(
                num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)
            ).table

        headers = [
            "Metric",
            "# of Apps Improved",
            "# Apps Degraded",
            "Overall Result",
            "Percentage Value",
        ]
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        for i, col in enumerate(columns[:-1]):
            table.cell(i + 1, 0).text = col
            table.cell(i + 1, 1).text = str(results[col]["upgraded"])
            table.cell(i + 1, 2).text = str(results[col]["downgraded"])
            table.cell(i + 1, 3).text = results[col]["overall_result"]
            table.cell(i + 1, 4).text = f"{results[col]['percentage']}%"

            if results[col]["overall_result"] == "Increase":
                table.cell(i + 1, 4).fill.solid()
                table.cell(i + 1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "APM Maturity Assessment Result"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        # -------------------------------------------------------------------
        # Nite: From here down – “Downgrade summary” slides for each APM area
        # Each block:
        #   - finds downgrades in df_analysis[column]
        #   - fills a grade breakdown table
        #   - updates “Declined/Decreased/Changed” rectangles from domain sheet
        # -------------------------------------------------------------------

        # -------------------------------------------------------------------
        # Nite: Slide 12 – App Agents Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[11]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["AppAgentsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "APM Agent - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "percentAgentsLessThan1YearOld": "Rectangle 11",
            "metricLimitNotHit": "Rectangle 10",
            "percentAgentsLessThan2YearsOld": "Rectangle 12",
            "percentAgentsReportingData": "Rectangle 13",
            "percentAgentsRunningSameVersion": "Rectangle 14",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_app_agents.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "declined" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.")

        # -------------------------------------------------------------------
        # Nite: Slide 13 – Machine Agents Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[12]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["MachineAgentsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Machine Agent - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "percentAgentsLessThan1YearOld": "Rectangle 8",
            "percentAgentsLessThan2YearsOld": "Rectangle 9",
            "percentAgentsReportingData": "Rectangle 10",
            "percentAgentsRunningSameVersion": "Rectangle 11",
            "percentAgentsInstalledAlongsideAppAgents": "Rectangle 12",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_machine_agents.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "declined" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 13.")

        # -------------------------------------------------------------------
        # Nite: Slide 14 – Business Transactions Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[13]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["BusinessTransactionsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Business Transactions - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "numberOfBTs": "Rectangle 17",
            "percentBTsWithLoad": "Rectangle 18",
            "btLockdownEnabled": "Rectangle 19",
            "numberCustomMatchRules": "Rectangle 20",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_BTs.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 14.")

        # -------------------------------------------------------------------
        # Nite: Slide 15 – Backends Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[14]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["BackendsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            applications_str = (
                ", ".join(str(app) for app in downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            logging.debug(f"Grade: {grade}, Applications: {applications_str}")

            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = applications_str
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Backends - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "percentBackendsWithLoad": "Rectangle 10",
            "backendLimitNotHit": "Rectangle 11",
            "numberOfCustomBackendRules": "Rectangle 12",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_Backends.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.")

        # -------------------------------------------------------------------
        # Nite: Slide 16 – Service Endpoints Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[15]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["ServiceEndpointsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Service Endpoints - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "numberOfCustomServiceEndpointRules": "Rectangle 10",
            "serviceEndpointLimitNotHit": "Rectangle 11",
            "percentServiceEndpointsWithLoadOrDisabled": "Rectangle 12",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_ServiceEndpoints.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 16.")

        # -------------------------------------------------------------------
        # Nite: Slide 17 – Error Configuration Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[16]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["ErrorConfigurationAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Error Configuration - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "successPercentageOfWorstTransaction": "Rectangle 10",
            "numberOfCustomRules": "Rectangle 11",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_ErrorConfiguration.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 17.")

        # -------------------------------------------------------------------
        # Nite: Slide 18 – Health Rules & Alerting Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[17]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["HealthRulesAndAlertingAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Health Rules & Alerting - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "numberOfHealthRuleViolations": "Rectangle 10",
            "numberOfDefaultHealthRulesModified": "Rectangle 11",
            "numberOfActionsBoundToEnabledPolicies": "Rectangle 12",
            "numberOfCustomHealthRules": "Rectangle 13",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_HealthRulesAndAlerting.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 18.")

        # -------------------------------------------------------------------
        # Nite: Slide 19 – Data Collectors Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[18]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["DataCollectorsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Data Collectors - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "numberOfDataCollectorFieldsConfigured": "Rectangle 10",
            "numberOfDataCollectorFieldsCollectedInSnapshots": "Rectangle 11",
            "numberOfDataCollectorFieldsCollectedInAnalytics": "Rectangle 12",
            "biqEnabled": "Rectangle 13",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_DataCollectors.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 19.")

        # -------------------------------------------------------------------
        # Nite: Slide 20 – Dashboards Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[19]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["DashboardsAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Dashboards - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "numberOfDashboards": "Rectangle 10",
            "percentageOfDashboardsModifiedLast6Months": "Rectangle 11",
            "numberOfDashboardsUsingBiQ": "Rectangle 12",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_Dashboards.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "decreased" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 20.")

        # -------------------------------------------------------------------
        # Nite: Slide 21 – Overhead Downgrade Summary + rectangles
        # -------------------------------------------------------------------
        slide = prs.slides[20]

        all_grades = ["platinum", "gold", "silver", "bronze"]
        grades_for_table = ["gold", "silver", "bronze"]
        downgrade_data = {
            grade: {"applications": [], "number_of_apps": 0, "percentage": 0}
            for grade in grades_for_table
        }

        for _, row in df_analysis.iterrows():
            current_value = row["OverheadAPM"]
            app_name = row["name"]
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            if "→" in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")
                try:
                    previous_value, current_grade = current_value.split("→")
                    previous_value = previous_value.strip().lower()
                    current_grade = current_grade.strip().lower().split(" ")[0]

                    logging.debug(
                        f"Extracted: Previous Value: {previous_value}, "
                        f"Current Grade: {current_grade}"
                    )

                    if previous_value in all_grades and current_grade in all_grades:
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(
                                f"Adding {app_name} to {current_grade} downgrade list"
                            )
                            downgrade_data[current_grade]["applications"].append(app_name)
                            downgrade_data[current_grade]["number_of_apps"] += 1
                        else:
                            logging.debug(
                                f"Not a downgrade for {app_name}: "
                                f"{previous_value} → {current_grade}"
                            )
                    else:
                        logging.debug(
                            f"Invalid grades for downgrade comparison: "
                            f"{previous_value}, {current_grade}"
                        )
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(
                    f"No Downgrade for App: {app_name} - Current Value: {current_value}"
                )

        for grade in grades_for_table:
            logging.debug(
                f"Applications for {grade}: {downgrade_data[grade]['applications']}"
            )

        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]["percentage"] = (
                len(downgrade_data[grade]["applications"]) / total_apps * 100
                if total_apps
                else 0
            )

        downgrade_placeholder = find_table_placeholder_by_name(
            slide, "Table Placeholder 1"
        )
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(
                slide, "Table Placeholder 1", len(grades_for_table) + 1, 4
            )
        else:
            logging.warning(
                "Downgrade table placeholder not found. Adding manually."
            )
            table = slide.shapes.add_table(
                len(grades_for_table) + 1,
                4,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(4),
            ).table

        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        for i, grade in enumerate(grades_for_table):
            table.cell(i + 1, 0).text = grade.capitalize()
            table.cell(i + 1, 1).text = (
                ", ".join(downgrade_data[grade]["applications"])
                if downgrade_data[grade]["applications"]
                else "None"
            )
            table.cell(i + 1, 2).text = str(
                downgrade_data[grade]["number_of_apps"]
            )
            table.cell(i + 1, 3).text = (
                f"{downgrade_data[grade]['percentage']:.2f}%"
            )

        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder and hasattr(title_placeholder, "text_frame"):
            title_placeholder.text = "Overhead - Downgrade Summary"
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

        columns_and_rectangles = {
            "developerModeNotEnabledForAnyBT": "Rectangle 10",
            "findEntryPointsNotEnabled": "Rectangle 11",
            "aggressiveSnapshottingNotEnabled": "Rectangle 12",
            "developerModeNotEnabledForApplication": "Rectangle 13",
        }
        declined_counts = {key: 0 for key in columns_and_rectangles}

        for _, row in df_Overhead.iterrows():
            for column, _rectangle in columns_and_rectangles.items():
                if "changed" in str(row[column]).lower():
                    declined_counts[column] += 1

        for column, rectangle_name in columns_and_rectangles.items():
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"
            else:
                logging.warning(f"{rectangle_name} not found on Slide 21.")

        # -------------------------------------------------------------------
        # Nite: Finally, save the finished deck
        # -------------------------------------------------------------------
        prs.save(powerpoint_output_path)
        logging.debug(f"PowerPoint saved to {powerpoint_output_path}.")
        logging.warning(f">>> SAVED PPT TO {powerpoint_output_path} <<<")

    except Exception as e:
        logging.error(f"Error generating PowerPoint: {e}", exc_info=True)
        raise
