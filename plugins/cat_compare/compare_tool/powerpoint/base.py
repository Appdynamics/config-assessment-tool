"""
base.py
-------
This module provides a base class for building PowerPoint presentations.

Purpose:
- Acts as a thin wrapper around the `python-pptx` library.
- Provides convenience methods for creating and customizing slides.

Key Features:
- `PPTBuilder` class:
  - Loads a PowerPoint template or creates a blank presentation.
  - Provides utility methods for common slide operations.
- Simplifies the process of generating PowerPoint presentations programmatically.
"""

# compare_tool/powerpoint/base.py

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

log = logging.getLogger(__name__)
log.info("base.py imported")


class PPTBuilder:
    """
    Thin wrapper around python-pptx Presentation that provides
    convenience methods for common slide types.
    """

    def __init__(self, template_path: Optional[str] = None) -> None:
        if template_path and Path(template_path).exists():
            log.info("Loading PowerPoint template from %s", template_path)
            self.prs = Presentation(template_path)
        else:
            if template_path:
                log.warning(
                    "Template %s not found, using blank presentation", template_path
                )
            self.prs = Presentation()

    # --- Basic slide helpers -------------------------------------------------

    def add_title_slide(self, title: str, subtitle: Optional[str] = None) -> None:
        layout = self.prs.slide_layouts[0]  # usually title slide
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        if subtitle is not None:
            # Most title layouts have subtitle as placeholder index 1
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle

    def add_title_only_slide(self, title: str) -> "object":
        layout = self.prs.slide_layouts[5]  # title-only in default template
        slide = self.prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        return slide

    def add_bullet_slide(self, title: str, bullets: list[str]) -> None:
        """
        Add a slide with a title and a bullet list.
        """
        layout = self.prs.slide_layouts[1]  # title + content
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        if not bullets:
            return

        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    def add_table_slide(self, title: str, df, subtitle: Optional[str] = None) -> None:
        """
        Add a slide with a title and a table from a pandas DataFrame.
        """
        layout = self.prs.slide_layouts[5]  # title only
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        rows, cols = df.shape
        rows += 1  # header row

        # Simple layout – centre of slide
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(0.8 + 0.3 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        from . import slides  # local import to avoid cyclic issues

        # Header
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            slides.style_header_cell(cell)

        # Body
        for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = "" if value is None else str(value)

        slides.autofit_table_columns(table)

    def save(self, path: str) -> None:
        out = Path(path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(out)
        log.info("Saved presentation to %s", out)
