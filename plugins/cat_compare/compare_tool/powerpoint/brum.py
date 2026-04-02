# compare_tool/powerpoint/brum_mrum.py

"""
brum.py
-------
This module handles the generation of PowerPoint presentations for BRUM (Business Resource Utilization Management) comparisons.

Purpose:
- Creates PowerPoint slides specific to BRUM comparison results.
- Uses templates to dynamically populate slides with BRUM-related data.

Key Features:
- Reads BRUM comparison data and generates slides with metrics, charts, and summaries.
- Saves the generated PowerPoint presentation to the specified output path.
"""

from __future__ import annotations

import logging
import os
import re
from typing import Optional, Dict, Any

import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

log = logging.getLogger(__name__)

# --------------------------------------------------------------------
# Small helper utilities so this module is self-contained
# --------------------------------------------------------------------

PINK = RGBColor(226, 0, 116)  # AppD-ish pink, used for arrows


def autosize_col_to_header(*args, **kwargs):
    """
    Legacy helper – safe no-op here.
    Your original version adjusted column width to header text;
    we don't strictly need it for correctness.
    """
    return


def set_arrow_cell(cell, direction: str, color=PINK, size_pt: int = 36, font_name: str = "Calibri"):
    """
    Render a big arrow glyph into a table cell.
    This is a simple approximation of your original helper.
    """
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = direction
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.name = font_name


# --------------------------------------------------------------------
# BRUM POWERPOINT
# --------------------------------------------------------------------

def generate_powerpoint_from_brum(
    comparison_result_path: str,
    powerpoint_output_path: str,
    current_file_path: str,
    previous_file_path: str,
    config: Optional[Dict[str, Any]] = None,
) -> str:
    """
    Legacy BRUM generator, adapted to:
      - take an explicit `config` dict (instead of global)
      - avoid interactive input() for missing templates
      - be self-contained in this module
    """

    cfg = config or {}
    log.debug("[BRUM] Generating PowerPoint presentation...")

    try:
        # ------------------------------------------------------------------
        # Template resolution – same spirit as APM
        # ------------------------------------------------------------------
        template_folder = cfg.get("TEMPLATE_FOLDER", "templates")
        template_path = os.path.join(template_folder, "template_brum.pptx")

        if not os.path.exists(template_path):
            env_path = os.getenv("TEMPLATE_BRUM_PATH", "")
            if env_path and os.path.exists(env_path):
                template_path = env_path
            else:
                log.warning(
                    "[BRUM] template_brum.pptx not found at %s; "
                    "falling back to a blank presentation.",
                    template_path,
                )
                prs = Presentation()
        else:
            prs = Presentation(template_path)
            log.debug("[BRUM] Template loaded from: %s", template_path)

        # ------------------------------------------------------------------
        # Load Excel data
        # ------------------------------------------------------------------
        df_current_analysis = pd.read_excel(current_file_path, sheet_name="Analysis")
        number_of_apps = (
            df_current_analysis["name"].dropna().astype(str).str.strip().ne("").sum()
        )
        log.info("[BRUM] Number of applications in the current 'Analysis' sheet: %s", number_of_apps)

        current_summary_df = pd.read_excel(current_file_path, sheet_name="Summary")
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name="Summary")

        summary_df = pd.read_excel(comparison_result_path, sheet_name="Summary")
        log.debug("[BRUM] Loaded Summary sheet successfully.")
        log.debug("[BRUM] Summary DataFrame head:\n%s", summary_df.head())

        df_analysis = pd.read_excel(comparison_result_path, sheet_name="Analysis")

        df_network_requests = pd.read_excel(comparison_result_path, sheet_name="NetworkRequestsBRUM")
        df_health_rules = pd.read_excel(comparison_result_path, sheet_name="HealthRulesAndAlertingBRUM")
        df_overall_brum = pd.read_excel(comparison_result_path, sheet_name="OverallAssessmentBRUM")

        # ------------------------------------------------------------------
        # Placeholders helpers
        # ------------------------------------------------------------------
        def find_table_placeholder_by_name(slide, name):
            for shape in slide.shapes:
                if shape.is_placeholder and shape.name == name:
                    return shape
            return None

        def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
            placeholder = find_table_placeholder_by_name(slide, placeholder_name)
            if not placeholder:
                log.error("[BRUM] Placeholder '%s' not found on the slide.", placeholder_name)
                return None
            left, top, width, height = (
                placeholder.left,
                placeholder.top,
                placeholder.width,
                placeholder.height,
            )
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            return table_shape.table

        def find_slides_with_placeholder(prs_obj, placeholder_name):
            return [s for s in prs_obj.slides if find_table_placeholder_by_name(s, placeholder_name)]

        def slide_title_text(slide):
            for shp in slide.shapes:
                if getattr(shp, "name", "") == "Title 2" and hasattr(shp, "text_frame") and shp.text_frame:
                    return shp.text_frame.text.strip()
            for shp in slide.shapes:
                if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.text:
                    return shp.text_frame.text.strip()
            return ""

        used_slide_ids = set()

        def choose_slide_for_section(prefer_titles, required_placeholders=("Table Placeholder 1",), exclude_ids=None):
            exclude_ids = exclude_ids or set()
            # Try by title first
            for s in prs.slides:
                if id(s) in exclude_ids:
                    continue
                title = slide_title_text(s)
                if any(pt.lower() in title.lower() for pt in prefer_titles):
                    if all(find_table_placeholder_by_name(s, p) for p in required_placeholders):
                        return s
            # Fallback: first slide with required placeholders not excluded
            for s in prs.slides:
                if id(s) in exclude_ids:
                    continue
                if all(find_table_placeholder_by_name(s, p) for p in required_placeholders):
                    return s
            return None

        # ------------------------------------------------------------------
        # BRUM Key Callouts
        # ------------------------------------------------------------------
        try:
            curr_overall_df = pd.read_excel(current_file_path, sheet_name="OverallAssessmentBRUM")
        except Exception:
            curr_overall_df = pd.DataFrame()

        try:
            prev_overall_df = pd.read_excel(previous_file_path, sheet_name="OverallAssessmentBRUM")
        except Exception:
            prev_overall_df = pd.DataFrame()

        def last_percent(df, col):
            if df.empty or col not in df.columns:
                return None
            s = pd.to_numeric(
                df[col].astype(str).str.replace("%", ""), errors="coerce"
            ).dropna()
            return float(s.iloc[-1]) if not s.empty else None

        curr_gold = last_percent(curr_overall_df, "percentageTotalGoldOrBetter")
        prev_gold = last_percent(prev_overall_df, "percentageTotalGoldOrBetter")
        curr_plat = last_percent(curr_overall_df, "percentageTotalPlatinum")
        prev_plat = last_percent(prev_overall_df, "percentageTotalPlatinum")

        def count_changes(df, col):
            if col not in df.columns:
                return 0, 0
            s = df[col].astype(str)
            return (
                s.str.contains("Upgraded", case=False, na=False).sum(),
                s.str.contains("Downgraded", case=False, na=False).sum(),
            )

        oa_up, oa_down = count_changes(df_analysis, "OverallAssessmentBRUM")
        net_up, net_down = count_changes(df_analysis, "NetworkRequestsBRUM")
        hra_up, hra_down = count_changes(df_analysis, "HealthRulesAndAlertingBRUM")

        def arrow(curr, prev):
            if curr is None or prev is None:
                return "→"
            return "↑" if curr > prev else "↓" if curr < prev else "→"

        def fmt_change(prev, curr, suffix="%"):
            if prev is None or curr is None:
                return "Change observed."
            if curr > prev:
                return f"Increase from {prev:.1f}{suffix}→{curr:.1f}{suffix}"
            if curr < prev:
                return f"Decrease from {prev:.1f}{suffix}→{curr:.1f}{suffix}"
            return f"No change ({curr:.1f}{suffix})."

        slides_with_ph = find_slides_with_placeholder(prs, "Table Placeholder 1")
        key_callouts_slide = (
            slides_with_ph[0] if slides_with_ph else (prs.slides[0] if len(prs.slides) else None)
        )
        if key_callouts_slide:
            used_slide_ids.add(id(key_callouts_slide))

        if key_callouts_slide:
            headers = [
                "AppD Maturity Progression & Engagement",
                "Commentary",
                "Outcomes",
                "Change/Status Since Last",
            ]
            rows = [
                [
                    "B/S/G/P Model Adoption & Maturity Status (BRUM).",
                    f"BRUM analysis coverage across {int(number_of_apps)} applications.",
                    f"Overall BRUM upgrades: {oa_up}, downgrades: {oa_down}.",
                    "↑" if oa_up > oa_down else "↓" if oa_down > oa_up else "→",
                ],
                [
                    "Gold Status Apps (BRUM).",
                    "Change in gold-or-better coverage across the portfolio.",
                    fmt_change(prev_gold, curr_gold),
                    arrow(curr_gold, prev_gold),
                ],
                [
                    "Platinum Status Apps (BRUM).",
                    "Top-tier BRUM maturity presence across applications.",
                    f"{curr_plat:.1f}% platinum." if curr_plat is not None else "Platinum presence observed.",
                    arrow(curr_plat, prev_plat),
                ],
                [
                    "Maturity Partnership (BRUM).",
                    "Ongoing improvements in BRUM instrumentation and alerting.",
                    f"Network Requests (↑{net_up}/↓{net_down}), Health Rules (↑{hra_up}/↓{hra_down}).",
                    "↑"
                    if (net_up + hra_up) > (net_down + hra_down)
                    else "↓"
                    if (net_down + hra_down) > (net_up + hra_up)
                    else "→",
                ],
            ]

            table = insert_table_at_placeholder(
                key_callouts_slide, "Table Placeholder 1", len(rows) + 1, len(headers)
            )
            if table:
                # Header row
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(12)

                # Data rows
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, value in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        if c_idx == 3 and value in ("↑", "↓", "→"):
                            set_arrow_cell(cell, value, color=PINK, size_pt=36, font_name="Calibri")
                        else:
                            cell.text = str(value)
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(12)
        else:
            log.error("[BRUM] No slide available to place Key Callouts table.")

        # ------------------------------------------------------------------
        # BRUM maturity badge & notes (uses current Analysis)
        # ------------------------------------------------------------------
        try:
            def _brum_grade_token(v):
                if not isinstance(v, str):
                    v = "" if pd.isna(v) else str(v)
                m = re.search(r"(platinum|gold|silver|bronze)", v, re.I)
                return m.group(1).lower() if m else None

            def _overall_maturity_from_df_brum(df):
                col = "OverallAssessmentBRUM"
                if df is None or col not in df.columns:
                    return None, {"bronze": 0, "silver": 0, "gold": 0, "platinum": 0}, 0
                counts = {"bronze": 0, "silver": 0, "gold": 0, "platinum": 0}
                rated = 0
                for v in df[col]:
                    t = _brum_grade_token(v)
                    if t in counts:
                        counts[t] += 1
                        rated += 1
                if rated == 0:
                    return None, counts, 0
                rank = {"bronze": 0, "silver": 1, "gold": 2, "platinum": 3}
                best = max(counts.items(), key=lambda kv: (kv[1], rank[kv[0]]))
                return best[0].title(), counts, rated

            def _ideal_text_rgb_local(rgb):
                r, g, b = rgb[0], rgb[1], rgb[2]
                bright = (r * 299 + g * 587 + b * 114) / 1000.0
                return RGBColor(255, 255, 255) if bright < 140 else RGBColor(31, 31, 31)

            def _color_oval_for_maturity_local(slide_obj, shape_name, tier):
                palette = {
                    "Bronze": RGBColor(205, 127, 50),
                    "Silver": RGBColor(166, 166, 166),
                    "Gold": RGBColor(255, 192, 0),
                    "Platinum": RGBColor(190, 190, 200),
                }
                if slide_obj is None or tier not in palette:
                    return
                target = next(
                    (sh for sh in slide_obj.shapes if getattr(sh, "name", "") == shape_name),
                    None,
                )
                if not target:
                    return
                target.fill.solid()
                target.fill.fore_color.rgb = palette[tier]
                if hasattr(target, "text_frame") and target.text_frame:
                    fg = _ideal_text_rgb_local(palette[tier])
                    for p in target.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = fg

            overall_tier, counts, rated = _overall_maturity_from_df_brum(df_current_analysis)
            if overall_tier and key_callouts_slide is not None:
                _color_oval_for_maturity_local(key_callouts_slide, "Oval 10", overall_tier)

                total_apps = int(number_of_apps) if number_of_apps else 0
                coverage = (rated / total_apps * 100.0) if total_apps else 0.0

                def pct(n):
                    return (n / rated * 100.0) if rated else 0.0

                pb, ps, pg, pp_ = (
                    pct(counts["bronze"]),
                    pct(counts["silver"]),
                    pct(counts["gold"]),
                    pct(counts["platinum"]),
                )

                notes = key_callouts_slide.notes_slide
                tf = notes.notes_text_frame
                tf.clear()
                tf.paragraphs[0].text = (
                    "Overall tier selection: majority of app ratings in Analysis; "
                    "ties prefer higher tier (Platinum > Gold > Silver > Bronze)."
                )
                p2 = tf.add_paragraph()
                p2.text = (
                    f"Status is {overall_tier} based on rated distribution — "
                    f"Platinum {pp_:.1f}% ({counts['platinum']}), "
                    f"Gold {pg:.1f}% ({counts['gold']}), "
                    f"Silver {ps:.1f}% ({counts['silver']}), "
                    f"Bronze {pb:.1f}% ({counts['bronze']})."
                )
                p3 = tf.add_paragraph()
                p3.text = f"Rated coverage this run: {coverage:.1f}% ({rated}/{total_apps})."
        except Exception as e:
            log.warning("[BRUM] Maturity badge/notes skipped: %s", e)

        # ------------------------------------------------------------------
        # Applications Improved (2nd slide with Table Placeholder 1)
        # ------------------------------------------------------------------
        improved = []
        cols_map = [
            ("NetworkRequestsBRUM", "Network Requests"),
            ("HealthRulesAndAlertingBRUM", "Health Rules & Alerting"),
            ("OverallAssessmentBRUM", "Overall"),
        ]

        for _, r in df_analysis.iterrows():
            app = str(r.get("name", "") or "").strip()
            if not app:
                continue
            areas = []
            for col, label in cols_map:
                if col in df_analysis.columns:
                    val = r.get(col, "")
                    if isinstance(val, str) and "upgraded" in val.lower():
                        areas.append(label)
            if areas:
                improved.append((app, ", ".join(areas)))

        improved.sort(key=lambda x: x[0].lower())

        improved_slide = None
        for s in slides_with_ph:
            if id(s) not in used_slide_ids:
                improved_slide = s
                break
        if improved_slide:
            used_slide_ids.add(id(improved_slide))
            headers = ["Application", "Improvement Areas"]
            row_count = max(1, len(improved)) + 1
            table = insert_table_at_placeholder(
                improved_slide, "Table Placeholder 1", row_count, len(headers)
            )
            if table:
                for c, h in enumerate(headers):
                    table.cell(0, c).text = h
                    table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
                if improved:
                    for idx, (app, areas) in enumerate(improved, start=1):
                        table.cell(idx, 0).text = app
                        table.cell(idx, 1).text = areas
                        table.cell(idx, 0).text_frame.paragraphs[0].font.size = Pt(12)
                        table.cell(idx, 1).text_frame.paragraphs[0].font.size = Pt(12)
                else:
                    table.cell(1, 0).text = "No applications improved in this period."
                    table.cell(1, 1).text = ""
        else:
            log.warning("[BRUM] No slide found for Improved Applications.")

        # ------------------------------------------------------------------
        # Summary slide: Previous / Current / Comparison
        # (Table Placeholder 4 / 3 / 1)
        # ------------------------------------------------------------------
        def find_slide_with_all_placeholders(prs_obj, names):
            for s in prs_obj.slides:
                if all(find_table_placeholder_by_name(s, n) for n in names):
                    return s
            return None

        summary_slide = find_slide_with_all_placeholders(
            prs, ["Table Placeholder 1", "Table Placeholder 3", "Table Placeholder 4"]
        )
        if not summary_slide:
            for s in prs.slides:
                if any(
                    find_table_placeholder_by_name(s, n)
                    for n in ["Table Placeholder 1", "Table Placeholder 3", "Table Placeholder 4"]
                ):
                    summary_slide = s
                    break
        if summary_slide:
            used_slide_ids.add(id(summary_slide))

            def fill_table_from_df(table, df):
                # Header
                for c, col in enumerate(df.columns):
                    table.cell(0, c).text = str(col)
                    table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
                # Rows
                for r_idx, row in df.iterrows():
                    for c_idx, val in enumerate(row):
                        table.cell(r_idx + 1, c_idx).text = "" if pd.isna(val) else str(val)
                        table.cell(r_idx + 1, c_idx).text_frame.paragraphs[0].font.size = Pt(12)

            ph4 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 4")
            if ph4:
                table_prev = insert_table_at_placeholder(
                    summary_slide,
                    "Table Placeholder 4",
                    len(previous_summary_df) + 1,
                    len(previous_summary_df.columns),
                )
                if table_prev:
                    fill_table_from_df(table_prev, previous_summary_df)

            ph3 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 3")
            if ph3:
                table_curr = insert_table_at_placeholder(
                    summary_slide,
                    "Table Placeholder 3",
                    len(current_summary_df) + 1,
                    len(current_summary_df.columns),
                )
                if table_curr:
                    fill_table_from_df(table_curr, current_summary_df)

            ph1 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 1")
            if ph1:
                table_comp = insert_table_at_placeholder(
                    summary_slide,
                    "Table Placeholder 1",
                    len(summary_df) + 1,
                    len(summary_df.columns),
                )
                if table_comp:
                    fill_table_from_df(table_comp, summary_df)
        else:
            log.error("[BRUM] No suitable slide found for Summary tables.")

        # ------------------------------------------------------------------
        # Overall Assessment slide
        # ------------------------------------------------------------------
        overall_slide = choose_slide_for_section(
            prefer_titles=["Overall Assessment", "Overall BRUM Assessment"],
            required_placeholders=("Table Placeholder 1",),
            exclude_ids=used_slide_ids,
        )
        if overall_slide:
            used_slide_ids.add(id(overall_slide))
            headers = [
                "Metric",
                "# of Apps Improved",
                "# Apps Degraded",
                "Overall Result",
                "Percentage Value",
            ]
            table = insert_table_at_placeholder(overall_slide, "Table Placeholder 1", 2, len(headers))
            if table:
                for i, h in enumerate(headers):
                    table.cell(0, i).text = h
                    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)

                overall_result = (
                    "Increase" if oa_up > oa_down else "Decrease" if oa_down > oa_up else "Even"
                )
                percentage_value = (
                    0
                    if overall_result == "Even"
                    else round((oa_up / max(1, oa_up + oa_down)) * 100)
                )

                table.cell(1, 0).text = "Overall BRUM Assessment"
                table.cell(1, 1).text = str(oa_up)
                table.cell(1, 2).text = str(oa_down)
                table.cell(1, 3).text = overall_result
                table.cell(1, 4).text = f"{percentage_value}%"

                p = table.cell(1, 4).text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                if overall_result == "Increase":
                    run.font.color.rgb = RGBColor(0, 176, 80)
                elif overall_result == "Decrease":
                    run.font.color.rgb = RGBColor(192, 0, 0)
        else:
            log.warning("[BRUM] Could not find slide for Overall Assessment.")

                # ============================
        # BRUM Entity Comparison (Slide 8) — Table Placeholder 1
        # Headers: Metric | # of Apps Improved | # Apps Degraded | Overall Result | Percentage Value
        # Rows: NetworkRequestsBRUM, HealthRulesAndAlertingBRUM
        # ============================

        def result_and_percentage(up, down):
            # Match overall slide logic: Percentage = improved ratio.
            if up > down:
                return "Increase", round((up / max(1, up + down)) * 100)
            if down > up:
                return "Decrease", round((up / max(1, up + down)) * 100)
            return "Even", 0

        # Prefer a slide by title; otherwise fall back to slide index 7 (8th slide),
        # or the next slide that has "Table Placeholder 1" not yet used.
        entity_slide = choose_slide_for_section(
            prefer_titles=[
                "BRUM Comparison",
                "Entity Comparison",
                "Comparison Result",
                "BRUM Entity Comparison",
                "BRUM Maturity Assessment Result",
            ],
            required_placeholders=("Table Placeholder 1",),
            exclude_ids=used_slide_ids,
        )
        if not entity_slide and len(prs.slides) > 7:
            entity_slide = prs.slides[7]
            # Ensure it has the table placeholder; if not, try any unused slide with the placeholder.
            if not find_table_placeholder_by_name(entity_slide, "Table Placeholder 1"):
                for s in prs.slides:
                    if id(s) in used_slide_ids:
                        continue
                    if find_table_placeholder_by_name(s, "Table Placeholder 1"):
                        entity_slide = s
                        break

        if entity_slide and find_table_placeholder_by_name(entity_slide, "Table Placeholder 1"):
            used_slide_ids.add(id(entity_slide))

            headers = ["Metric", "# of Apps Improved", "# Apps Degraded", "Overall Result", "Percentage Value"]

            # Compute results for each BRUM area using the already-computed counters
            # from df_analysis (net_up/down, hra_up/down).
            net_result, net_pct = result_and_percentage(net_up, net_down)
            hra_result, hra_pct = result_and_percentage(hra_up, hra_down)

            rows = [
                ["Network Requests (BRUM)", str(net_up), str(net_down), net_result, f"{net_pct}%"],
                ["Health Rules & Alerting (BRUM)", str(hra_up), str(hra_down), hra_result, f"{hra_pct}%"],
            ]

            table = insert_table_at_placeholder(entity_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                # Header row
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(12)

                # Data rows
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, value in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = value
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(12)
                        # Color Percentage Value based on Overall Result.
                        if c_idx == 4:
                            result_text = rows[r_idx - 1][3]
                            run = p.runs[0] if p.runs else p.add_run()
                            if result_text == "Increase":
                                run.font.color.rgb = RGBColor(0, 176, 80)   # Green
                            elif result_text == "Decrease":
                                run.font.color.rgb = RGBColor(192, 0, 0)   # Red
        else:
            logging.warning("[BRUM] Entity comparison slide not found or missing 'Table Placeholder 1'.")

        # ============================
        # Slide 11: Network Requests Deep Dive (BRUM)
        # - Rectangles: 8..12 with key metrics
        # - Table Placeholder 1: Grade summary with declined/downgraded apps
        # ============================

        # Locate slide 11 (0-based index 10), with fallback by title/placeholder.
        deep_dive_slide = prs.slides[10] if len(prs.slides) > 10 else None
        if deep_dive_slide is None or not find_table_placeholder_by_name(deep_dive_slide, "Table Placeholder 1"):
            deep_dive_slide = choose_slide_for_section(
                prefer_titles=[
                    "Network Requests Deep Dive",
                    "Network Requests",
                    "BRUM Network Requests",
                ],
                required_placeholders=("Table Placeholder 1",),
                exclude_ids=used_slide_ids,
            )
        if deep_dive_slide:
            used_slide_ids.add(id(deep_dive_slide))
        else:
            logging.warning("[BRUM] Network Requests Deep Dive slide not found.")
            deep_dive_slide = None

        # Helper to safely set text for a named rectangle/textbox on the slide.
        def set_shape_text(slide, shape_name, text):
            if slide is None:
                return False
            for shp in slide.shapes:
                if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame") and shp.text_frame:
                    shp.text_frame.clear()
                    shp.text_frame.text = str(text)
                    return True
            logging.debug("[BRUM] Shape '%s' not found on deep dive slide.", shape_name)
            return False

        # Column resolution helpers.
        def first_present_col(df, candidates):
            for c in candidates:
                if c in df.columns:
                    return c
            return None

        # ============================
        # Rectangles 8..12 — count declines per metric from df_network_requests.
        # ============================
        if deep_dive_slide:

            logging.debug("[BRUM][Slide11] df_network_requests columns: %s", list(df_network_requests.columns))

            # Resolve application column.
            def resolve_app_col(df):
                col = first_present_col(
                    df,
                    ["name", "Name", "applicationName", "Application Name", "Application"],
                )
                if col:
                    return col
                for c in df.columns:
                    if str(c).lower() in ("app", "application", "application name"):
                        return c
                # fallback: first object/text column
                for c in df.columns:
                    try:
                        if df[c].dtype == object:
                            return c
                    except Exception:
                        continue
                return None

            app_col_nr_eff = resolve_app_col(df_network_requests)
            logging.debug("[BRUM][Slide11] Rectangles: resolved app_col=%s", app_col_nr_eff)

            # Transition parsing helpers: "prev → curr (Declined)" etc.
            def parse_transition_tokens(val):
                s = str(val or "").strip()
                if "→" in s:
                    prev, curr = s.split("→", 1)
                    return prev.strip(), curr.strip().split("(")[0].strip()
                return None, None

            def token_to_bool(tok):
                t = str(tok or "").strip().lower()
                if t in {"true", "yes", "y", "1"}:
                    return True
                if t in {"false", "no", "n", "0"}:
                    return False
                return None

            def token_to_num(tok):
                try:
                    return float(str(tok).strip())
                except Exception:
                    return None

            def is_bool_decline_cell(val):
                prev, curr = parse_transition_tokens(val)
                if prev is not None and curr is not None:
                    pb = token_to_bool(prev)
                    cb = token_to_bool(curr)
                    if pb is True and cb is False:
                        return True
                s = str(val or "").lower()
                return ("declined" in s or "downgraded" in s or "decreased" in s or "reduced" in s) and "false" in s

            def is_num_decline_cell(val):
                prev, curr = parse_transition_tokens(val)
                if prev is not None and curr is not None:
                    pn = token_to_num(prev)
                    cn = token_to_num(curr)
                    if pn is not None and cn is not None and cn < pn:
                        return True
                s = str(val or "").lower()
                return any(k in s for k in ("declined", "decreased", "reduced", "down", "↓"))

            # Case-insensitive resolver for metric column names.
            def resolve_metric_col(df, candidates):
                col = first_present_col(df, candidates)
                if col:
                    return col
                # Case/spacing-insensitive.
                norm = lambda x: "".join(str(x).lower().split())
                cand_norms = {norm(c): c for c in candidates}
                for c in df.columns:
                    if norm(c) in cand_norms:
                        return c
                # Heuristic fallback by keyword overlap.
                key_tokens = set()
                for c in candidates:
                    key_tokens.update(
                        [
                            t
                            for t in str(c)
                            .lower()
                            .replace("#", "")
                            .replace("_", " ")
                            .split()
                            if t
                        ]
                    )
                best = None
                best_score = 0
                for c in df.columns:
                    lc = str(c).lower()
                    score = sum(1 for t in key_tokens if t in lc)
                    if score > best_score:
                        best, best_score = c, score
                return best if best_score >= max(2, int(len(key_tokens) * 0.4)) else None

            # Resolve metric columns.
            col_collecting = resolve_metric_col(
                df_network_requests,
                [
                    "CollectingDataPastOneDay",
                    "CollectingDataPast1Day",
                    "CollectingDataPastDay",
                    "CollectingData",
                ],
            )
            col_limit_nothit = resolve_metric_col(
                df_network_requests,
                [
                    "NetworkRequestLimitNotHit",
                    "NetworkRequestsLimitNotHit",
                    "LimitNotHit",
                    "RequestLimitNotHit",
                    "networkRequestLimitNotHit",
                ],
            )
            col_custom_rules = resolve_metric_col(
                df_network_requests,
                [
                    "CustomMatchRulesCount",
                    "# Custom Match Rules",
                    "NumCustomMatchRules",
                    "CustomMatchRules",
                    "customMatchRulesCount",
                ],
            )
            col_bt_corr = resolve_metric_col(
                df_network_requests,
                [
                    "HasBTCorrelation",
                    "BTCorrelation",
                    "BusinessTransactionCorrelation",
                ],
            )
            col_ces_include = resolve_metric_col(
                df_network_requests,
                [
                    "HasCustomEventServiceIncludeRule",
                    "CustomEventServiceIncludeRule",
                    "HasCESIncludeRule",
                ],
            )

            logging.debug(
                "[BRUM][Slide11] Rectangles: metric columns -> collecting=%s, limitNotHit=%s, customRules=%s, btCorr=%s, cesInclude=%s",
                col_collecting,
                col_limit_nothit,
                col_custom_rules,
                col_bt_corr,
                col_ces_include,
            )

            # Count declines per metric by scanning rows directly (no Analysis gating).
            def count_metric_declines(df, app_col, metric_col, is_bool, label):
                if not app_col or not metric_col:
                    logging.warning(
                        "[BRUM][Slide11] Metric '%s' missing column (%s) or app_col (%s).",
                        label,
                        metric_col,
                        app_col,
                    )
                    return 0, []
                apps = []
                for _, r in df.iterrows():
                    app = str(r.get(app_col, "") or "").strip()
                    val = r.get(metric_col, "")
                    if (is_bool and is_bool_decline_cell(val)) or (
                        (not is_bool) and is_num_decline_cell(val)
                    ):
                        apps.append(app)
                logging.info(
                    "[BRUM][Slide11] Rectangles: %s declines=%d (apps sample: %s)",
                    label,
                    len(apps),
                    apps[:10],
                )
                return len(apps), apps

            collecting_cnt, collecting_apps = count_metric_declines(
                df_network_requests,
                app_col_nr_eff,
                col_collecting,
                True,
                "Collecting Data Past One Day",
            )
            limit_not_hit_cnt, limit_not_hit_apps = count_metric_declines(
                df_network_requests,
                app_col_nr_eff,
                col_limit_nothit,
                True,
                "Network Request Limit Not Hit",
            )
            custom_rules_cnt, custom_rules_apps = count_metric_declines(
                df_network_requests,
                app_col_nr_eff,
                col_custom_rules,
                False,
                "# Custom Match Rules",
            )
            bt_corr_cnt, bt_corr_apps = count_metric_declines(
                df_network_requests,
                app_col_nr_eff,
                col_bt_corr,
                True,
                "Has BT Correlation",
            )
            ces_include_cnt, ces_include_apps = count_metric_declines(
                df_network_requests,
                app_col_nr_eff,
                col_ces_include,
                True,
                "Has Custom Event Service Include Rule",
            )

            # Explicit debug for 'Concerto' if present.
            if app_col_nr_eff and col_limit_nothit:
                concerto_rows = df_network_requests[
                    df_network_requests[app_col_nr_eff]
                    .astype(str)
                    .str.strip()
                    .str.lower()
                    == "concerto"
                ]
                if not concerto_rows.empty:
                    logging.debug(
                        "[BRUM][Slide11] Concerto networkRequestLimitNotHit cell(s): %s",
                        concerto_rows[col_limit_nothit].astype(str).tolist(),
                    )

            # Write counts into the rectangles.
            set_shape_text(deep_dive_slide, "Rectangle 8", str(collecting_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 9", str(limit_not_hit_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 10", str(custom_rules_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 11", str(bt_corr_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 12", str(ces_include_cnt))

        # ============================
        # Table — Declined-only (Network Requests) in Table Placeholder 1.
        # ============================
        if deep_dive_slide and find_table_placeholder_by_name(
            deep_dive_slide, "Table Placeholder 1"
        ):

            logging.debug(
                "[BRUM][Slide11] df_network_requests columns: %s",
                list(df_network_requests.columns),
            )
            logging.debug(
                "[BRUM][Slide11] df_analysis columns: %s", list(df_analysis.columns)
            )

            # Canonical grade order for rank comparisons.
            all_grades = ["platinum", "gold", "silver", "bronze"]
            table_grades = ["Gold", "Silver", "Bronze"]  # display order

            def norm_grade(s):
                s = str(s).strip().lower()
                for g in all_grades:
                    if g in s:
                        return g
                return None

            def parse_transition(val):
                """
                Returns (prev_grade_norm, curr_grade_norm) or (None, None) if not parsable.
                Supports:
                  - 'Gold → Silver'
                  - 'Declined to Bronze', 'Downgraded to Silver', 'Now Gold'
                """
                s = str(val or "").strip()
                if not s:
                    return (None, None)

                # Arrow format
                if "→" in s:
                    parts = s.split("→", 1)
                    prev = norm_grade(parts[0])
                    curr = norm_grade(parts[1])
                    return (prev, curr)

                # Phrasal formats
                low = s.lower()
                import re

                m = re.search(
                    r"(?:declined|downgraded)\s+(?:to\s+)?(platinum|gold|silver|bronze)",
                    low,
                )
                if m:
                    return (None, m.group(1))  # only current known
                m = re.search(
                    r"(?:now|is\s+now|became)\s+(platinum|gold|silver|bronze)", low
                )
                if m:
                    return (None, m.group(1))
                # If any grade appears, treat as current grade.
                g = norm_grade(low)
                if g:
                    return (None, g)
                return (None, None)

            # Resolve app and grade columns and also prepare row lookup.
            app_col_nr = first_present_col(
                df_network_requests,
                ["name", "Name", "applicationName", "Application Name", "Application"],
            ) or next(
                (
                    c
                    for c in df_network_requests.columns
                    if str(c).lower() in ("app", "application", "application name")
                ),
                None,
            )
            grade_col_nr = first_present_col(
                df_network_requests,
                [
                    "NetworkRequestsGrade",
                    "networkRequestsGrade",
                    "BRUMNetworkRequestsGrade",
                    "Network Requests Grade",
                    "Grade",
                    "grade",
                ],
            )

            # Row lookup by app.
            def row_for_app(app):
                if not app_col_nr:
                    return None
                match = df_network_requests[
                    df_network_requests[app_col_nr].astype(str).str.strip()
                    == str(app)
                ]
                return match.iloc[0] if not match.empty else None

            # Fallback grade resolver that scans all values in the NR row.
            def resolve_grade_for_app(app):
                # Prefer explicit grade column.
                if grade_col_nr:
                    r = row_for_app(app)
                    if r is not None:
                        g = norm_grade(r.get(grade_col_nr, ""))
                        if g:
                            return g.capitalize()
                # Scan entire row values for grade keywords.
                r = row_for_app(app)
                if r is not None:
                    for v in r.values:
                        g = norm_grade(v)
                        if g:
                            return g.capitalize()
                # Scan Analysis text for grade keywords.
                if "NetworkRequestsBRUM" in df_analysis.columns:
                    txt = df_analysis.loc[
                        df_analysis["name"].astype(str).str.strip() == str(app),
                        "NetworkRequestsBRUM",
                    ]
                    if not txt.empty:
                        _, cg = parse_transition(txt.iloc[0])
                        if cg:
                            return cg.capitalize()
                return None

            logging.debug(
                "[BRUM][Slide11] Resolved app_col_nr=%s, grade_col_nr=%s",
                app_col_nr,
                grade_col_nr,
            )

            # Build per-grade totals from NR sheet if possible; else infer from row scan or Analysis.
            totals_by_grade = {g: 0 for g in table_grades}
            grade_by_app = {}

            if app_col_nr:
                for _, r in df_network_requests.iterrows():
                    app = str(r.get(app_col_nr, "") or "").strip()
                    g = None
                    if grade_col_nr:
                        g = norm_grade(r.get(grade_col_nr, ""))
                    if not g:
                        # scan row values
                        for v in r.values:
                            g = norm_grade(v)
                            if g:
                                break
                    disp = g.capitalize() if g else None
                    if app and disp in totals_by_grade:
                        totals_by_grade[disp] += 1
                        grade_by_app[app] = disp

            if not any(v > 0 for v in totals_by_grade.values()) and "NetworkRequestsBRUM" in df_analysis.columns:
                # Infer totals from Analysis mentions.
                inferred_totals = {g: 0 for g in table_grades}
                for _, r in df_analysis.iterrows():
                    _, cg = parse_transition(r.get("NetworkRequestsBRUM", ""))
                    disp = cg.capitalize() if cg else None
                    if disp in inferred_totals:
                        inferred_totals[disp] += 1
                if any(inferred_totals.values()):
                    totals_by_grade = inferred_totals
                    logging.warning(
                        "[BRUM][Slide11] Using inferred per-grade totals from Analysis: %s",
                        totals_by_grade,
                    )

            # A) Declines detected from Analysis.
            declined_apps_analysis = set()
            if "NetworkRequestsBRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    val = r.get("NetworkRequestsBRUM", "")
                    prev_g, curr_g = parse_transition(val)
                    is_dg = False
                    low = str(val).lower()
                    if "declined" in low or "downgraded" in low:
                        is_dg = True
                    elif (
                        prev_g
                        and curr_g
                        and prev_g in all_grades
                        and curr_g in all_grades
                    ):
                        is_dg = all_grades.index(prev_g) < all_grades.index(curr_g)
                    if is_dg and app:
                        declined_apps_analysis.add(app)

            logging.info(
                "[BRUM][Slide11] Declined apps from Analysis: %d",
                len(declined_apps_analysis),
            )
            logging.debug(
                "[BRUM][Slide11] Declined apps (Analysis) sample: %s",
                list(sorted(declined_apps_analysis))[:20],
            )

            # B) Declines detected from metric columns (same detectors as rectangles).
            metric_cols_and_types = [
                (col_collecting, True),
                (col_limit_nothit, True),
                (col_custom_rules, False),
                (col_bt_corr, True),
                (col_ces_include, True),
            ]
            metric_declined_apps = set()
            if app_col_nr:
                for col, is_bool in metric_cols_and_types:
                    if not col:
                        continue
                    for _, r in df_network_requests.iterrows():
                        app = str(r.get(app_col_nr, "") or "").strip()
                        val = r.get(col, "")
                        if (is_bool and is_bool_decline_cell(val)) or (
                            (not is_bool) and is_num_decline_cell(val)
                        ):
                            metric_declined_apps.add(app)

            logging.info(
                "[BRUM][Slide11] Declined apps from NR metrics: %d",
                len(metric_declined_apps),
            )
            logging.debug(
                "[BRUM][Slide11] Declined apps (Metrics) sample: %s",
                list(sorted(metric_declined_apps))[:20],
            )

            # Union of Analysis- and Metric-detected declines.
            declined_union = declined_apps_analysis.union(metric_declined_apps)
            logging.info(
                "[BRUM][Slide11] Total declined apps (union): %d", len(declined_union)
            )

            # Group apps under destination grade.
            declined_by_grade = {g: [] for g in table_grades}
            missing_grade = []
            for app in sorted(declined_union):
                dest = grade_by_app.get(app)
                if not dest:
                    dest = resolve_grade_for_app(app)
                if dest in declined_by_grade:
                    declined_by_grade[dest].append(app)
                else:
                    missing_grade.append(app)

            if missing_grade:
                logging.warning(
                    "[BRUM][Slide11] %d declined apps have no resolvable grade (not shown in table): %s",
                    len(missing_grade),
                    missing_grade[:20],
                )

            logging.debug(
                "[BRUM][Slide11] Declined-by-grade counts: %s",
                {k: len(v) for k, v in declined_by_grade.items()},
            )
            logging.debug(
                "[BRUM][Slide11] Final totals_by_grade: %s", totals_by_grade
            )

            # Build table rows in Gold, Silver, Bronze order.
            headers = [
                "Grade",
                "Application Names",
                "Number of Applications",
                "Percentage Declined",
            ]
            rows = []
            for g in table_grades:
                names = sorted(declined_by_grade[g], key=str.lower)
                num_apps = len(names)
                denom = totals_by_grade.get(g, 0)
                pct = round((num_apps / denom) * 100) if denom > 0 else 0
                logging.info(
                    "[BRUM][Slide11] Grade=%s declined=%d total=%d pct=%d%%",
                    g,
                    num_apps,
                    denom,
                    pct,
                )
                rows.append(
                    [g, "\n".join(names) if names else "—", str(num_apps), f"{pct}%"]
                )

            # Insert and render the single table.
            table = insert_table_at_placeholder(
                deep_dive_slide, "Table Placeholder 1", len(rows) + 1, len(headers)
            )
            if table:
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(12)

                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = val
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(12)
                        if c_idx == 3:
                            run = p.runs[0] if p.runs else p.add_run()
                            pct_num = int(str(val).replace("%", "") or 0)
                            run.font.color.rgb = (
                                RGBColor(192, 0, 0)
                                if pct_num > 0
                                else RGBColor(0, 176, 80)
                            )
            else:
                logging.error(
                    "[BRUM][Slide11] Could not insert Network Requests Declined table; placeholder missing."
                )

        # ============================
        # Slide 12: Health Rules & Alerting Deep Dive (BRUM)
        # - Rectangles: 10..12 with key metrics
        # - Table Placeholder 1: Grade summary with declined/downgraded apps
        # ============================

        # Locate slide 12 (0-based index 11), fallback by title/placeholder.
        hra_deep_dive_slide = prs.slides[11] if len(prs.slides) > 11 else None
        if hra_deep_dive_slide is None or not find_table_placeholder_by_name(
            hra_deep_dive_slide, "Table Placeholder 1"
        ):
            hra_deep_dive_slide = choose_slide_for_section(
                prefer_titles=[
                    "Health Rules & Alerting Deep Dive",
                    "Health Rules Deep Dive",
                    "Health Rules & Alerting",
                    "Health Rules",
                ],
                required_placeholders=("Table Placeholder 1",),
                exclude_ids=used_slide_ids,
            )
        if hra_deep_dive_slide:
            used_slide_ids.add(id(hra_deep_dive_slide))
        else:
            logging.warning("[BRUM] Health Rules & Alerting Deep Dive slide not found.")
            hra_deep_dive_slide = None

        # Helpers (reuse patterns from Slide 11).
        def set_shape_text(slide, shape_name, text):
            if slide is None:
                return False
            for shp in slide.shapes:
                if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame") and shp.text_frame:
                    shp.text_frame.clear()
                    shp.text_frame.text = str(text)
                    return True
            logging.debug("[BRUM] Shape '%s' not found on HRA deep dive slide.", shape_name)
            return False

        def first_present_col(df, candidates):
            for c in candidates:
                if c in df.columns:
                    return c
            return None

        # ============================
        # Rectangles 10..12 — count declines per metric from df_health_rules.
        # ============================
        if hra_deep_dive_slide:

            logging.debug(
                "[BRUM][Slide12] df_health_rules columns: %s",
                list(df_health_rules.columns),
            )

            # Resolve application column.
            def resolve_app_col(df):
                col = first_present_col(
                    df,
                    ["name", "Name", "applicationName", "Application Name", "Application"],
                )
                if col:
                    return col
                for c in df.columns:
                    if str(c).lower() in ("app", "application", "application name"):
                        return c
                # fallback: first text/object column
                for c in df.columns:
                    try:
                        if df[c].dtype == object:
                            return c
                    except Exception:
                        continue
                return None

            app_col_hr_eff = resolve_app_col(df_health_rules)
            logging.debug("[BRUM][Slide12] Rectangles: resolved app_col=%s", app_col_hr_eff)

            # Transition parsing helpers: "prev → curr (Declined)" etc.
            def parse_transition_tokens(val):
                s = str(val or "").strip()
                if "→" in s:
                    prev, curr = s.split("→", 1)
                    return prev.strip(), curr.strip().split("(")[0].strip()
                return None, None

            def token_to_num(tok):
                try:
                    return float(str(tok).strip())
                except Exception:
                    return None

            def is_num_decline_cell(val):
                prev, curr = parse_transition_tokens(val)
                if prev is not None and curr is not None:
                    pn = token_to_num(prev)
                    cn = token_to_num(curr)
                    if pn is not None and cn is not None and cn < pn:
                        return True
                s = str(val or "").lower()
                return any(k in s for k in ("declined", "decreased", "reduced", "down", "↓"))

            # Case-insensitive resolver for metric column names.
            def resolve_metric_col(df, candidates):
                col = first_present_col(df, candidates)
                if col:
                    return col
                # Case/spacing-insensitive match.
                norm = lambda x: "".join(str(x).lower().split())
                cand_norms = {norm(c): c for c in candidates}
                for c in df.columns:
                    if norm(c) in cand_norms:
                        return c
                # Heuristic fallback by keyword overlap.
                key_tokens = set()
                for c in candidates:
                    key_tokens.update(
                        [
                            t
                            for t in str(c)
                            .lower()
                            .replace("#", "")
                            .replace("_", " ")
                            .split()
                            if t
                        ]
                    )
                best = None
                best_score = 0
                for c in df.columns:
                    lc = str(c).lower()
                    score = sum(1 for t in key_tokens if t in lc)
                    if score > best_score:
                        best, best_score = c, score
                return best if best_score >= max(2, int(len(key_tokens) * 0.4)) else None

            # Resolve metric columns (common headers + variants).
            col_violations = resolve_metric_col(
                df_health_rules,
                [
                    "NumberOfHealthRuleViolations",
                    "# Of Health Rule Violations",
                    "HealthRuleViolations",
                    "HealthRulesViolations",
                    "numHealthRuleViolations",
                ],
            )
            col_actions_bound = resolve_metric_col(
                df_health_rules,
                [
                    "NumberOfActionsBoundToEnabledPolicies",
                    "# Of Actions Bound To Enabled Policies",
                    "ActionsBoundToEnabledPolicies",
                    "ActionsBoundEnabledPolicies",
                    "ActionsBoundEnabledPoliciesCount",
                ],
            )
            col_custom_rules = resolve_metric_col(
                df_health_rules,
                [
                    "NumberOfCustomHealthRules",
                    "# Of Custom Health Rules",
                    "CustomHealthRulesCount",
                    "CustomHealthRules",
                ],
            )

            logging.debug(
                "[BRUM][Slide12] Rectangles: metric columns -> violations=%s, actionsBound=%s, customRules=%s",
                col_violations,
                col_actions_bound,
                col_custom_rules,
            )

            # Count numeric declines per metric.
            def count_metric_declines(df, app_col, metric_col, label):
                if not app_col or not metric_col:
                    logging.warning(
                        "[BRUM][Slide12] Metric '%s' missing column (%s) or app_col (%s).",
                        label,
                        metric_col,
                        app_col,
                    )
                    return 0, []
                apps = []
                for _, r in df.iterrows():
                    app = str(r.get(app_col, "") or "").strip()
                    val = r.get(metric_col, "")
                    if is_num_decline_cell(val):
                        apps.append(app)
                logging.info(
                    "[BRUM][Slide12] Rectangles: %s declines=%d (apps sample: %s)",
                    label,
                    len(apps),
                    apps[:10],
                )
                return len(apps), apps

            violations_cnt, violations_apps = count_metric_declines(
                df_health_rules,
                app_col_hr_eff,
                col_violations,
                "Number Of Health Rule Violations",
            )
            actions_bound_cnt, actions_bound_apps = count_metric_declines(
                df_health_rules,
                app_col_hr_eff,
                col_actions_bound,
                "Number Of Actions Bound To Enabled Policies",
            )
            custom_rules_cnt, custom_rules_apps = count_metric_declines(
                df_health_rules,
                app_col_hr_eff,
                col_custom_rules,
                "Number Of Custom Health Rules",
            )

            # Write counts into the rectangles.
            set_shape_text(hra_deep_dive_slide, "Rectangle 10", str(violations_cnt))
            set_shape_text(hra_deep_dive_slide, "Rectangle 11", str(actions_bound_cnt))
            set_shape_text(hra_deep_dive_slide, "Rectangle 12", str(custom_rules_cnt))

        # ============================
        # Table — Declined-only (Health Rules & Alerting) in Table Placeholder 1.
        # ============================
        if hra_deep_dive_slide and find_table_placeholder_by_name(
            hra_deep_dive_slide, "Table Placeholder 1"
        ):

            logging.debug(
                "[BRUM][Slide12] df_health_rules columns: %s",
                list(df_health_rules.columns),
            )
            logging.debug(
                "[BRUM][Slide12] df_analysis columns: %s", list(df_analysis.columns)
            )

            # Canonical grade order for rank comparisons.
            all_grades = ["platinum", "gold", "silver", "bronze"]
            table_grades = ["Gold", "Silver", "Bronze"]

            def norm_grade(s):
                s = str(s).strip().lower()
                for g in all_grades:
                    if g in s:
                        return g
                return None

            def parse_transition(val):
                s = str(val or "").strip()
                if not s:
                    return (None, None)
                if "→" in s:
                    parts = s.split("→", 1)
                    prev = norm_grade(parts[0])
                    curr = norm_grade(parts[1])
                    return (prev, curr)
                low = s.lower()
                import re

                m = re.search(
                    r"(?:declined|downgraded)\s+(?:to\s+)?(platinum|gold|silver|bronze)",
                    low,
                )
                if m:
                    return (None, m.group(1))
                m = re.search(
                    r"(?:now|is\s+now|became)\s+(platinum|gold|silver|bronze)", low
                )
                if m:
                    return (None, m.group(1))
                g = norm_grade(low)
                if g:
                    return (None, g)
                return (None, None)

            # Resolve app and grade columns and prepare row lookup.
            app_col_hr = first_present_col(
                df_health_rules,
                ["name", "Name", "applicationName", "Application Name", "Application"],
            ) or next(
                (
                    c
                    for c in df_health_rules.columns
                    if str(c).lower() in ("app", "application", "application name")
                ),
                None,
            )
            grade_col_hr = first_present_col(
                df_health_rules,
                [
                    "HealthRulesAndAlertingGrade",
                    "HealthRulesGrade",
                    "BRUMHealthRulesGrade",
                    "Health Rules Grade",
                    "Grade",
                    "grade",
                ],
            )

            def row_for_app_hr(app):
                if not app_col_hr:
                    return None
                match = df_health_rules[
                    df_health_rules[app_col_hr].astype(str).str.strip() == str(app)
                ]
                return match.iloc[0] if not match.empty else None

            def resolve_grade_for_app_hr(app):
                # Prefer explicit grade column.
                if grade_col_hr:
                    r = row_for_app_hr(app)
                    if r is not None:
                        g = norm_grade(r.get(grade_col_hr, ""))
                        if g:
                            return g.capitalize()
                # Scan entire row values for grade keywords.
                r = row_for_app_hr(app)
                if r is not None:
                    for v in r.values:
                        g = norm_grade(v)
                        if g:
                            return g.capitalize()
                # Scan Analysis text for grade keywords.
                if "HealthRulesAndAlertingBRUM" in df_analysis.columns:
                    txt = df_analysis.loc[
                        df_analysis["name"].astype(str).str.strip() == str(app),
                        "HealthRulesAndAlertingBRUM",
                    ]
                    if not txt.empty:
                        _, cg = parse_transition(txt.iloc[0])
                        if cg:
                            return cg.capitalize()
                return None

            logging.debug(
                "[BRUM][Slide12] Resolved app_col_hr=%s, grade_col_hr=%s",
                app_col_hr,
                grade_col_hr,
            )

            # Build per-grade totals from HRA sheet if possible; else infer.
            totals_by_grade_hr = {g: 0 for g in table_grades}
            grade_by_app_hr = {}

            if app_col_hr:
                for _, r in df_health_rules.iterrows():
                    app = str(r.get(app_col_hr, "") or "").strip()
                    g = None
                    if grade_col_hr:
                        g = norm_grade(r.get(grade_col_hr, ""))
                    if not g:
                        for v in r.values:
                            g = norm_grade(v)
                            if g:
                                break
                    disp = g.capitalize() if g else None
                    if app and disp in totals_by_grade_hr:
                        totals_by_grade_hr[disp] += 1
                        grade_by_app_hr[app] = disp

            if not any(v > 0 for v in totals_by_grade_hr.values()) and "HealthRulesAndAlertingBRUM" in df_analysis.columns:
                inferred_totals = {g: 0 for g in table_grades}
                for _, r in df_analysis.iterrows():
                    _, cg = parse_transition(r.get("HealthRulesAndAlertingBRUM", ""))
                    disp = cg.capitalize() if cg else None
                    if disp in inferred_totals:
                        inferred_totals[disp] += 1
                if any(inferred_totals.values()):
                    totals_by_grade_hr = inferred_totals
                    logging.warning(
                        "[BRUM][Slide12] Using inferred per-grade totals from Analysis: %s",
                        totals_by_grade_hr,
                    )

            # A) Declines detected from Analysis.
            declined_apps_analysis_hr = set()
            if "HealthRulesAndAlertingBRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    val = r.get("HealthRulesAndAlertingBRUM", "")
                    prev_g, curr_g = parse_transition(val)
                    is_dg = False
                    low = str(val).lower()
                    if "declined" in low or "downgraded" in low:
                        is_dg = True
                    elif (
                        prev_g
                        and curr_g
                        and prev_g in all_grades
                        and curr_g in all_grades
                    ):
                        is_dg = all_grades.index(prev_g) < all_grades.index(curr_g)
                    if is_dg and app:
                        declined_apps_analysis_hr.add(app)

            logging.info(
                "[BRUM][Slide12] Declined apps from Analysis: %d",
                len(declined_apps_analysis_hr),
            )
            logging.debug(
                "[BRUM][Slide12] Declined apps (Analysis) sample: %s",
                list(sorted(declined_apps_analysis_hr))[:20],
            )

            # B) Declines detected from HRA metric columns (use same parser as rectangles).
            metric_cols_hr = [col_violations, col_actions_bound, col_custom_rules]
            metric_declined_apps_hr = set()
            if app_col_hr:
                for col in metric_cols_hr:
                    if not col:
                        continue
                    for _, r in df_health_rules.iterrows():
                        app = str(r.get(app_col_hr, "") or "").strip()
                        val = r.get(col, "")
                        if is_num_decline_cell(val):
                            metric_declined_apps_hr.add(app)

            logging.info(
                "[BRUM][Slide12] Declined apps from HRA metrics: %d",
                len(metric_declined_apps_hr),
            )
            logging.debug(
                "[BRUM][Slide12] Declined apps (Metrics) sample: %s",
                list(sorted(metric_declined_apps_hr))[:20],
            )

            # Union of Analysis- and Metric-detected declines.
            declined_union_hr = declined_apps_analysis_hr.union(
                metric_declined_apps_hr
            )
            logging.info(
                "[BRUM][Slide12] Total declined apps (union): %d",
                len(declined_union_hr),
            )

            # Group apps under destination grade.
            declined_by_grade_hr = {g: [] for g in table_grades}
            missing_grade_hr = []
            for app in sorted(declined_union_hr):
                dest = grade_by_app_hr.get(app)
                if not dest:
                    dest = resolve_grade_for_app_hr(app)
                if dest in declined_by_grade_hr:
                    declined_by_grade_hr[dest].append(app)
                else:
                    missing_grade_hr.append(app)

            if missing_grade_hr:
                logging.warning(
                    "[BRUM][Slide12] %d declined apps have no resolvable grade (not shown in table): %s",
                    len(missing_grade_hr),
                    missing_grade_hr[:20],
                )

            logging.debug(
                "[BRUM][Slide12] Declined-by-grade counts: %s",
                {k: len(v) for k, v in declined_by_grade_hr.items()},
            )
            logging.debug(
                "[BRUM][Slide12] Final totals_by_grade_hr: %s", totals_by_grade_hr
            )

            # Build table in Gold, Silver, Bronze order.
            headers = [
                "Grade",
                "Application Names",
                "Number of Applications",
                "Percentage Declined",
            ]
            rows = []
            for g in table_grades:
                names = sorted(declined_by_grade_hr[g], key=str.lower)
                num_apps = len(names)
                denom = totals_by_grade_hr.get(g, 0)
                pct = round((num_apps / denom) * 100) if denom > 0 else 0
                logging.info(
                    "[BRUM][Slide12] Grade=%s declined=%d total=%d pct=%d%%",
                    g,
                    num_apps,
                    denom,
                    pct,
                )
                rows.append(
                    [g, "\n".join(names) if names else "—", str(num_apps), f"{pct}%"]
                )

            table = insert_table_at_placeholder(
                hra_deep_dive_slide,
                "Table Placeholder 1",
                len(rows) + 1,
                len(headers),
            )
            if table:
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(12)

                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = val
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(12)
                        if c_idx == 3:
                            run = p.runs[0] if p.runs else p.add_run()
                            pct_num = int(str(val).replace("%", "") or 0)
                            run.font.color.rgb = (
                                RGBColor(192, 0, 0)
                                if pct_num > 0
                                else RGBColor(0, 176, 80)
                            )
            else:
                logging.error(
                    "[BRUM][Slide12] Could not insert Health Rules & Alerting Declined table; placeholder missing."
                )

        # ============================
        # Populate "TextBox 7" with number of BRUM applications (slide index 5)
        # ============================

        def set_textbox_value(prs_obj, shape_name, text, fallback_slide_index=5):
            for s in prs_obj.slides:
                for shp in s.shapes:
                    if getattr(shp, "name", "") == shape_name and hasattr(
                        shp, "text_frame"
                    ):
                        shp.text_frame.clear()
                        shp.text_frame.text = str(text)
                        return True
            if len(prs_obj.slides) > fallback_slide_index:
                s = prs_obj.slides[fallback_slide_index]
                for shp in s.shapes:
                    if getattr(shp, "name", "") == shape_name and hasattr(
                        shp, "text_frame"
                    ):
                        shp.text_frame.clear()
                        shp.text_frame.text = str(text)
                        return True
            return False

        if not set_textbox_value(prs, "TextBox 7", number_of_apps, fallback_slide_index=5):
            logging.warning(
                "[BRUM] 'TextBox 7' not found; BRUM application count not written."
            )

        # Save the presentation when all slide content is populated.
        prs.save(powerpoint_output_path)
        logging.debug(f"[BRUM] PowerPoint saved to: {powerpoint_output_path}")
        return powerpoint_output_path
    
    except Exception as e:
        log.error("[BRUM] Error generating PowerPoint: %s", e, exc_info=True)
        # Re-raise so upstream code (CLI / API) sees the failure.
        raise